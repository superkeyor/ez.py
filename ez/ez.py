from __future__ import print_function
from __future__ import unicode_literals
from __future__ import absolute_import

__author__ = ""

__doc__ = """This module is for easy interaction with linux, Mac OS X, Windows shell 
(Should Support both Python 2 and 3)...
=============================================
To see your python version
in terminal: python -V
or in python: import sys; print (sys.version)
=============================================
Install:
https://pypi.python.org/pypi/ez
pip install ez
dependencies:
pip install chardet psutil xlsxwriter xlwt xlrd pytz tzlocal pyperclip gmail

Almost all commands support the usage of '~', '..', '.', '?', '*' in path (ls,fls only support regular expression).
Symbolic link itself is the target of file operations; the actual file should be safe.

see also pyg.py
debug(1/0)
    # 0 = everything will be actually executed
    # 1 = simulate operations of cp, mv, execute; other commands will be actually performed.
          will print out simulated commands, useful for debugging and for counting files when necessary.
error(msg)

readx(path, sheet=0, r=[1,], c=None)  # Read xlsx, xls file into a list
savex(path, data, header=None, delimiter=",", sheet_name='Sheet1') # Write a list of list to a xlsx (xlsxwriter), xls(xlwt), csv file

fullpath(path) fp()
pwd() or cwd()  # Returns current working director.
csd() # Returns full path of current script directory, i.e. the directory where the running script is. 
csf() # Returns current script name without ext.
parentdir(path) pr() # Returns the parent directory of a path.
joinpath(path1[, path2[, ...]])   jp() # Returns the joined path. Supports vectorization.
splitpath(path) sp() # Returns a list of path elements: [path, file, ext]. Supports vectorization.
cd(path), ce(), cf()    # Changes to a new working directory.
stepfolder(-1)

trim(string,how[,chars]) quote(string)
join(sep,string1,string2), join(sep,array) # Glues together strings with sep. Supports vectorization.
sort(array)
replace(theList,theItem,replacement), remove(theList,theItem)

ls([path[, regex]], full=True, dotfile=False)    # Returns a list of all (including hidden) files with their full paths in path, filtered by regular expression.
lsd([path[, regex]], full=False, dotfolder=False)
fls([path[, regex, dotf=False]])   # Returns a list of files with their full paths in flattened path (i.e. walk each subdirectory).
# the filter only works for short file name not for full file name, i.e. the file name itself not its full path
# regular expression is case-sensitive
# usage: ls(); ls(cwd()); ls(cwd(), "\.py$")

mkdir("path/to/a/directory")    # Makes a directory (also any one of the "path", "to", "a" directories if not exits).
rn(old, new) # Renames old to new.
exists(path)    # Returns the existence of path (0 or 1).
rm(path)    # Deletes a file or folder. Supports wildcards, vectorization.
cp(source, destination)  # Copies source file(s) or folder to destination. Supports wildcards, vectorization.
mv(source, destination)  # Moves source file(s) or folder to destination. Supports wildcards, vectorization.

sprintf(formatString, *args, **kwargs)
evaluate(exp)
# Executes a shell command
execute0 # captures output (subprocess.Popen), mnemonic: extra 0 for output
execute  # no capture output (subprocess.call)
esp(0)   # execute sprintf shell commands
espR(0)  # execute sprintf R codes
espA(0)  # execute sprintf applescript codes
with nooutput():
    print 'this is will not be printed in stdout'
pprint(text,color='green') # color print; ppprint() # "pretty-print" arbitrary Python data structures
beep()  # Beeps to notify user.
where(name) # Prints where a module is and in which module a function is. where('python') returns which python is being used.
help(name)/doc(name) # name is a string, Prints the doc string of a module/class/function
    when write a module, add:
    __doc__ = three double quotes blabla three double quotes         <-----this is module's docstring, use explicit

    when write a function/class:
    def function(arg):
        three double quotes Returns, blabla three double quotes      <-----this is function's doctoring, use implicit
        return sth
ver(package_name) version(package_name), see a package's version.  package_name could be 'python'
whos(name),whos() list imported functions/packages

logon(file="log.txt", mode='a', status=True, timestamp=True), logoff()

tree([path[, sum=True, save=None, sort=True, case=True]) # Prints a directory tree structure. 
    sum=True (default) prints only folders, i.e., print less to show the big structure
    sum=False prints files plus folders

[starts, ends] = regexp(string, pattern); regexp(string, pattern, method='split/match'), regexpi
regexprep(string, pattern, replace, count=0), regexprepi

iff(expression, result1, result2), ifelse()
clear(module, recursive=False)

num(string)
isempty(s)
Randomize(x), randomize(x) # Sets a randomization seed.
RandomizeArray(list=[])   randomizearray(list=[])  # Shuffles a list in place.
Random(a,b) # Returns a random integer N such that a <= N <= b.
RandomChoice(seq), randomchoice(seq) # Returns a random element from sequence
Permute(iterable=[]) permute(iterable=[]) # Returns permutations in a list

unique(seq), union(seq1,seq2), intersect(seq1,seq2), setdiff(seq1,seq2) in original order
seq could be a list
    note: setdiff(seq1,seq2) may not be equal to setdiff(seq2,seq1)
            >>> unique('abracadaba')
            ['a', 'b', 'r', 'c', 'd']
            >>> unique('simsalabim')
            ['s', 'i', 'm', 'a', 'l', 'b']
            >>>
            >>> setdiff('abracadaba','simsalabim')
            ['r', 'c', 'd']
            >>> setdiff('simsalabim','abracadaba')
            ['s', 'i', 'm', 'l']
duplicate(seq) # returns a list of duplicated elements in original order
    # e.g.
    # a = [1,5,2,3,2,1,5,6,5,5,5]
    # duplicate(a) # yields [2, 1, 5]

JDict() # Jerry's dictionary, customized ordered dictionary class with convient attributes and methods, see help(JDict)
Moment(timezone)    # Generates the current datetime in specified timezone, or local naive datetime if omitted.

SetClip(content), setclip(content)   # Copy/Write something to current clipboard
content = GetClip(), content = getclip()   # Read out content from current clipboard and assign to a variable

lines(path='.', pattern='\.py$|.ini$|\.c$|\.h$|\.m$', recursive=True) # Counts lines of codes, counting empty lines as well.
keygen(length=8, complexity=3)  # generate a random key
hashes(filename): # Calculate/Print a file's md5 32; sha1 32; can handle big files in a memory efficient way
pinyinauthor()
encoding_detect(), encoding_convert()
hanzifreq()
gmail(), mail(), Mail()
    To avoid typing email password each time, place a file named pysecrets.py with
    EMAIL = 'someone@gmail.com'
    PASSWORD = 'abcdefghik'
    or better pysecrets.pyc
    in the site-packages/ez folder, check with ez.where('ez') or os.path.dirname(ez.__file__)
    The functions will no longer need email/password and become like this
    Mail(to, subject, body, attach=None)

"""
# for __all__
import fire,requests,sqlite3,json,io,os,sys,platform,string,random,re,datetime,tzlocal,pandas as pd,numpy as np,urllib.parse as urlparse;from bs4 import BeautifulSoup

# reference: abspath for ../ ./, expanduser for ~, glob to resolve wildcards, fnmatch.translate wildcards to re
import os, sys, platform, string, random, shutil, re, subprocess, glob, fnmatch
from os.path import abspath, basename, dirname, splitext, isfile, isdir, realpath, expanduser
try:
    import configparser
except:
    import ConfigParser

_DEBUG_MODE = 0
def ShellDebug(debugMode=1):
    global _DEBUG_MODE
    _DEBUG_MODE = debugMode
debug = ShellDebug    

# secretes
try: # local
    from . pysecrets import *
except:
    pass
    # try: # heroku
    #     if '/app/fz' in sys.path: sys.path.remove('/app/fz')
    #     sys.path.insert(0,'/app/fz')
    #     from pysecrets import *
    # except:
    #     pass

def clorox(path='~/Dropbox',delete=False):
    """
    Clean up certain files. Why clorox? -My favorite(?) cleaning product brand.
    Mac: ^\._; Windows: ^\._|\.DS_Store|^Icon.$
    delete: actually delete; if False, return file lists
    """
    if platform.system()=='Darwin':
        remover='^\._'  
    else:
        remover='^\._|\.DS_Store|^Icon.$'
    
    stains = fls(path,remover,dotf=True)

    if delete:
        if platform.system()=='Darwin': 
            execute(f'dot_clean "{path}"')
        else:
            rm(stains)
    return stains
clx=clorox

def retry(*args,timeout=20,n=None,min=4/60,max=10):
    """
    args: positional arguments, supports both @retry and @retry() as valid syntax
          @retry as a decorator
          retry(pd.read_html)(another_url)
          f = retry(lambda x: pd.read_html(x)); f(another_url)
    timeout: in minutes
    n: number of retries; if specified, ignore timeout
          4->10 min
          10->20 min
          12->30 min
          15->60 min (after 15 tries, wait the max)
    min: in minutes, min time to wait between retries
    max: in minutes, max time to wait between retries
    """
    # debug: ez.retry(lambda x:x/0)()
    # https://stackoverflow.com/a/62132401/2292993
    # https://tenacity.readthedocs.io/en/latest/
    # https://github.com/jd/tenacity
    # use parameter "after" not "before", which prints even when first try
    # https://github.com/jd/tenacity/blob/master/tenacity/__init__.py
    if len(args) == 1 and callable(args[0]):
        return retry(timeout=timeout,n=n,min=min,max=max)(args[0])
    else:
        import tenacity
        import functools
        if n is None:
            stop=tenacity.stop_after_delay(timeout*60)
        else:
            stop=tenacity.stop_after_attempt(n)
        myretry=functools.partial(
            tenacity.retry,
            reraise=True,
            stop=stop,wait=tenacity.wait_random_exponential(multiplier=1, min=min*60, max=max*60),
            retry=tenacity.retry_if_exception_type(),after=lambda retry_state: print(f"Retry... {retry_state.attempt_number}")
        )()
        return myretry

def rget(url,headers={},timeout=20,n=None,*args,**kwargs):
    """
    see retry for parameters
    """
    return retry(lambda : requests.get(url,headers=headers,*args,**kwargs).json(), timeout=timeout, n=n)()

def rgetraw(url,headers={},timeout=20,n=None,*args,**kwargs):
    return retry(lambda : requests.get(url,headers=headers,*args,**kwargs), timeout=timeout, n=n)()

def rgetfile(url,file,headers={},timeout=20,n=None,*args,**kwargs):
    r = retry(lambda : requests.get(url,headers=headers,*args,**kwargs), timeout=timeout, n=n)()
    with open(file, 'wb') as f: f.write(r.content) 

# json is recommended for config
# def ReadConfig(item):
#    """Read a variable from the config.ini file"""
#    config = ConfigParser.RawConfigParser()
#    config.read("config.ini")
#    return config.get("Default",item).strip("\"\" ")

from os import getcwd as pwd
from os import getcwd as cwd

def error(msg):
    raise Exception(msg)

def fullpath(path):
    """
    fullpath(path) # Returns the full path by resolving ~ and relative path.
    note: no trailing / returned, at least on mac os x
    """
    # https://stackoverflow.com/a/40311142/2292993
    # os.path.abspath returns the absolute path, but does NOT resolve symlinks
    # os.path.realpath will first resolve any symbolic links in the path, and then return the absolute path.
    # both abspath and realpath imply os.path.normpath
    # neither abspath or realpath will resolve ~ to the user's home directory
    # abspath and realpath: if fullpath provided, simply return fullpath. if not, resolved relative to pwd
    return os.path.abspath(os.path.expanduser(path))
fp = fullpath

def csd():
    """(),Returns full path of current script directory, i.e. the directory where the running script is (not the location of imported file).
    if in interactive mode, return current working directory
    It should always be the best to call csd() at the top of a script, 
    during which the running script has not change any working directory yet. 
    Therefore the internal call of sys.argv[0], if given not in full path, can be resolved to correct full path
    """
    # https://stackoverflow.com/a/22424821/2292993
    import __main__ as main
    is_interactive = not hasattr(main, '__file__')
    if is_interactive:
        return os.getcwd()
    else:
        # for difference between __file__ and sys.argv[0]
        # see https://stackoverflow.com/questions/5851588/difference-between-file-and-sys-argv0
        # os.path.split returns (head,tail), here for path = , same effect as splitpath
        path = os.path.split(os.path.abspath(sys.argv[0]))[0]
        # hack when a script is packed into an app, which returns xxx.app/Contents/Resources
        return os.path.abspath(os.path.join(path,os.pardir,os.pardir,os.pardir)) if path.endswith('.app/Contents/Resources') else path

def stepfolder(step=-1):
    """
    folder = stepfolder(step)
    step: regex of folder, eg '01', if multiple match, return only the first matched
          integer -1 (default, backward) +2 (forward) 2 (same as +2)
    folder: target folder path
    Usage:
          Under a projFolder, there might be 01Original, 02Slicing, 03Motion...
          Use this function to go to a specific "step folder"
    """
    import __main__ as main
    is_interactive = not hasattr(main, '__file__')
    if is_interactive:
        theCSD = os.getcwd()
    else:
        path = splitpath(os.path.abspath(sys.argv[0]))[0]
        # hack when a script is packed into an app, which returns xxx.app/Contents/Resources
        theCSD = os.path.abspath(os.path.join(path,os.pardir,os.pardir,os.pardir)) if path.endswith('.app/Contents/Resources') else path
    projFolder = parentdir(theCSD)
    if isinstance(step,str):
        folder = lsd(projFolder,step)
        folder = folder[0]
        folder = joinpath(projFolder,folder)
    elif isinstance(step,int):
        steps = lsd(projFolder,'^\d\d')  # a list of all folders like 01Original, 06Set
        currentStep = splitpath(theCSD)[1]
        try:
            currentStepNum = steps.index(currentStep)
            targetStep = steps[currentStepNum+step]
        except:
            targetStep = currentStep
        folder = joinpath(projFolder,targetStep)
    return folder

def csf():
    """(),Returns current script file name without ext, i.e. the name of the running script without ext.
    different from csd(), even if working directory changes, csf() still return correct sys.argv[0]
    """
    import __main__ as main
    is_interactive = not hasattr(main, '__file__')
    if is_interactive:
        path = os.getcwd()
    else:
        # os.path.split returns (head, tail)
        file = os.path.split(os.path.abspath(sys.argv[0]))[1]
    return file

def parentdir(path):
    """parentdir(path), Returns the parent directory of a path."""
    path = fullpath(path)
    dir = os.path.dirname(path)
    ext = os.path.splitext(path)[1]
    file = os.path.basename(path)[0:-len(ext)] if len(ext) != 0 else os.path.basename(path)
    # if isfile
    if len(ext) != 0:
        return os.path.abspath(os.path.join(os.path.abspath(os.path.join(path, os.pardir)), os.pardir))
    else:
        return os.path.abspath(os.path.join(path, os.pardir))
pr = parentdir

def joinpath(*args):
    """joinpath(path1[, path2[, ...]])   # Returns the joined path. Supports vectorization.
    """
    vectorization = False
    # to see whether there is an list/tuple input
    for arg in args:
        if type(arg) in [list, tuple]:
            vectorization = True
            length = len(arg)
            break
    if vectorization:
        # fill arg which is a string
        # args is by default an unmutable tuple
        args = list(args)
        for j in range(0,len(args)):
            arg = args[j]
            if type(arg) in [str]:
                args[j] = [arg]*length

        results = []
        for i in range(0,length):
            result = ''
            for j in range(0,len(args)):
                result = os.path.join(result, args[j][i])
            results.append(result)
        return results
    else:
        return os.path.join(*args)
jp = joinpath

def splitpath(path):
    """splitpath(path), Split path into [dir, file, ext]. e.g., file=easyshell, ext=.py
    Supports vectorization."""

    # os.path.split only returns (head, tail), so do not use

    # vectorization
    if type(path) in [list, tuple]:
        dirs = []; files = []; exts = []
        for p in path:
            [dir, file, ext] = splitpath(p)
            dirs.append(dir); files.append(file); exts.append(ext)
        return [dirs, files, exts]

    path = path
    path = fullpath(path)
    dir = os.path.dirname(path)
    ext = os.path.splitext(path)[1]
    file = os.path.basename(path)[0:-len(ext)] if len(ext) != 0 else os.path.basename(path)
    return [dir, file, ext]
sp = splitpath

def extpath(path,newext):
    """
    path: 'a/b/c.txt', or a list of path
    newext: '.md' or a list of newext
    returns new filepath with changed extension: 'a/b/c.md'
    Supports vectorization.
    """
    if type(path) in [list, tuple]:
        [pth, fname, ext] = splitpath(path)
        result = []
        for p,f,e in zip(pth,fname,newext):
            result.apend(joinpath(p,f+newext))
        return result

    [pth, fname, ext] = splitpath(path)
    return joinpath(pth,fname+newext)
ep=extpath

def cleanpath(path,allow_unicode=False):
    """
    clean file name only (not whole path!, but takes whole path as parameter)

    Removal or replacement of CGI escaped ASCII characters, i.e. %20 becomes " ".
    Convert to ASCII if 'allow_unicode' is False. 
    Remove anything that is not an alphanumeric, ".", " ", "_" and "-"
    Convert consecutive " ", "_" and "-" to single one.
    Strip leading and trailing " ", "_" and "-".
    """
    # inspired by
    # https://github.com/django/django/blob/main/django/utils/text.py  
    # slugify get_valid_filename
    # http://detox.sourceforge.net/
    import unicodedata, urllib.parse

    [pth, fname, ext] = splitpath(path)
    fname = urllib.parse.unquote_plus(fname)

    if allow_unicode:
        fname = unicodedata.normalize('NFKC', fname)
    else:
        fname = unicodedata.normalize('NFKD', fname).encode('ascii', 'ignore').decode('ascii')
    
    # (?u) switch on unicode
    fname = re.sub(r'(?u)[^\w.\s-]', '', fname)
    expression = '(?<=[(%s)])(%s)*|^(%s)+|(%s)+$' % ('\s|_|\-','\s|_|\-','\s|_|\-','\s|_|\-')
    fname = re.sub(expression, "", fname, count=0)
    if fname in {'', '.', '..'}:
        raise Exception(f"Could not clean path '{path}'")
    path = joinpath(pth,fname+ext)
    return path
#mnemonic: x = clean
xp=cleanpath

def trim(s, how=4, chars=None):
    """Merge multiple spaces to single space in the middle, and remove trailing/leading spaces
    trim(s, how=4 [,chars])
        s: a string 
        how: a num 1=left only; 
                   2=right only; 
                   3=left and right; 
                   4 (default)=left and right and merge middle
        chars: if not given (default), space, horizontal tab, line feed, carriage return
               if given and not None, remove characters in chars instead. e.g. \s|_|\-
    eg, "Hi        buddy        what's up    Bro"  --> "Hi buddy what's up bro"
        trim(s,4,'\nx')
        " Hi        buddy        what's up    Bro\nx" --> " Hi        buddy        what's up    Bro"
    """
    if (how==1):
        s = str.lstrip(s,chars)
    elif (how==2):
        s = str.rstrip(s,chars)
    elif (how==3):
        s = str.strip(s,chars)
    elif (how==4):
        if chars is None: chars='\s|\t|\r|\n|\r\n'
        expression = '(?<=[(%s)])(%s)*|^(%s)+|(%s)+$' % (chars,chars,chars,chars)
        # http://stackoverflow.com/a/25734388/2292993
        s = re.sub(expression, "", s, count=0)
    return(s)

def quote(string):
    """ add a double quote around a single string/number/logic
    Returns a string: eg, '"3"', '"abc"', regardless of whether the string parameter is passed with 'abc', or "abc"
    """
    # https://stackoverflow.com/a/20056615/2292993
    return '"{}"'.format(string)

def join(sep='',*args):
    """glue together strings/array with sep
    1) join('','string1','string2',...)
    2) Input only supports vectorization
            join('|',['string1','string2'],['string3','string4'],...) -> [string1|string3,string2|string4]
    3) numeric/mixed array works too (the built-in '.'.join() only supports string array)
            join(',',['this','is',1,'array'])
            join('_',[1,2,2]) >> 1_2_2
    """
    # len(args) only refers to *args itself
    # in the case of (sep), len(args) == 0
    assert len(args) >= 1, "Give me more inputs"
    if len(args) == 1:
        theList = args[0]
        theList = [str(x) for x in theList]
        return sep.join(theList)
    else:        
        vectorization = False
        # to see whether there is an list/tuple input
        for arg in args:
            if type(arg) in [list, tuple]:
                vectorization = True
                length = len(arg)
                break
        if vectorization:
            # fill arg which is a string
            # args is by default an unmutable tuple
            args = list(args)
            for j in range(0,len(args)):
                arg = args[j]
                if type(arg) in [str]:
                    args[j] = [arg]*length
    
            results = []
            for i in range(0,length):
                # a bit tricy because of result + = result and .join
                result = args[0][i]
                for j in range(1,len(args)):
                    result = sep.join([result, args[j][i]])
                results.append(result)
            return results
        else:
            return sep.join(args)

def replace(theList, theItem, replacement):
    """replace(theList, theItem, replacement)
    replace all occurance of theItem in theList with replacement
    theItem could be a value, e.g., 3, 'cat' or a condition '>3' '!="cat"' or a (lambda) function
    supports numeric list, string list, or mixed
    Be careful when theItem is condition and theList is mixed (e.g., 'cat'>3 is true)
    
    Returns the changed list but the passed-in list is NOT changed.
    
    e.g., 
    replace([0,-1,1],'<0',0)
    replace([0,-1,1],-1,0)
    replace([0,-1,1],'!=0',0)
    replace([0,-1,1],'-1',1)   # <-- will not replace, because -1 and '-1' not the same
    replace(['c','a','t','cat'],'!="cat"','cat')  #<--notice the quotes around "cat", otherwise report error
    replace([0,-1,1],lambda x: x<0,0)
    """
    # http://www.python-course.eu/deep_copy.php
    from copy import deepcopy
    theList = deepcopy(theList)
    # check if it is a function
    # http://stackoverflow.com/questions/624926/how-to-detect-whether-a-python-variable-is-a-function
    # http://stackoverflow.com/questions/3655842/how-to-test-whether-a-variable-holds-a-lambda
    # http://www.diveintopython.net/power_of_introspection/lambda_functions.html
    if hasattr(theItem,'__call__'):
        for index, item in enumerate(theList):
            if theItem(item):   # pass item to the function
                theList[index] = replacement
    else:
        if type(theItem) not in [str]:
            cnd = ' == theItem'
        else:
            if theItem[0] in ['<','=','!','>']:   # startswith
                cnd = theItem
            else:
                cnd = ' == theItem'
        for index, item in enumerate(theList):
            if eval('item' + cnd):
                theList[index] = replacement
    return theList

def remove(theList, theItem):
    """remove(theList, theItem)
    remove all occurance of theItem in theList
    theItem could be a value, e.g., 3, 'cat' or a condition '>3' '!="cat"' or a (lambda) function
    supports numeric list, string list, or mixed
    Be careful when theItem is condition and theList is mixed (e.g., 'cat'>3 is true)
    
    Returns the changed list but the passed-in list is NOT changed.
    
    e.g., 
    ([0,-1,1],'<0')
    ([0,-1,1],-1)
    ([0,-1,1],'!=0')
    ([0,-1,1],'-1')   # <-- will not remove, because -1 and '-1' not the same
    (['c','a','t','cat'],'!="cat"')  #<--notice the quotes around "cat", otherwise report error
    ([0,-1,1],lambda x: x<0)
    """
    # http://www.python-course.eu/deep_copy.php
    from copy import deepcopy
    theList = deepcopy(theList)
    # check if it is a function
    # http://stackoverflow.com/questions/624926/how-to-detect-whether-a-python-variable-is-a-function
    # http://stackoverflow.com/questions/3655842/how-to-test-whether-a-variable-holds-a-lambda
    # http://www.diveintopython.net/power_of_introspection/lambda_functions.html
    
    # list remove,pop can cause index shift
    # http://stackoverflow.com/questions/497426/deleting-multiple-elements-from-a-list    
    if hasattr(theItem,'__call__'):
        for index, item in enumerate(theList):
            if theItem(item):   # pass item to the function
                theList[index] = '!!!'
    else:
        if type(theItem) not in [str]:
            cnd = ' == theItem'
        else:
            if theItem[0] in ['<','=','!','>']:   # startswith
                cnd = theItem
            else:
                cnd = ' == theItem'
        for index, item in enumerate(theList):
            if eval('item' + cnd):
                theList[index] = '!!!'
    
    for i in range(0,theList.count('!!!')):
        theList.remove('!!!') # remove
    return theList
    
def sort(*args, **kwargs):
    """wrapper of sorted, passed in list does not change, returns a sorted list"""
    return sorted(*args, **kwargs)

def cd(path=None):
    """cd(path), Changes to a new directory."""
    if path is None: path = csd()
    path = fullpath(path)
    os.chdir(path)
    print("Start working in " + os.getcwd())

def ce():
    """Changes to working directory in Finder"""
    applescript = '''
    tell app "Finder" to return the POSIX path of (target of window 1 as alias)
    '''
    def myesp(cmdString):
        import os, inspect, tempfile, subprocess
        caller = inspect.currentframe().f_back
        cmd =  cmdString % caller.f_locals
        
        fd, path = tempfile.mkstemp(suffix='.applescript')
        res = os.getcwd()
        try:
            with os.fdopen(fd, 'w') as tmp:
                tmp.write(cmd.replace('"','\"').replace("'","\'")+'\n\n')
            # res = subprocess.Popen('osascript ' + path, shell=True, executable="/bin/bash", stdout=subprocess.PIPE)
            # res.stdout.read()
            # https://stackoverflow.com/questions/41171791
            res = subprocess.run('osascript ' + path, shell=True, executable="/bin/bash", capture_output=True, text=True)
        finally:
            os.remove(path)
        return res.stdout.strip("\n")
    path = myesp(applescript)
    os.chdir(path)
    print("Start working in " + os.getcwd())

# def cf():
#     """Changes to working directory in Pather Finder"""
#     applescript = '''
#     tell app "Path Finder" to return the POSIX path of the target of the front finder window
#     '''
#     def myesp(cmdString):
#         import os, inspect, tempfile, subprocess
#         caller = inspect.currentframe().f_back
#         cmd =  cmdString % caller.f_locals
        
#         fd, path = tempfile.mkstemp(suffix='.applescript')
#         res = os.getcwd()
#         try:
#             with os.fdopen(fd, 'w') as tmp:
#                 tmp.write(cmd.replace('"','\"').replace("'","\'")+'\n\n')
#             # res = subprocess.Popen('osascript ' + path, shell=True, executable="/bin/bash", stdout=subprocess.PIPE)
#             # res.stdout.read()
#             # https://stackoverflow.com/questions/41171791
#             res = subprocess.run('osascript ' + path, shell=True, executable="/bin/bash", capture_output=True, text=True)
#         finally:
#             os.remove(path)
#         return res.stdout.strip("\n")
#     path = myesp(applescript)
#     os.chdir(path)
#     print("Start working in " + os.getcwd())

def cf():
    """Changes to working directory in QSpace"""
    cmd = '''
    item=$(osascript -e 'tell app "QSpace" to return the url of selected items')
    if [[ "$item" == "" ]]; then item=$(osascript -e 'tell app "QSpace" to return the url of root item'); fi
    if [[ -f "$item" ]]; then item=$( dirname "${item}" ); fi
    [[ -d "$item" ]] && echo "${item}"
    '''
    path = execute0(cmd,verbose=0)[0]
    os.chdir(path)
    print("Start working in " + os.getcwd())

def ls(path="./", regex=".*", full=True, dotfile=False, sort=True, case=True):
    """ls([path[, regex]], full=True, dotfile=False, sort=True)    # Returns a list of all (including hidden) files with their full paths in path, filtered by regular expression.
    case: if True, get ['Ant', 'Bat', 'Cat', 'Goat', 'Lion', 'ant', 'bat', 'cat']
          if false, get ['ant', 'Ant', 'bat', 'Bat', 'cat', 'Cat', 'Goat', 'Lion']
    to sort by modification time: when get the results, call results.sort(key=os.path.getmtime)
    """
    def _FilterList(flist, pattern_regex):
        # match_pattern = re.compile(pattern_regex, re.IGNORECASE).search
        match_pattern = re.compile(pattern_regex).search
        return list(filter(match_pattern, flist))
    def _ListingParse(path="./", pattern_regex=".*"):
        # ls() or ls('homebrew'), ls("homebrew", "\.py$")
        if os.path.isdir(path):
            path = path
            pattern_regex = pattern_regex
        # ls('*.py'), ls('homebrew/*.py')
        else:
            path, pattern = os.path.split(path)
            pattern_regex = fnmatch.translate(pattern)
        return path, pattern_regex

    path = fullpath(path)
    pattern_regex = regex
    path, pattern_regex = _ListingParse(path, pattern_regex)

    files = _FilterList(os.listdir(path), pattern_regex)
    if not dotfile: files = _FilterList(files,'^[^\.]')
    if full: 
        files = [os.path.join(path,file) for file in files]
        result = [file for file in files if os.path.isfile(file)]
    else:
        result = [file for file in files if os.path.isfile(os.path.join(path,file))]
    if sort: 
        if case:
            return sorted(result)
        else:
            # https://stackoverflow.com/questions/13954841/sort-list-of-strings-ignoring-upper-lower-case
            return sorted(result, key=lambda v: v.upper())
    else:
        return result

def lsd(path="./", regex=".*", full=False, dotfolder=False, sort=True, case=True):
    """lsd([path[, regex]], full=False, dotfolder=False, sort=True), Returns a list of all (including hidden) folders with their (optionally) full paths in path, filtered by regular expression.
    case: if True, get ['Ant', 'Bat', 'Cat', 'Goat', 'Lion', 'ant', 'bat', 'cat']
          if false, get ['ant', 'Ant', 'bat', 'Bat', 'cat', 'Cat', 'Goat', 'Lion']
    to sort by modification time: when get the results, call results.sort(key=os.path.getmtime)          
    """
    def _FilterList(flist, pattern_regex):
        # match_pattern = re.compile(pattern_regex, re.IGNORECASE).search
        match_pattern = re.compile(pattern_regex).search
        return list(filter(match_pattern, flist))
    def _ListingParse(path="./", pattern_regex=".*"):
        # ls() or ls('homebrew'), ls("homebrew", "\.py$")
        if os.path.isdir(path):
            path = path
            pattern_regex = pattern_regex
        # ls('*.py'), ls('homebrew/*.py')
        else:
            path, pattern = os.path.split(path)
            pattern_regex = fnmatch.translate(pattern)
        return path, pattern_regex

    path = fullpath(path)
    pattern_regex = regex
    path, pattern_regex = _ListingParse(path, pattern_regex)

    # os.listdir() bug on Redhat (maybe all linux) personally came across on Mon, Feb 05 2018
    # os.listdir() also list files if the dir is not working directory
    # so change working directory and change back
    olddir = os.getcwd()
    os.chdir(path)
    files = _FilterList(os.listdir(path), pattern_regex)
    if not dotfolder: files = _FilterList(files,'^[^\.]')
    if full: files = [os.path.join(path,file) for file in files]
    result = [file for file in files if not os.path.isfile(file)]
    os.chdir(olddir)
    if sort: 
        if case:
            return sorted(result)
        else:
            # https://stackoverflow.com/questions/13954841/sort-list-of-strings-ignoring-upper-lower-case
            return sorted(result, key=lambda v: v.upper())
    else:
        return result

def fls(path="./", regex=".*", dotf=False, sort=True, case=True):
    """fls([path[, regex=".*", dotf=False], sort=True])   # Returns a list of files with their full paths in flattened path (i.e. walk each subdirectory).
    case: if True, get ['Ant', 'Bat', 'Cat', 'Goat', 'Lion', 'ant', 'bat', 'cat']
          if false, get ['ant', 'Ant', 'bat', 'Bat', 'cat', 'Cat', 'Goat', 'Lion']
    """
    def _FilterList(flist, pattern_regex):
        # match_pattern = re.compile(pattern_regex, re.IGNORECASE).search
        match_pattern = re.compile(pattern_regex).search
        return list(filter(match_pattern, flist))
    def _ListingParse(path="./", pattern_regex=".*"):
        # ls() or ls('homebrew'), ls("homebrew", "\.py$")
        if os.path.isdir(path):
            path = path
            pattern_regex = pattern_regex
        # ls('*.py'), ls('homebrew/*.py')
        else:
            path, pattern = os.path.split(path)
            pattern_regex = fnmatch.translate(pattern)
        return path, pattern_regex

    path = fullpath(path)
    pattern_regex = regex
    path, pattern_regex = _ListingParse(path, pattern_regex)

    flatFiles = []
    for root, dirs, files in os.walk(path):
        # skip any folder whose name starts with . and its any subfolders
        if not dotf:
            # operationally, check if the path root has /.folder
            parts = root.split(os.sep)
            dotfolders = [part for part in parts if len(part)>=2 and part.startswith('.') and part!='..']
            if len(dotfolders)>0: continue

        filteredFiles = _FilterList(files, pattern_regex)
        if not dotf: filteredFiles = _FilterList(filteredFiles,'^[^\.]')
        ## If there is matched
        if filteredFiles:
            for file in filteredFiles:
                fileFullName = os.path.join(root, file)
                flatFiles.append(fileFullName)
    result = flatFiles
    if sort: 
        if case:
            return sorted(result)
        else:
            # https://stackoverflow.com/questions/13954841/sort-list-of-strings-ignoring-upper-lower-case
            return sorted(result, key=lambda v: v.upper())
    else:
        return result

def mkdir(path):
    """mkdir("path/to/a/directory") , Makes a directory (also any one of the "path", "to", "a" directories if not exits)."""
    path = fullpath(path)
    if not os.path.exists(path):
        os.makedirs(path)
        print("Created: " + path)

def exists(path):
    """Returns the existence of path (0 or 1, supports wildcards such as "../homebrew/*.pyc" but not regular expression)."""
    path = fullpath(path)
    paths = glob.glob(path)
    return True if paths else False

def rn(*args):
    """Rename a file or folder.
    rn(old, new)
    to parent folder must exist already; otherwise error
    old and new cannot be the same, otherwise error
    note: input could be a list/tuple, vectorization supported
    in case new name exists
          if old and new both folders, move old to new as subfolder
          if old and new both files, overwrite the new file with old file without prompt
    rn('a','b')-->rename folder a to folder b
    """
    vectorization = False
    # to see whether there is an list/tuple input
    for arg in args:
        if type(arg) in [list, tuple]:
            vectorization = True
            length = len(arg)
            break
    if vectorization:
        # fill arg which is a string
        # args is by default an unmutable tuple
        args = list(args)
        for j in range(0,len(args)):
            arg = args[j]
            if type(arg) in [str]:
                args[j] = [arg]*length

        for i in range(0,length):
            rn(args[0][i],args[1][i])
        return

    source = fullpath(args[0])
    destination = fullpath(args[1])
    if source == destination: raise Exception('Cannot rename to the same name')

    ext = os.path.splitext(destination)[1]
    # if destination dirlike
    if ext == '':
        # if destination dir exist
        if os.path.isdir(destination):
            # raise Exception('Cannot rename to exising folder name')
            mv(source,destination)
        else:
            os.rename(source,destination)
    else:
        if not os.path.isdir(os.path.dirname(destination)):
            raise Exception('Destination parent folder does not exist')
        else:
            os.rename(source,destination)

    # if os.path.isfile(source):
    #     os.rename(source, destination)
    #     print "Renamed file: " + "->".join([source, destination])
    # elif os.path.isdir(source):
    #     os.rename(source, destination)
    #     print "Renamed folder: " + "->".join([source, destination])
    # else:
    #     print source + " not renamed to " + destination

def trash(path):
    """Move a file or folder to trash. Requires: pip install Send2Trash
    trash(path)
    removes a folder recursively or a file; file supports wildcards
    if a path does not exist, nothing happens.
    note: input could be a list/tuple, vectorization supported
    """

    # vectorization
    if type(path) in [list,tuple]:
        for p in path:
            trash(p)
        return

    if not path: return  # trash('')
    path = fullpath(path)
    paths = glob.glob(path)
    if not paths: return # wildcard found nothing to remove

    from send2trash import send2trash

    for path in paths:
        if os.path.isfile(path):
            send2trash(path)
            # print "Trashed file: " + path
        elif os.path.isdir(path):
            send2trash(path)
            # print "Trashed folder: " + path
        # else:
            # print path + " not trashed"

def rm(path):
    """Deletes a file or folder.
    rm(path)
    removes a folder recursively or a file; file supports wildcards
    if a path does not exist, nothing happens.
    note: input could be a list/tuple, vectorization supported
    """
    # vectorization
    if type(path) in [list,tuple]:
        for p in path:
            rm(p)
        return

    if not path: return  # rm('')
    path = fullpath(path)
    paths = glob.glob(path)
    if not paths: return # wildcard found nothing to remove

    for path in paths:
        if os.path.isfile(path):
            os.remove(path)
            # print "Removed file: " + path
        elif os.path.isdir(path):
            if os.path.islink(path):
                # remove symbolic link to folder
                os.remove(path)
            else:
                shutil.rmtree(path)
            # print "Removed folder: " + path
        # else:
            # print path + " not removed"

def cp(source, destination, ignores=None, debugMode=False):
    """
    Copies a source file or folder to destination.
    accepts ignore, which is a list of ignore_patterns, supports wildcards (not regex)
    e.g. ['*.pyc', 'tmp*']
    This will copy everything except .pyc files and files or directories whose name starts with tmp.
    other examples ['*.py', '*.sh', 'specificfile.file'],       ['.git*'],      ['.*']

    support vectorization
    destination folder does not have to exist already
    e.g.,
    1) both works cp('a.txt','folder'), cp('a.txt','folder/b.txt')
    the former copy still has the same name 'a.txt', the latter copy new name 'b.txt'
    also cp(['a.txt','b.txt']),'folder')
    2) folder: cp('a','b')-->if b not exists, cp contents of a to b; if b exist, a becomes subfolder of b
    kinda combines rn and mv
    """
    if debugMode:
        debug_mode_in_effect = True
    else:
        if _DEBUG_MODE:
            debug_mode_in_effect = True
        else:
            debug_mode_in_effect = False

    vectorization = False
    for arg in [source,destination]:
        if type(arg) in [list, tuple]:
            vectorization = True
            length = len(arg)
            break
    if vectorization:
        args = [source,destination]
        for j in range(0,len(args)):
            arg = args[j]
            if type(arg) in [str]:
                args[j] = [arg]*length

        for i in range(0,length):
            cp(args[0][i],args[1][i],ignores=ignores)
        return

    if (not source) or (not destination): return  # cp('','')
    source = fullpath(source)
    destination = fullpath(destination)
    sources = glob.glob(source)
    if not sources: return  # wildcard found nothing to copy

    for source in sources:
        if os.path.isfile(source):
            if not debug_mode_in_effect:
                # prepare destination dir if not exist
                ext = os.path.splitext(destination)[1]
                # if destination dirlike
                if ext == '':
                    # if destination dir not exist
                    if not os.path.isdir(destination):
                        os.makedirs(destination)
                else:
                    destDir = os.path.dirname(destination)
                    if not os.path.isdir(destDir):
                        os.makedirs(destDir)

                shutil.copy(source, destination)
                # print "Copied file: " + "->".join([source, destination])
            else:
                pprint("Simulation! Copied file: " + "->".join([source, destination]), 'yellow')
        elif os.path.isdir(source):
            if not debug_mode_in_effect:
                # if destination dir not exist, shutil.copytree will auto create
                # shutil.copytree will raise error if destination dir exists
                if os.path.isdir(destination):
                    destination = os.path.join(destination,os.path.basename(source))
                if ignores:
                    shutil.copytree(source, destination, ignore=shutil.ignore_patterns(*ignores))
                else:
                    shutil.copytree(source, destination)
                # print "Copied folder: " + "->".join([source, destination])
            else:
                pprint("Simulation! Copied folder: " + "->".join([source, destination]), 'yellow')
        # else:
        #     print source + " not copied to " + destination

def mv(source, destination, debugMode=False):
    """Moves a source file or folder to destination.
    support vectorization
    destination parent folder does not have to exist already
    mv('a.txt','folder'), mv('a.txt','folder/a.txt'), mv('a.txt','folder/b.txt')
    mv('a','b')-->get b/a, b now has a as subfolder, regardless of b exists or not
                  use ez.rn('a','b') to change name a->b
    """
    if debugMode:
        debug_mode_in_effect = True
    else:
        if _DEBUG_MODE:
            debug_mode_in_effect = True
        else:
            debug_mode_in_effect = False

    vectorization = False
    for arg in [source,destination]:
        if type(arg) in [list, tuple]:
            vectorization = True
            length = len(arg)
            break
    if vectorization:
        args = [source,destination]
        for j in range(0,len(args)):
            arg = args[j]
            if type(arg) in [str]:
                args[j] = [arg]*length

        for i in range(0,length):
            mv(args[0][i],args[1][i])
        return

    if (not source) or (not destination): return  # mv('','')
    source = fullpath(source)
    destination = fullpath(destination)
    sources = glob.glob(source)
    if not sources: return  # wildcard found nothing to copy

    # prepare destination dir if not exist
    ext = os.path.splitext(destination)[1]
    # if destination dirlike
    if ext == '':
        # if destination dir not exist
        if not os.path.isdir(destination):
            os.makedirs(destination)
    else:
        destDir = os.path.dirname(destination)
        if not os.path.isdir(destDir):
            os.makedirs(destDir)

    for source in sources:
        if not debug_mode_in_effect:
            shutil.move(source, destination)
            # print "Moved: " + "->".join([source, destination])
        elif debug_mode_in_effect:
            pprint("Simulation! Moved: " + "->".join([source, destination]), 'yellow')
        # else:
        #     print source + " not moved to " + destination

def lns(source, destination):
    """Creates a soft link/symbolic link/shortcut."""
    source = fullpath(source)
    destination = fullpath(destination)

    os.symlink(source, destination)
    print("Symbolic link: " + "->".join([source, destination]))

def execute0(cmd, verbose=3, save=None, saveMode='a', redirect=None, redirectMode='a', shell='bash', debugMode=False, *args, **kwargs):
    """Executes a bash command. can capture output
    verbose: any screen display here does not affect returned values
            0 = nothing to display
            1 = only the actual command
            2 = only the command output
            3 = both the command itself and output
    save: None, or a file path to save the cmd (shebang prepended), can still save even if error occurs (for debugging) or simulation
    saveMode: 'a' (append) or 'w' (overwrite), ignored if save=None.
    redirect: None, or a file path to save the redirected cmd execution output. compatible with logon(); works also cmd itself has redirection (eg, tee)
    redirectMode: 'a' (append) or 'w' (overwrite), ignored if redirect=None.
    shell: 'bash', 'tcsh', 'sh' (internally converted to '/bin/'+shell)
    return: ...regardless of verbose...
            returns shell output as a list with each elment is a line of string (whitespace stripped both sides) from output
            if error occurs, return None, also always print out the error message to screen
            if no output or all empty output, return [] 
               note execute0('printf "\n\n"')-->[]; but execute0('printf "\n\n3"')-->['', '', '3']
    note: if use this function interactively, one can return _ = execute0() to a dummy variable
          alternatively, in ipython, execute0(); (add semicolon) to suppress the returned contents
          seems to recognize cmd='echo $PATH', but not alias in .bash_profile
    """
    if debugMode:
        debug_mode_in_effect = True
    else:
        if _DEBUG_MODE:
            debug_mode_in_effect = True
        else:
            debug_mode_in_effect = False
    if not debug_mode_in_effect:
        if verbose in [1,3]: pprint("Command: " + cmd + "\n> > > > > > > > > > > > > > > > > > > > > > > > > > > > > > > > ")

        # https://askubuntu.com/a/731237
        if redirect:
            if redirectMode == 'a':
                if shell in ['tcsh']:
                    # cmdSuffix = ' |& tee -a "' + redirect + '"'
                    cmdSuffix = ' 2>&1 | tee -a "' + redirect + '"'
                else:
                    cmdSuffix = ' 2>&1 | tee -a "' + redirect + '"'
            else:
                if shell in ['tcsh']:
                    # cmdSuffix = ' |& tee "' + redirect + '"'
                    cmdSuffix = ' 2>&1 | tee "' + redirect + '"'
                else:
                    cmdSuffix = ' 2>&1 | tee "' + redirect + '"'
        else:
            cmdSuffix = ''

        if saveMode=='w': rm(save)
        if save:
            if os.path.exists(save):
                with open(save, 'a') as tmp:
                    tmp.write(cmd.replace('"','\"').replace("'","\'")+'\n\n') 
            else:
                with open(save, 'a') as tmp:
                    tmp.write('#!/bin/'+('tcsh -xef' if shell in ['tcsh'] else shell)+'\n\n'+cmd.replace('"','\"').replace("'","\'")+'\n\n')
            subprocess.call('chmod +x '+save, shell=True)
            print(('\nCommand saved at '+save+'\n'))

        # https://stackoverflow.com/a/40139101/2292993
        def _execute_cmd(cmd):
            if os.name == 'nt' or platform.system() == 'Windows':
                # set stdin, out, err all to PIPE to get results (other than None) after run the Popen() instance
                # shell=True can do shell pipes, filename wildcards, environment variable expansion, and expansion of ~
                import tempfile
                tmpfd, tmpPath = tempfile.mkstemp(suffix='.bat')
                with os.fdopen(tmpfd, 'w') as tmp:
                    tmp.write('@echo off'+'\n\n'+cmd.replace('"','\"').replace("'","\'")+'\n\n')                
                p = subprocess.Popen(tmpPath, stdout=subprocess.PIPE, stderr=subprocess.PIPE, shell=True)
            else:
                # tcsh -xef
                # -x : echo commands to terminal before executing them
                # -e : terminate script when encountering any error
                # -f : do not process user's ~/.cshrc file
                # for subprocess(), if cmd is a string, set shell=True; if a list, set shell=False
                # also, The executable argument specifies a replacement program to execute. It is very seldom needed. 
                # If shell=True, on Unix the executable argument specifies a replacement shell for the default /bin/sh
                # this shell is not effected by the actual shell terminal used when execute python
                # when using executable, i.e., subprocess.Popen(executable="/bin/"+shell), you cannot pass arg to it. that's why I use prefix+cmd to implement tcsh -xef
                # p = subprocess.Popen('/bin/'+('tcsh -xef' if shell in ['tcsh'] else shell)+' -c "'+cmd+'"'+cmdSuffix, stdout=subprocess.PIPE, stderr=subprocess.PIPE, shell=True)

                # well, -c seems not to be able to handle $var: eg, error for /bin/bash -c "x=33333333333333; echo $x"
                # another issue with -c is sometimes "Argument list too long" "getconf ARG_MAX"
                # save the command to a temp file, then execute
                # the temp file solution can bypass both above issues
                import tempfile
                tmpfd, tmpPath = tempfile.mkstemp(suffix='.sh')
                with os.fdopen(tmpfd, 'w') as tmp:
                    tmp.write('#!/bin/'+('tcsh -xef' if shell in ['tcsh'] else shell)+'\n\n'+cmd.replace('"','\"').replace("'","\'")+'\n\n')
                # p = subprocess.Popen('/bin/'+('tcsh -xef' if shell in ['tcsh'] else shell)+' "'+tmpPath+'"'+cmdSuffix, stdout=subprocess.PIPE, stderr=subprocess.PIPE, shell=True, executable="/bin/bash")
                # some command line tools (e.g., ffmpeg) direct their outputs as stderr; so I simply merge stderr to stdout
                p = subprocess.Popen('/bin/'+('tcsh -xef' if shell in ['tcsh'] else shell)+' "'+tmpPath+'"'+cmdSuffix, stdout=subprocess.PIPE, stderr=subprocess.STDOUT, shell=True, executable="/bin/bash")
            
            # rationale: I capture everything, and then I decide what/if to manually print to screen

            # the Popen() instance starts running once instantiated (??)
            # additionally, communicate(), or poll() and wait process to terminate
            # communicate() accepts optional input as stdin to the pipe (requires setting stdin=subprocess.PIPE above), return out, err as tuple
            # if communicate(), the results are buffered in memory, not good for realtime display on screen
            
            # Read stdout from subprocess until the buffer is empty !
            # if error occurs, the stdout is '', which means the below loop is essentially skipped
            # A prefix of 'b' or 'B' is ignored in Python 2; 
            # it indicates that the literal should become a bytes literal in Python 3 
            # (e.g. when code is automatically converted with 2to3).
            # return iter(p.stdout.readline, b'')
            for line in iter(p.stdout.readline, b''):
                # # Windows has \r\n, Unix has \n, Old mac has \r
                # if line not in ['','\n','\r','\r\n']: # Don't print blank lines
                yield line.decode('utf-8')
            while p.poll() is None:                                                                                                                                        
                sleep(.1) #Don't waste CPU-cycles
            
            # err = p.stderr.read() 
            # if p.returncode != 0:
            #     # responsible for logging STDERR 
            #     if verbose in [2,3]: print("Error: " + err.decode('utf-8'))
            #     yield None
            
            # test: cmd="0/0"
            # err = p.stdout.read()  # stdout.read() was read earlier but now is empty
            if p.returncode != 0:
                # responsible for logging STDERR 
                if verbose in [2,3]: pprint("Error occured",color='red')
                yield None
            # delete temp file
            try:
                os.remove(tmpPath)
            except:
                pass

        out = []
        for line in _execute_cmd(cmd):
            # error did not occur earlier
            if line is not None:
                # print() default end is '\n'; to avoid a newline (by print itself) being printed
                if verbose in [2,3]: print(line, end='')
                out.append(line.strip())
            else:
                # error occured earlier
                out = None

        # post-process for returning value
        if out is None:
            return None
        else:
            # https://stackoverflow.com/a/3845453/2292993
            # filter() check if one or all empty string
            if len(out)==0 or len([_f for _f in out if _f])==0:
                return []
            else:
                return out
    else:
        pprint("Simulation! Execute command: " + cmd + "\n< < < < < < < < < < < < < < < < < < < < < < < < < < < < < < < < ", 'yellow')
        if saveMode=='w': rm(save)
        if save:
            if os.path.exists(save):
                with open(save, 'a') as tmp:
                    tmp.write(cmd.replace('"','\"').replace("'","\'")+'\n\n') 
            else:
                with open(save, 'a') as tmp:
                    tmp.write('#!/bin/'+('tcsh -xef' if shell in ['tcsh'] else shell)+'\n\n'+cmd.replace('"','\"').replace("'","\'")+'\n\n')
            subprocess.call('chmod +x '+save, shell=True)
            print(('\nCommand saved at '+save+'\n'))
        return None

def execute(cmd, verbose=3, save=None, saveMode='a', redirect=None, redirectMode='a', shell='bash', debugMode=False, *args, **kwargs):
    """
    a wrapper of subprocess.call, cannot capture output, does not return the output to a python variable
    Executes a shell command.
    verbose: any screen display here does not affect returned values
            0 = nothing to display
            1 = only the actual command
            2 = only the command output
            3 = both the command itself and output
    save: None, or a file path to save the cmd (shebang prepended), can still save even if error occurs (for debugging)
    saveMode: 'a' (append) or 'w' (overwrite), ignored if save=None.
    redirect: None, or a file path to save the redirected cmd execution output. compatible with logon(); works also cmd itself has redirection (eg, tee)
    redirectMode: 'a' (append) or 'w' (overwrite), ignored if redirect=None.
    note: seems to recognize cmd='echo $PATH', but not alias in .bash_profile
    """
    if debugMode:
        debug_mode_in_effect = True
    else:
        if _DEBUG_MODE:
            debug_mode_in_effect = True
        else:
            debug_mode_in_effect = False

    if not debug_mode_in_effect:
        if verbose in [1,3]: pprint("Command: " + cmd + "\n> > > > > > > > > > > > > > > > > > > > > > > > > > > > > > > > ")
        if verbose in [0,1]: 
            output=False
        else:
            output=True

        # https://askubuntu.com/a/731237
        if redirect:
            if redirectMode == 'a':
                if shell in ['tcsh']:
                    # cmdSuffix = ' |& tee -a "' + redirect + '"'
                    cmdSuffix = ' 2>&1 | tee -a "' + redirect + '"'
                else:
                    cmdSuffix = ' 2>&1 | tee -a "' + redirect + '"'
            else:
                if shell in ['tcsh']:
                    # cmdSuffix = ' |& tee "' + redirect + '"'
                    cmdSuffix = ' 2>&1 | tee "' + redirect + '"'
                else:
                    cmdSuffix = ' 2>&1 | tee "' + redirect + '"'
        else:
            cmdSuffix = ''

        # save cmd
        if saveMode=='w': rm(save)
        if save:
            if os.path.exists(save):
                with open(save, 'a') as tmp:
                    tmp.write(cmd.replace('"','\"').replace("'","\'")+'\n\n') 
            else:
                with open(save, 'a') as tmp:
                    tmp.write('#!/bin/'+('tcsh -xef' if shell in ['tcsh'] else shell)+'\n\n'+cmd.replace('"','\"').replace("'","\'")+'\n\n')
            subprocess.call('chmod +x '+save, shell=True)
            print(('\nCommand saved at '+save+'\n'))

        if os.name == 'nt' or platform.system() == 'Windows':
            import tempfile
            tmpfd, tmpPath = tempfile.mkstemp(suffix='.bat')
            with os.fdopen(tmpfd, 'w') as tmp:
                tmp.write('@echo off'+'\n\n'+cmd.replace('"','\"').replace("'","\'")+'\n\n')                
            if output:
                subprocess.call(tmpPath, shell=True, stderr=subprocess.STDOUT)
            else:
                subprocess.call(tmpPath, shell=True, stdout=open(os.devnull, "w"), stderr=subprocess.STDOUT)
        else:
            if output:
                # subprocess.call('/bin/'+('tcsh -xef' if shell in ['tcsh'] else shell)+' -c "'+cmd+'"'+cmdSuffix, shell=True, executable="/bin/bash") 
                import tempfile
                tmpfd, tmpPath = tempfile.mkstemp(suffix='.sh')
                try:
                    with os.fdopen(tmpfd, 'w') as tmp:
                        tmp.write('#!/bin/'+('tcsh -xef' if shell in ['tcsh'] else shell)+'\n\n'+cmd.replace('"','\"').replace("'","\'")+'\n\n')
                    subprocess.call('/bin/'+('tcsh -xef' if shell in ['tcsh'] else shell)+' "'+tmpPath+'"'+cmdSuffix, shell=True, executable="/bin/bash", stderr=subprocess.STDOUT)
                finally:
                    os.remove(tmpPath)
            else:
                # subprocess.call('/bin/'+('tcsh -xef' if shell in ['tcsh'] else shell)+' -c "'+cmd+'"'+cmdSuffix, shell=True, executable="/bin/bash", stdout=open(os.devnull, "w"), stderr=subprocess.STDOUT)
                import tempfile
                tmpfd, tmpPath = tempfile.mkstemp(suffix='.sh')
                try:
                    with os.fdopen(tmpfd, 'w') as tmp:
                        tmp.write('#!/bin/'+('tcsh -xef' if shell in ['tcsh'] else shell)+'\n\n'+cmd.replace('"','\"').replace("'","\'")+'\n\n')
                    subprocess.call('/bin/'+('tcsh -xef' if shell in ['tcsh'] else shell)+' "'+tmpPath+'"'+cmdSuffix, shell=True, executable="/bin/bash", stdout=open(os.devnull, "w"), stderr=subprocess.STDOUT)
                finally:
                    os.remove(tmpPath)
        print("")

    else:
        pprint("Simulation! Execute command: " + cmd + "\n< < < < < < < < < < < < < < < < < < < < < < < < < < < < < < < < ", 'yellow')
        if saveMode=='w': rm(save)
        if save:
            if os.path.exists(save):
                with open(save, 'a') as tmp:
                    tmp.write(cmd.replace('"','\"').replace("'","\'")+'\n\n') 
            else:
                with open(save, 'a') as tmp:
                    tmp.write('#!/bin/'+('tcsh -xef' if shell in ['tcsh'] else shell)+'\n\n'+cmd.replace('"','\"').replace("'","\'")+'\n\n')
            subprocess.call('chmod +x '+save, shell=True)
            print(('\nCommand saved at '+save+'\n'))

def esp0(cmdString, verbose=3, save=None, saveMode='a', redirect=None, redirectMode='a', shell='bash', skipdollar=0, debugMode=False, *args, **kwargs):
    """
    Execute a SPrintf, can capture output
    a shortcut for execute0(sprintf(cmdString))
    return: ...regardless of verbose...
        returns shell output as a list with each elment is a line of string (whitespace stripped both sides) from output
        if error occurs, return None, also always print out the error message to screen
        if no output or all empty output, return [] 
           note execute0('printf "\n\n"')-->[]; but execute0('printf "\n\n3"')-->['', '', '3']
    verbose: any screen display here does not affect returned values
            0 = nothing to display
            1 = only the actual command
            2 = only the command output
            3 = both the command itself and output
    save: None, or a file path to save the cmd (shebang prepended), can still save even if error occurs (for debugging)
    saveMode: 'a' (append) or 'w' (overwrite), ignored if save=None.
    redirect: None, or a file path to save the redirected cmd execution output. compatible with logon(); works also cmd itself has redirection (eg, tee)
    redirectMode: 'a' (append) or 'w' (overwrite), ignored if redirect=None.
    if skipdollar=1 (1/0), $ (but not others) syntax will be entirely skipped, useful for R codes (df$col), or certain bash codes
    note: seems to recognize cmd='echo $PATH', but not alias in .bash_profile
    Example:
            var_to_be_in_bash = 'blabla'
            cmd = '''
echo $var_to_be_in_bash
echo "new line"
'''
            ez.esp(cmd)
            # Command: echo blabla
            # Actual output: blabla    
    """
    if debugMode:
        debug_mode_in_effect = True
    else:
        if _DEBUG_MODE:
            debug_mode_in_effect = True
        else:
            debug_mode_in_effect = False
    # # caller's caller
    # caller = inspect.currentframe().f_back.f_back
    import inspect
    caller = inspect.currentframe().f_back
    # for esp1()
    if kwargs: 
        if kwargs['insideCalling']:
            caller = inspect.currentframe().f_back.f_back
    cmd = sprintf(cmdString, caller.f_locals, skipdollar=skipdollar)
    return execute0(cmd, verbose=verbose, save=save, saveMode=saveMode, redirect=redirect, redirectMode=redirectMode, shell=shell, debugMode=debug_mode_in_effect, *args, **kwargs)


def esp(cmdString, verbose=3, save=None, saveMode='a', redirect=None, redirectMode='a', shell='bash', skipdollar=0, debugMode=False, *args, **kwargs):
    """
    a shortcut for execute(sprintf(cmdString)), cannot capture output
    Execute a SPrintf, but does not return the output to a python variable
    execute, esp (subprocess.call) seem to work better with AFNI commands, while execute1/2, esp1/2 (based on subprocess.Popen) sometimes fail
    verbose: any screen display here does not affect returned values
            0 = nothing to display
            1 = only the actual command
            2 = only the command output
            3 = both the command itself and output
    save: None, or a file path to save the cmd (shebang prepended), can still save even if error occurs (for debugging)
    saveMode: 'a' (append) or 'w' (overwrite), ignored if save=None.
    redirect: None, or a file path to save the redirected cmd execution output. compatible with logon(); works also cmd itself has redirection (eg, tee)
    redirectMode: 'a' (append) or 'w' (overwrite), ignored if redirect=None.
    if skipdollar=1 (1/0), $ (but not others) syntax will be entirely skipped, useful for R codes (df$col), or certain bash codes
    note: seems to recognize cmd='echo $PATH', but not alias in .bash_profile
    Example:
            var_to_be_in_bash = 'blabla'
            cmd = '''
echo $var_to_be_in_bash
echo "new line"
'''
            ez.esp(cmd)
            # Command: echo blabla
            # Actual output: blabla
    """
    if debugMode:
        debug_mode_in_effect = True
    else:
        if _DEBUG_MODE:
            debug_mode_in_effect = True
        else:
            debug_mode_in_effect = False
    # # caller's caller
    # caller = inspect.currentframe().f_back.f_back
    import inspect
    caller = inspect.currentframe().f_back
    cmd = sprintf(cmdString, caller.f_locals, skipdollar=skipdollar)
    execute(cmd, verbose=verbose, save=save, saveMode=saveMode, redirect=redirect, redirectMode=redirectMode, shell=shell, debugMode=debug_mode_in_effect, *args, **kwargs)

def espR0(cmdString, verbose=3, save=None, saveMode='a', redirect=None, redirectMode='a', shell='bash', skipdollar=1, debugMode=False, *args, **kwargs):
    """
    a shortcut for execute0(sprintf(cmdString)), can capture output
    write cmdString (R codes) to a temp file, then call "Rscript temp.R", finally remove the temp file
    Execute a SPrintf    
    return: ...regardless of verbose...
        returns shell output as a list with each elment is a line of string (whitespace stripped both sides) from output
        if error occurs, return None, also always print out the error message to screen
        if no output or all empty output, return [] 
           note execute0('printf "\n\n"')-->[]; but execute0('printf "\n\n3"')-->['', '', '3']
    verbose: any screen display here does not affect returned values
            0 = nothing to display
            1 = only the actual command
            2 = only the command output
            3 = both the command itself and output
    save: None, or a file path to save the cmd (shebang prepended), can still save even if error occurs (for debugging) or simulation
    saveMode: 'a' (append) or 'w' (overwrite), ignored if save=None.
    redirect: None, or a file path to save the redirected cmd execution output. compatible with logon(); works also cmd itself has redirection (eg, tee)
    redirectMode: 'a' (append) or 'w' (overwrite), ignored if redirect=None.
    if skipdollar=1 (1/0), $ (but not others) syntax will be entirely skipped, useful for R codes (df$col), or certain bash codes
    note: seems to recognize cmd='echo $PATH', but not alias in .bash_profile
    Example: 
            ez.espR('iris$Species')
            # print out iris$Species, Levels: setosa versicolor virginica    
    """
    # # caller's caller
    # caller = inspect.currentframe().f_back.f_back
    import inspect
    caller = inspect.currentframe().f_back
    # for esp()
    if kwargs: 
        if kwargs['insideCalling']:
            caller = inspect.currentframe().f_back.f_back
    cmd = sprintf(cmdString,caller.f_locals,skipdollar=skipdollar)
    
    if debugMode:
        debug_mode_in_effect = True
    else:
        if _DEBUG_MODE:
            debug_mode_in_effect = True
        else:
            debug_mode_in_effect = False
    
    if not debug_mode_in_effect:
        # save R source code
        if saveMode=='w': rm(save)
        if save:
            if os.path.exists(save):
                with open(save, 'a') as tmp:
                    tmp.write(cmd.replace('"','\"').replace("'","\'")+'\n\n') 
            else:
                with open(save, 'a') as tmp:
                    tmp.write('#!/bin/Rscript \n\n'+cmd.replace('"','\"').replace("'","\'")+'\n\n')
            print(('\nCommand saved at '+save+'\n'))
        
        import tempfile
        # create temp file with specified suffix
        fd, path = tempfile.mkstemp(suffix='.R')
        try:
            with os.fdopen(fd, 'w') as tmp:
                tmp.write('#!/bin/Rscript \n\n'+cmd.replace('"','\"').replace("'","\'")+'\n\n')
            # not save this command line
            result = execute0('Rscript --no-save --no-restore ' + path, verbose=verbose, save=None, saveMode='a', redirect=redirect, redirectMode=redirectMode, shell=shell, debugMode=debug_mode_in_effect, *args, **kwargs)
        # delete it when it is done (can still delete after return)
        # A finally clause is always executed before leaving the try statement, whether an exception has occurred or not. 
        finally:
            os.remove(path)
        return result
    else:
        pprint("Simulation! Execute command: " + cmd + "\n< < < < < < < < < < < < < < < < < < < < < < < < < < < < < < < < ", 'yellow')
        if saveMode=='w': rm(save)
        if save:
            if os.path.exists(save):
                with open(save, 'a') as tmp:
                    tmp.write(cmd.replace('"','\"').replace("'","\'")+'\n\n') 
            else:
                with open(save, 'a') as tmp:
                    tmp.write('#!/bin/Rscript \n\n'+cmd.replace('"','\"').replace("'","\'")+'\n\n')
            print(('\nCommand saved at '+save+'\n'))
        return None

def espR(cmdString, verbose=3, save=None, saveMode='a', redirect=None, redirectMode='a', shell='bash', skipdollar=1, debugMode=False, *args, **kwargs):
    """
    a shortcut for execute(sprintf(cmdString)), cannot capture output
    write cmdString (R codes) to a temp file, then call "Rscript temp.R", finally remove the temp file
    Execute a SPrintf, but does not return the output to a python variable
    cmdString: R codes
    verbose: any screen display here does not affect returned values
            0 = nothing to display
            1 = only the actual command
            2 = only the command output
            3 = both the command itself and output
    save: None, or a file path to save the cmd (shebang prepended), can still save even if error occurs (for debugging)
    saveMode: 'a' (append) or 'w' (overwrite), ignored if save=None.
    redirect: None, or a file path to save the redirected cmd execution output. compatible with logon(); works also cmd itself has redirection (eg, tee)
    redirectMode: 'a' (append) or 'w' (overwrite), ignored if redirect=None.
    if skipdollar=1 (1/0), $ (but not others) syntax will be entirely skipped, useful for R codes (df$col), or certain bash codes
    note: seems to recognize cmd='echo $PATH', but not alias in .bash_profile
    Example: 
            ez.espR('iris$Species')
            # print out iris$Species, Levels: setosa versicolor virginica
    """
    # # caller's caller
    # caller = inspect.currentframe().f_back.f_back
    import inspect
    caller = inspect.currentframe().f_back
    # for esp()
    if kwargs: 
        if kwargs['insideCalling']:
            caller = inspect.currentframe().f_back.f_back
    cmd = sprintf(cmdString,caller.f_locals,skipdollar=skipdollar)
    
    if debugMode:
        debug_mode_in_effect = True
    else:
        if _DEBUG_MODE:
            debug_mode_in_effect = True
        else:
            debug_mode_in_effect = False
    
    if not debug_mode_in_effect:
        # save R source code
        if saveMode=='w': rm(save)
        if save:
            if os.path.exists(save):
                with open(save, 'a') as tmp:
                    tmp.write(cmd.replace('"','\"').replace("'","\'")+'\n\n') 
            else:
                with open(save, 'a') as tmp:
                    tmp.write('#!/bin/Rscript \n\n'+cmd.replace('"','\"').replace("'","\'")+'\n\n')
            print(('\nCommand saved at '+save+'\n'))

        import tempfile
        # create temp file with specified suffix
        fd, path = tempfile.mkstemp(suffix='.R')
        try:
            with os.fdopen(fd, 'w') as tmp:
                tmp.write('#!/bin/Rscript \n\n'+cmd.replace('"','\"').replace("'","\'")+'\n\n')
            # not save this command line
            execute('Rscript --no-save --no-restore ' + path, verbose=verbose, save=None, saveMode='a', redirect=redirect, redirectMode=redirectMode, shell=shell, debugMode=debug_mode_in_effect, *args, **kwargs)
        # delete it when it is done (can still delete after return)
        # A finally clause is always executed before leaving the try statement, whether an exception has occurred or not. 
        finally:
            os.remove(path)
    else:
        pprint("Simulation! Execute command: " + cmd + "\n< < < < < < < < < < < < < < < < < < < < < < < < < < < < < < < < ", 'yellow')
        if saveMode=='w': rm(save)
        if save:
            if os.path.exists(save):
                with open(save, 'a') as tmp:
                    tmp.write(cmd.replace('"','\"').replace("'","\'")+'\n\n') 
            else:
                with open(save, 'a') as tmp:
                    tmp.write('#!/bin/Rscript \n\n'+cmd.replace('"','\"').replace("'","\'")+'\n\n')
            print(('\nCommand saved at '+save+'\n'))
        return None

def espA(cmdString, verbose=0, save=None, saveMode='a', redirect=None, redirectMode='a', shell='bash', skipdollar=1, debugMode=False, *args, **kwargs):
    """
    typical usage (generally no need to tweak other parameters):
    key = 'm'
    cmd = '''
    tell application "System Events" to keystroke "%(key)s" using {command down}
    '''
    ez.espA(cmd)

    a shortcut for execute(sprintf(cmdString)), cannot capture output
    write cmdString (applescript codes) to a temp file, then call "osascript temp.applescript", finally remove the temp file
    Execute a SPrintf, but does not return the output to a python variable
    cmdString: applescript codes
    verbose: any screen display here does not affect returned values
            0 = nothing to display
            1 = only the actual command
            2 = only the command output
            3 = both the command itself and output
    save: None, or a file path to save the cmd, can still save even if error occurs (for debugging)
    saveMode: 'a' (append) or 'w' (overwrite), ignored if save=None.
    redirect: None, or a file path to save the redirected cmd execution output. compatible with logon(); works also cmd itself has redirection (eg, tee)
    redirectMode: 'a' (append) or 'w' (overwrite), ignored if redirect=None.
    if skipdollar=1 (1/0), $ (but not others) syntax will be entirely skipped, useful for {} in applescript such as, keystroke "m" using {command down}
    note: seems to recognize cmd='echo $PATH', but not alias in .bash_profile
    Example: 
            ez.espA('''tell application "System Events" to keystroke "m" using {command down}''')
    """

    # a simplified, standalone version
    # def espA(cmdString):
    #     """
    #     key = 'm'
    #     cmd = '''
    #     tell application "System Events" to keystroke "%(key)s" using {command down}
    #     '''
    #     ez.espA(cmd)
    #     """
    #     import os, inspect, tempfile, subprocess
    #     caller = inspect.currentframe().f_back
    #     cmd =  cmdString % caller.f_locals
        
    #     fd, path = tempfile.mkstemp(suffix='.applescript')
    #     try:
    #         with os.fdopen(fd, 'w') as tmp:
    #             tmp.write(cmd.replace('"','\"').replace("'","\'")+'\n\n')
    #         subprocess.call('osascript ' + path, shell=True, executable="/bin/bash")
    #     finally:
    #         os.remove(path)
    #     return None


    # # caller's caller
    # caller = inspect.currentframe().f_back.f_back
    import inspect
    caller = inspect.currentframe().f_back
    # for esp()
    if kwargs: 
        if kwargs['insideCalling']:
            caller = inspect.currentframe().f_back.f_back
    cmd = sprintf(cmdString,caller.f_locals,skipdollar=skipdollar)
    
    if debugMode:
        debug_mode_in_effect = True
    else:
        if _DEBUG_MODE:
            debug_mode_in_effect = True
        else:
            debug_mode_in_effect = False
    
    if not debug_mode_in_effect:
        # save R source code
        if saveMode=='w': rm(save)
        if save:
            if os.path.exists(save):
                with open(save, 'a') as tmp:
                    tmp.write(cmd.replace('"','\"').replace("'","\'")+'\n\n') 
            else:
                with open(save, 'a') as tmp:
                    tmp.write(cmd.replace('"','\"').replace("'","\'")+'\n\n')
            print(('\nCommand saved at '+save+'\n'))

        import tempfile
        # create temp file with specified suffix
        fd, path = tempfile.mkstemp(suffix='.applescript')
        try:
            with os.fdopen(fd, 'w') as tmp:
                tmp.write(cmd.replace('"','\"').replace("'","\'")+'\n\n')
            # not save this command line
            execute('osascript ' + path, verbose=verbose, save=None, saveMode='a', redirect=redirect, redirectMode=redirectMode, shell=shell, debugMode=debug_mode_in_effect, *args, **kwargs)
        # delete it when it is done (can still delete after return)
        # A finally clause is always executed before leaving the try statement, whether an exception has occurred or not. 
        finally:
            os.remove(path)
    else:
        pprint("Simulation! Execute command: " + cmd + "\n< < < < < < < < < < < < < < < < < < < < < < < < < < < < < < < < ", 'yellow')
        if saveMode=='w': rm(save)
        if save:
            if os.path.exists(save):
                with open(save, 'a') as tmp:
                    tmp.write(cmd.replace('"','\"').replace("'","\'")+'\n\n') 
            else:
                with open(save, 'a') as tmp:
                    tmp.write(cmd.replace('"','\"').replace("'","\'")+'\n\n')
            print(('\nCommand saved at '+save+'\n'))
        return None

from contextlib import contextmanager
@contextmanager
def nooutput():
    """Temporarily suppress stdout, stderr in a context
    with nooutput():
        print 'this is will not be printed in stdout'
    """
    with open(os.devnull, "w") as devnull:
        old_stdout = sys.stdout
        sys.stdout = devnull
        old_stderr = sys.stderr
        sys.stderr = devnull
        try:  
            yield
        finally:
            sys.stdout = old_stdout
            sys.stderr = old_stderr

from pprint import pprint as ppprint
def pprint(text,color='green'):
    """
    # color print to terminal (may not work on any terminal)
    (text,color='green')
    color: string, case insentive, out of black, red, green, yellow, blue, magenta, cyan, white
    """
    # http://blog.mathieu-leplatre.info/colored-output-in-console-with-python.html
    color = color.lower()
    colors = {'black':0, 'red':1, 'green':2, 'yellow':3, 'blue':4, 'magenta':5, 'cyan':6, 'white':7}
    color = colors[color]
    def _has_colors(stream):
        if not hasattr(stream, "isatty"):
            return False
        if not stream.isatty():
            return False # auto color only on TTYs
        try:
            import curses
            curses.setupterm()
            return curses.tigetnum("colors") > 2
        except:
            # guess false in case of error
            return False
    if _has_colors(sys.stdout):
        seq = "\x1b[1;%dm" % (30+color) + text + "\x1b[0m"
        sys.stdout.write(seq+'\n')
    else:
        sys.stdout.write(text+'\n')

def beep():
    """Beeps to notify user."""
    sys.stdout.write("\a")
    sys.stdout.flush()

def where(name):
    """
    name without or with package name, i.e., ez.help, help
    where(name), Prints where a module is and in which module a function is.
    where('python') returns which python is being used and version info."""
    name_no_prefix = name.split('.')[-1]
    if name_no_prefix == 'python':
        from distutils.sysconfig import get_python_lib
        print(get_python_lib())
        print((sys.version))
    else:
        try:
            print(sys.modules[name_no_prefix])
        except:
            for module in list(sys.modules.keys()):
                try:
                    import inspect
                    caller = inspect.currentframe().f_back
                    if name_no_prefix in eval('dir('+ module + ')', caller.f_locals):
                        print('function ' + name_no_prefix + '() in ' + module)
                        # print sys.modules[module]
                        # print ''
                except:
                    pass

from shutil import which

def doc(package_prefixed_name):
    """
    name with a package name, i.e., ez.help
    help(name)/doc(name) # name is a string, Prints the doc string of a module/class/function
    when write a module, add:
    __doc__ = three double quotes blabla three double quotes         <-----this is module's docstring, use explicit

    when write a function/class:
    def function(arg):
        three double quotes Returns, blabla three double quotes      <-----this is function's doctoring, use implicit
        return sth
    """
    import inspect
    caller = inspect.currentframe().f_back
    print(eval(package_prefixed_name + '.__doc__', caller.f_locals))
help = doc

def getos():
    """
    returns Windows, Linux or Darwin
    """
    # https://stackoverflow.com/a/58071295/2292993
    return platform.system()

def ver(package_name='python'):
    """
    ver(package_name) version(package_name), see a package's version.  package_name could be 'python'
    """
    print(package_name + ' version installed:')
    if package_name == 'python' or package_name == 'ez':
        print((sys.version))
        HERE = os.path.dirname(os.path.abspath(__file__))
        with open(os.path.join(HERE,'version.py'),'r') as f:
            ezv=f.readline()
        ezv = ezv.split(" = ")[1].strip("'").strip('"')
        print(f'\nLib: {HERE}')
        print(f'ez: {ezv}')
    else:
        # https://docs.python.org/2.7/reference/simple_stmts.html#exec
        theNameSpace = {}
        exec('import ' + package_name, theNameSpace)
        print(theNameSpace[package_name].__version__)
version = ver

def evaluate(exp):
    """
    evaluate(exp), evaluate a statement or expression, combining python built-in eval() and exec()
    python eval() only works for expression--returns something, exec() for statement--not returns anything
    an expression:  x+1,      x = eval('x+1')
    a statement:    x = x+1   exec('x = x+1')
    """
    import inspect
    caller = inspect.currentframe().f_back
    try:
        return eval(exp,caller.f_locals)
    except SyntaxError:
        # https://docs.python.org/2.7/reference/simple_stmts.html#exec
        exec(exp,caller.f_locals)

def who(name=''):
    """
    whos(name),whos() list imported functions/packages
    """
    import inspect
    caller = inspect.currentframe().f_back
    theList = eval('dir('+ name + ')', caller.f_locals)
    theList = sorted(theList)
    theNameSpace = name if name != '' else 'Current'
    print(theNameSpace + ' namespace has the following existing functions/modules:\n')
    # http://stackoverflow.com/questions/19863388/modify-print-function-for-multiple-columns-python
    def pretty_print(theList, ncols):
        columns = len(theList)//200+ncols
        lines = ("".join(s.ljust(20) for s in theList[i:i+columns-1])+theList[i:i+columns][-1] for i in range(0, len(theList), columns))
        return "\n".join(lines)
    print(pretty_print(theList,4))
whos = who

# def sedawk(path, regex=".*", search=None, replace=None, recursion=True):
#     """Searches and replaces the matched content in every file with regular expression."""
#     path = fullpath(path)
#     pattern_regex = regex
#     files = fls(path, pattern_regex) if recursion else ls(path, pattern_regex)
#     for file in files:
#         with open(file, 'r') as fileHandler: content = fileHandler.read()
#         if re.search(search,content):
#             content = re.sub(search, replace, content)
#             with open(file, 'w') as fileHandler: fileHandler.write(content)
#             print "Updated " + file

def _setlog(file="log.txt", mode='a', status=True, timestamp=True):
    """
    (file="log.txt", mode='a', status=True, timestamp=True)
    convenient shortcut: logon(), logoff()

    file: could be .csv to generate excel file, or be different to genereate multiple files
    mode: a=append; w=overwrite
    status: True  logging on, Prints output to both terminal and a file (log.txt, default name) globally
            False logging off, Prints output only to terminal
    timestamp: True/False, insert timestamp at the beginning and end in the log file

    Note, use this function carefully, because it changes the sys.stdout globally.
    """
    import sys, datetime

    class Logger(object):
        def __init__(self, file):
            self.file = file
            sys.stdout = sys.__stdout__
            print("log on with " + fullpath(self.file))
            self.terminal = sys.stdout
            self.log = open(file, mode)
            if timestamp:
                self.log.write("++++++++++log on at " + datetime.datetime.now().strftime("%Y-%m-%d_%H-%M-%S") + "\n")
                self.log.flush()

        def flush(self):
            self.terminal.flush()
            try:
                self.log.flush()
            except:  # in case self.log closed
                pass
            
        def write(self, message):
            self.terminal.write(message)
            self.log.write(message)
            self.log.flush()

        def off(self):
            if timestamp:
                self.log.write("++++++++++log off at " + datetime.datetime.now().strftime("%Y-%m-%d_%H-%M-%S") + "\n")
            self.log.write('\n')
            self.log.flush()
            self.log.close()
            sys.stdout = sys.__stdout__
            print("log off with " + fullpath(self.file))

    if status:
        # restore first if it has been changed
        try:
            sys.stdout.off()
        except AttributeError:
            pass

        sys.stdout = Logger(file)
    else:
        try:
            sys.stdout.off()
        except AttributeError:
            pass

def logon(file="log.txt", mode='a', status=True, timestamp=True):
    """
    only works for stdout, not stderr
    
    (file="log.txt", mode='a', status=True, timestamp=True)
    wrapper of log()
    logon(file="log.txt", mode='a', status=True, timestamp=True)

    Usage: logon(), logoff()

    file: could be .csv to generate excel file, or be different to genereate multiple files
    mode: a=append; w=overwrite
    status: True  logging on, Prints output to both terminal and a file (log.txt, default name) globally
            False logging off, Prints output only to terminal
    timestamp: True/False, insert timestamp at the beginning and end in the log file

    Note, use this function carefully, because it changes the sys.stdout globally.
    """
    _setlog(file=file,mode=mode,status=status,timestamp=timestamp)

def logoff():
    """
    wrapper of log()
    logon(), logoff()
    turn off the file redirection, does not accept any parameter, auto recognize log file
    """
    _setlog(status=False)

# +++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++
# print directory tree structure starts
# modified from http://code.activestate.com/recipes/577091/
# +++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++
def tree(path='./', sum=True, save=None, sort=True, case=True):
    """
    tree([path[, sum=True]) # Prints a directory tree structure. 
    sum=True (default) prints only folders, i.e., print less to show the big structure
    sum=False prints files plus folders
    save=None or a file path, eg, 'tree.log'
    sort=True, sort listed folders, call sorted()
    case: if True, get ['Ant', 'Bat', 'Cat', 'Goat', 'Lion', 'ant', 'bat', 'cat']
          if false, get ['ant', 'Ant', 'bat', 'Bat', 'cat', 'Cat', 'Goat', 'Lion']
    """
    import sys, os
    global PRINT_FILES; PRINT_FILES = not sum
    path = os.path.abspath(os.path.expanduser(path))

    def walk(root, dirs, files, prefix=''):
        if PRINT_FILES and files:
            file_prefix = prefix + ('|' if dirs else ' ') + '   '
            for name in files:
                print((file_prefix + name))
            print(file_prefix)
        dir_prefix, walk_prefix = prefix + '+---', prefix + '|   '
        for pos, neg, name in enumerate2(dirs):
            if neg == -1:
                dir_prefix, walk_prefix = prefix + '\\---', prefix + '    '
            path = os.path.join(root, name)
            try:
                dirs, files = listdir(path)[:2]
                print((dir_prefix + name + '\t(' + str(len(files)) +  ' files)'))
            except:
                print((dir_prefix + name))
            else:
                walk(path, dirs, files, walk_prefix)

    def listdir(path):
        dirs, files, links = [], [], []
        names = os.listdir(path)
        if sort: 
            if case:
                names = sorted(names)
            else:
                # https://stackoverflow.com/questions/13954841/sort-list-of-strings-ignoring-upper-lower-case
                names = sorted(names, key=lambda v: v.upper())
        for name in names:
            path_name = os.path.join(path, name)
            if os.path.isdir(path_name):
                dirs.append(name)
            elif os.path.isfile(path_name):
                files.append(name)
            elif os.path.islink(path_name):
                links.append(name)
        return dirs, files, links

    def enumerate2(sequence):
        length = len(sequence)
        for count, value in enumerate(sequence):
            yield count, count - length, value

    if save is not None: logon(file=save, mode='w', status=True, timestamp=False)
    dirs, files = listdir(path)[:2]
    print(path)
    walk(path, dirs, files)
    if not dirs:
        print('No subfolders exist')
    if save is not None: logoff()
# +++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++
# print directory tree structure ends
# +++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++

def regexp(string, pattern, method='find'):
    """regexp(string, pattern, method='find'), returns starting/ending indices in two lists
    [starts, ends] = regexp('abcfffgdefff','f{3}') -->[[3,9],[5,11]]  0 based
    if not find, returns [[],[]]

    method='split'
    ('Words, words, words.','\W+') --> ['Words', 'words', 'words', '']
    if not split, returns original string

    method='match'
    text = "He was carefully disguised but captured quickly by police."
    (text, '\w+ly') --> ['carefully', 'quickly']
    if not match, returns []

    seems like always a good idea to use r'' with regular expression, eg,
    newValue=ez.regexprep(oldValue,r'^(\d{1,2})(,|.)',r'\1|',1)

    python regular expression notes:
    1) special trick:
    matching string not starting with my:     ^(?!my)\w+$

    group back reference \ 1 \ 2. (Perl is $1 $2)   
    >>> re.sub(r'(\w+) (\w+)',r'\2 \1','Joe Bob')
        'Bob Joe'

    either or 
    (img|hdr)
    note: [ab] either character a or b, only work for single character

    2) general usage:
    . one character
    ^ $  start end
    [^ ] exclude
    [] include
    * >= 0
    + >= 1
    ? 0 or 1
    {n} exact n
    {m,n} m to n
    {n,} n or more
    \d  \D = [^0-9]
    \w [a-zA-Z0-9_]  \W = [^a-zA-Z0-9_]
    \s white space, tab etc
    letters only [a-zA-Z]
    \ t \ r \ n
    """
    # re.match() checks for a match only at the beginning of the string,
    # while re.search() checks for a match anywhere in the string
    import re
    if method == 'find':
        theMatch = re.finditer(pattern, string)
        starts=[]; ends=[]
        for ind, matchObj in enumerate(theMatch):
             starts.append(matchObj.start())
             ends.append(matchObj.end()-1)
        return [starts, ends]
    elif method == 'split':
        return re.split(pattern, string)
    elif method == 'match':
        return re.findall(pattern, string)

def regexpi(string, pattern, method='find'):
    """ignore cases
    regexp(string, pattern, method='find'), returns starting/ending indices in two lists
    [starts, ends] = regexp('abcfffgdefff','f{3}') -->[[3,9],[5,11]]  0 based
    if not find, returns [[],[]]

    method='split'
    ('Words, words, words.','\W+') --> ['Words', 'words', 'words', '']
    if not split, returns original string

    method='match'
    text = "He was carefully disguised but captured quickly by police."
    (text, '\w+ly') --> ['carefully', 'quickly']
    if not match, returns []

    seems like always a good idea to use r'' with regular expression, eg,
    newValue=ez.regexprep(oldValue,r'^(\d{1,2})(,|.)',r'\1|',1)

    python regular expression notes:
    1) special trick:
    matching string not starting with my:     ^(?!my)\w+$

    group back reference \1 \2. (Perl is $1 $2)   
    >>> re.sub(r'(\w+) (\w+)',r'\2 \1','Joe Bob')
        'Bob Joe'

    either or 
    (img|hdr)
    note: [ab] either character a or b, only work for single character

    2) general usage:
    . one character
    ^ $  start end
    [^ ] exclude
    [] include
    * >= 0
    + >= 1
    ? 0 or 1
    {n} exact n
    {m,n} m to n
    {n,} n or more
    \d  \D = [^0-9]
    \w [a-zA-Z0-9_]  \W = [^a-zA-Z0-9_]
    \s white space, tab etc
    letters only [a-zA-Z]
    \ t \ r \ n    
    """
    # re.match() checks for a match only at the beginning of the string,
    # while re.search() checks for a match anywhere in the string
    import re
    if method == 'find':
        theMatch = re.finditer(pattern, string, flags=re.IGNORECASE)
        starts=[]; ends=[]
        for ind, matchObj in enumerate(theMatch):
             starts.append(matchObj.start())
             ends.append(matchObj.end()-1)
        return [starts, ends]
    elif method == 'split':
        return re.split(pattern, string, flags=re.IGNORECASE)
    elif method == 'match':
        return re.findall(pattern, string, flags=re.IGNORECASE)

def regexprep(string, pattern, replace, count=0):
    """(string, pattern, replace, count=0)
    ('Baked Beans And Spam', '\sAnd\s', ' & ') --> 'Baked Beans & Spam'

    count: maximum number of pattern occurrences to be replaced; count must be a non-negative integer.
    If omitted or zero, all occurrences will be replaced

    seems like always a good idea to use r'' with regular expression, eg,
    newValue=ez.regexprep(oldValue,r'^(\d{1,2})(,|.)',r'\1|',1)

    python regular expression notes:
    1) special trick:
    matching string not starting with my:     ^(?!my)\w+$

    group back reference \1 \2. (Perl is $1 $2)   
    >>> re.sub(r'(\w+) (\w+)',r'\2 \1','Joe Bob')
        'Bob Joe'

    either or 
    (img|hdr)
    note: [ab] either character a or b, only work for single character

    2) general usage:
    . one character
    ^ $  start end
    [^ ] exclude
    [] include
    * >= 0
    + >= 1
    ? 0 or 1
    {n} exact n
    {m,n} m to n
    {n,} n or more
    \d  \D = [^0-9]
    \w [a-zA-Z0-9_]  \W = [^a-zA-Z0-9_]
    \s white space, tab etc
    letters only [a-zA-Z]
    \ t \ r \ n    
    """
    return re.sub(pattern, replace, string, count=count)

def regexprepi(string, pattern, replace, count=0):
    """case-sensitive
    (string, pattern, replace, count=0)
    ('Baked Beans And Spam', '\sAND\s', ' & ') --> 'Baked Beans & Spam'

    count: maximum number of pattern occurrences to be replaced; count must be a non-negative integer.
    If omitted or zero, all occurrences will be replaced

    seems like always a good idea to use r'' with regular expression, eg,
    newValue=ez.regexprep(oldValue,r'^(\d{1,2})(,|.)',r'\1|',1)

    python regular expression notes:
    1) special trick: 
    matching string not starting with my:     ^(?!my)\w+$

    group back reference \1 \2. (Perl is $1 $2)   
    >>> re.sub(r'(\w+) (\w+)',r'\2 \1','Joe Bob')
        'Bob Joe'

    either or 
    (img|hdr)
    note: [ab] either character a or b, only work for single character

    2) general usage:
    . one character
    ^ $  start end
    [^ ] exclude
    [] include
    * >= 0
    + >= 1
    ? 0 or 1
    {n} exact n
    {m,n} m to n
    {n,} n or more
    \d  \D = [^0-9]
    \w [a-zA-Z0-9_]  \W = [^a-zA-Z0-9_]
    \s white space, tab etc
    letters only [a-zA-Z]
    \ t \ r \ n    
    """
    return re.sub(pattern, replace, string, count=count, flags=re.IGNORECASE)

def sprintf(formatString, *args, **kwargs):
    """
    note: if specify skipdollar=1, $ (but not others) syntax will be entirely skipped, useful for R codes (df$col), or certain bash codes
    see usage 3&4

    (formatString, *args, **kwargs)
    data dictionary is passed as *args, right now **kwargs only for skipdollar

    # -------------------------------------------------------------------------
    Examples:
    language = 'Python'; number = 2
    theDict = {"language": "Matlab", "number": 1}
    
    
    
    # usage 1: you have to pass unnamed/positional args manually
    # evaluated first before other methods in the following
    # call formatString % args
    s = sprintf('%02d\n is bigger than\n %02d',4,3)
    s = sprintf('%02d\n is bigger than\n %02d',[4,3])
    s = sprintf('%02d\n is bigger than\n %02d',(4,3))
    s = sprintf('%s has %03d quote types', language, number)
    s = sprintf('%s', language)
    


    # usage 2: if '%(' found in the string, call formatString % args
    # the following methods will be skipped
    s = sprintf('%(language)s has %(number)03d quote types.', {"language": "Python", "number": 2})
    s = sprintf('%(language)s has %(number)03d quote types.', theDict)
    s = sprintf('%(language)s has %(number)03d quote types.', locals()) # locals() returns a dictionary
    s = sprintf('%(language)s has %(number)03d quote types.') # auto get dictionary from locals()
    # 
    # this usage is very useful for complex string mixing different styles and/or env/bash variables
    s = sprintf('family="springer"; echo $family $HOME %(language)s')
    


    # usage 3: if ${var} or $var exist, replace them with regex both to {var}
    # then send to formatString.format() for evaluation -- see usage 4 for details
    # skipdollar will skip usage 3&4 entirely, sprintf('$packt {family}',skipdollar=1)
    # 
    # note: the bash $ is not natively supported by python
    # here I hack it to support the bash style, so that cmd can be directly used by both python and bash
    # but won't support $number:03d (ie, skipped) for consistency with bash
    # $0 ${0} $() ${0} {} {0,1..3} {1..$(2)} $number:03d ${number:03d} {number:03d} etc
    # other than [a-zA-Z_]+\w*? will all be skipped
    s = sprintf('$language has $number quote types.', theDict)
    s = sprintf('$language has $number quote types.')
    s = sprintf('${language}_ has $number quote types.')
    s = sprintf('$language')
    s = sprintf('$PATH')    # <--existing env variables (eg, $PATH) will not be replaced but kept as is (for later bash)
    s = sprintf('${PATH}')  # <--existing env variables (eg, ${PATH}) will not be replaced but kept as is (for later bash)
                            # <--env vars checked via os.environ (which returns a dictionary)
    #
    # better do not mix different styles, ie, %s $var {var} when formating  <--except ${var}
    s = sprintf('${language} has {number:03d} quote types.')  # -> 'Python has {number:03d} quote types.'
    s = sprintf('${language} has {number} quote types.')      # -> 'Python has 2 quote types.'
    s = sprintf('${language} = {language}')                   # OK
    


    # usage 4: usage 3 converges on this method, calls formatString.format()
    # skipdollar will skip usage 3&4 entirely, sprintf('$packt {family}',skipdollar=1)
    s = sprintf('{language} has {number} quote types.', theDict) # <--auto unpack
    s = sprintf('{language} has {number} quote types.')          # <--auto search
    


    # usage 5: long string
    longString = '''
    Hello, %s
    
    Best,
    '''
    s = sprintf(longString, language)
    # -------------------------------------------------------------------------
    """  
    # args is a tuple even when only one arg is passed in
    import re, inspect
    
    if len(args) > 1:
        return formatString % args
    elif len(args) == 1:
        if type(args[0]) in [list, tuple]:
            return formatString % tuple(*args)
        elif type(args[0]) in [dict]:
            # syntax %()
            if re.search('%\(', formatString):
                return formatString % args[0]
            else:
                # kwargs is empty {} if not specified
                try:
                    if kwargs['skipdollar']!=1: kwargs.pop('skipdollar')
                except:
                    pass
                if 'skipdollar' not in list(kwargs.keys()):
                    # \w is [a-zA-Z0-9_], but I do not want pure number
                    # so [a-zA-Z_]+\w*? for valid variable naming
                    # ${number}, ${language}_ to {number}, {language}_
                    # but skip ${number:03d}
                    # but not replace ${PATH}
                    rs = re.findall('\$\{([a-zA-Z_]+\w*?)\}', formatString)
                    for r in rs:    
                        if r not in os.environ:
                            formatString = re.sub('\$\{('+r+')\}', r'{\1}', formatString)
                        else:
                            # trick the later .format() function ${PATH}   ->   |___|PATH|__|
                            formatString = re.sub('\$\{('+r+')\}', r'|___|\1|__|', formatString)
                    
                    # now $var -> {var}
                    rs = re.findall('\$([a-zA-Z_]+\w*?)', formatString)    
                    for r in rs:
                        if r not in os.environ:
                            formatString = re.sub('\$('+r+')', r'{\1}', formatString)
                    
                    # anything other than {[a-zA-Z_]+\w*?}--will be skipped
                    # do not know how to construct regex for invalid names, so...
                    allrs = re.findall('\{(.*?)\}', formatString)
                    validrs = re.findall('\{([a-zA-Z_]+\w*?)\}', formatString)
                    for r in allrs:
                        if r not in validrs:
                            # r may contain irregular char (eg, dot .) in this case, therefore re.sub may not work well
                            # formatString = re.sub('\{('+r+')\}', r'@__@\1@___@', formatString)
                            formatString = formatString.replace('{'+r+'}','@__@'+r+'@___@')
                    formatString = formatString.format(**args[0])
                    # replace back {invalid ...}
                    formatString = re.sub('@__@(.*?)@___@', r'{\1}', formatString)
                    # replace back env variable |___|PATH|__|  -->  ${PATH}
                    formatString = re.sub('\|___\|([a-zA-Z_]+\w*?)\|__\|', r'${\1}', formatString)

                    return formatString
                elif kwargs['skipdollar']==1:
                    return formatString
        else:
            # a single string or int
            return formatString % args         
    # no args passed in
    else:
        caller = inspect.currentframe().f_back
        return sprintf(formatString,caller.f_locals,**kwargs)
        
def iff(expression, result1, result2):
    """iff(expression, result1, result2)"""
    if expression:
        return result1
    else:
        return result2
ifelse=iff

def clear(module, recursive=False):
    """clear(module, recursive=False)
    remove a module from sys.modules so it cannot be searched.
    when recursive=True, remove the module and its submodules"""
    if recursive:
        for mod in list(sys.modules.keys()):
            if mod.startswith(module):
                del(sys.modules[mod])
    else:
        if module in sys.modules:
            del(sys.modules[module])

try:
    import pytz
    import tzlocal
except:
    pass

import os, sys, platform, string, random, re
import datetime

from random import seed as Randomize
from random import seed as randomize
from random import shuffle as RandomizeArray
from random import shuffle as randomizearray
from random import randint as Random
from random import choice as RandomChoice
from random import choice as randomchoice
def Permute(iterable=[]):
  """Permute(iterable=[])
  Returns permutations in a list
  e.g., ([1,2,3]) --> [(1, 2, 3), (1, 3, 2), (2, 1, 3), (2, 3, 1), (3, 1, 2), (3, 2, 1)]
  """
  from itertools import permutations
  return list(permutations(iterable))
permute = Permute

from time import sleep
wait = sleep

def pause():
    """
    print("Are you going to do sth?")
    if ez.pause():
        do sth
    """
    try:
        input("Press <Enter> to continue..., or Ctrl+c to exit.")
        return True
    except KeyboardInterrupt:
        # raise Exception('User cancelled.')
        return False

def round2(x,y=None):
    """
    always round 0.5 up (Python 2 rounding behavior)
    Python 3's way (called "round half to even" or "banker's rounding") is considered the standard rounding method these days
    https://stackoverflow.com/questions/10825926/python-3-x-rounding-behavior
    """
    return round(x+1e-15,y)

def num(s,force=True):
    """num(s)
    num(3),num(3.7)-->3
    num('3')-->3, num('3.7')-->3.7
    num('3,700')-->ValueError
    num('3a'),num('a3'),-->ValueError
    num('3e4') --> 30000.0
    num(' '),num('') -->ValueError
    if force, not raise error, but returns pd.NA
    """
    try:
        return int(s)
    except ValueError:
        try:
            return float(s)
        except ValueError:
            if force: 
                import pandas as pd
                return pd.NA
            raise ValueError('argument is not a string of number')

def isempty(s):
    """isempty(s), None, numpy.nan return True"""
    try:
        return len(s) == 0
    except TypeError:
        # None, numpy.nan
        return True

def unique(seq):
    """
    unique(seq), union(seq1,seq2), intersect(seq1,seq2), setdiff(seq1,seq2) in original order
    seq could be a list
    note: setdiff(seq1,seq2) may not be equal to setdiff(seq2,seq1)
            >>> unique('abracadaba')
            ['a', 'b', 'r', 'c', 'd']
            >>> unique('simsalabim')
            ['s', 'i', 'm', 'a', 'l', 'b']
            >>>
            >>> setdiff('abracadaba','simsalabim')
            ['r', 'c', 'd']
            >>> setdiff('simsalabim','abracadaba')
            ['s', 'i', 'm', 'l']
    duplicate(seq) # returns a list of duplicated elements in original order
    # e.g.
    # a = [1,5,2,3,2,1,5,6,5,5,5]
    # duplicate(a) # yields [2, 1, 5]
    """
    # from .orderedset import OrderedSet
    return list(OrderedSet(seq))

def union(seq1,seq2):
    """
    unique(seq), union(seq1,seq2), intersect(seq1,seq2), setdiff(seq1,seq2) in original order
    seq could be a list
    note: setdiff(seq1,seq2) may not be equal to setdiff(seq2,seq1)
            >>> unique('abracadaba')
            ['a', 'b', 'r', 'c', 'd']
            >>> unique('simsalabim')
            ['s', 'i', 'm', 'a', 'l', 'b']
            >>>
            >>> setdiff('abracadaba','simsalabim')
            ['r', 'c', 'd']
            >>> setdiff('simsalabim','abracadaba')
            ['s', 'i', 'm', 'l']
    duplicate(seq) # returns a list of duplicated elements in original order
    # e.g.
    # a = [1,5,2,3,2,1,5,6,5,5,5]
    # duplicate(a) # yields [2, 1, 5]
    """
    # from .orderedset import OrderedSet
    return list(OrderedSet(seq1) | OrderedSet(seq2))

def intersect(seq1,seq2):
    """
    unique(seq), union(seq1,seq2), intersect(seq1,seq2), setdiff(seq1,seq2) in original order
    seq could be a list
    note: setdiff(seq1,seq2) may not be equal to setdiff(seq2,seq1)
            >>> unique('abracadaba')
            ['a', 'b', 'r', 'c', 'd']
            >>> unique('simsalabim')
            ['s', 'i', 'm', 'a', 'l', 'b']
            >>>
            >>> setdiff('abracadaba','simsalabim')
            ['r', 'c', 'd']
            >>> setdiff('simsalabim','abracadaba')
            ['s', 'i', 'm', 'l']
    duplicate(seq) # returns a list of duplicated elements in original order
    # e.g.
    # a = [1,5,2,3,2,1,5,6,5,5,5]
    # duplicate(a) # yields [2, 1, 5]
    """
    # from .orderedset import OrderedSet
    return list(OrderedSet(seq1) & OrderedSet(seq2))

def setdiff(seq1,seq2):
    """
    unique(seq), union(seq1,seq2), intersect(seq1,seq2), setdiff(seq1,seq2) in original order
    seq could be a list
    note: setdiff(seq1,seq2) may not be equal to setdiff(seq2,seq1)
            >>> unique('abracadaba')
            ['a', 'b', 'r', 'c', 'd']
            >>> unique('simsalabim')
            ['s', 'i', 'm', 'a', 'l', 'b']
            >>>
            >>> setdiff('abracadaba','simsalabim')
            ['r', 'c', 'd']
            >>> setdiff('simsalabim','abracadaba')
            ['s', 'i', 'm', 'l']
    duplicate(seq) # returns a list of duplicated elements in original order
    # e.g.
    # a = [1,5,2,3,2,1,5,6,5,5,5]
    # duplicate(a) # yields [2, 1, 5]
    """
    # from .orderedset import OrderedSet
    return list(OrderedSet(seq1) - OrderedSet(seq2))

def setdiff2(seq2,seq1):
    return list(OrderedSet(seq1) - OrderedSet(seq2))

def compare(lh,rh,value=False):

    print('\t\t\t\tUnion: {:4.0f}\n'.format(len(union(lh,rh))))
    print('\t\tLH: {:4.0f} /{:4.0f}\t\t\t\tRH: {:4.0f} /{:4.0f}\n'.format(len(lh),len(unique(lh)),len(rh),len(unique(rh))))
    print('\t\t\t\tInter: {:4.0f}\n'.format(len(intersect(lh,rh))))
    print('\n')
    print('\t\tLH>: {:4.0f}\t\t\t\t<RH: {:4.0f}\n'.format(len(setdiff(lh,rh)),len(setdiff(rh,lh))))

    if value:
        print('LH>:')
        ppprint(','.join(setdiff(lh,rh)))
        print('\n')
        print('RH>:')
        ppprint(','.join(setdiff(rh,lh)))

def duplicate(seq):
    """
    unique(seq), union(seq1,seq2), intersect(seq1,seq2), setdiff(seq1,seq2) in original order
    seq could be a list
    note: setdiff(seq1,seq2) may not be equal to setdiff(seq2,seq1)
            >>> unique('abracadaba')
            ['a', 'b', 'r', 'c', 'd']
            >>> unique('simsalabim')
            ['s', 'i', 'm', 'a', 'l', 'b']
            >>>
            >>> setdiff('abracadaba','simsalabim')
            ['r', 'c', 'd']
            >>> setdiff('simsalabim','abracadaba')
            ['s', 'i', 'm', 'l']
    duplicate(seq) # returns a list of duplicated elements in original order
    # e.g.
    # a = [1,5,2,3,2,1,5,6,5,5,5]
    # duplicate(a) # yields [2, 1, 5]
    """
    # not working...
    # # from .orderedset import OrderedSet
    # seen = OrderedSet()
    # seen_add = seen.add
    # # adds all elements it doesn't know yet to seen and all other to seen_twice
    # seen_twice = OrderedSet( x for x in seq if x in seen or seen_add(x) )
    # # turn the set into a list (as requested)
    # return list( seen_twice )
    import pandas as pd
    df=pd.DataFrame({'col':seq})
    df=(df.loc[df.duplicated(subset=['col'],keep='first').tolist()]).drop_duplicates()
    return df.col.tolist()

# As of Python version 3.7, dictionaries are ordered. In Python 3.6 and earlier, dictionaries are unordered.
from collections import OrderedDict
class JDict(OrderedDict):
    """Jerry's dictionary
    customized ordered dictionary class with convient attributes and methods

    Methods:
    sort([reverse=False])
        sort the dictionary by key in (reverse) order
        returns a new sorted JDict, does not modify the original/itself
    update({k:v})
        adds new keys/vals or updates existing keys/vals, cannot delete
        v could be interger, string, list, dictionary
            if interger, string, assign the same key, updates val
            if list, assign the same key, append val to the existing list
            if dictionary, assign the same key, recursive call
            once, a type of v (e.g. list) determined at the first time, cannot be changed
            v could be a dictionary of list, dictionary or whatever

        e.g.
        a = JDict(); b = JDict(); c = JDict(); d = JDict()
        a.update({'k1':'v1'}); print(a)
        a.update({'k1':'v'}); print(a)
        a.update({'k2':'v'}); print(a)
        a.update({'k1':'v1', 'k2':'v2'}); print(a)

        b.update({'k1': ['v1']}); print(b   # <-- seems no necessary to use ['v1',], side note: for tuple (2,))
        b.update({'k1': ['v', 'v']}); print(b)
        b.update({'k2': ['v']}); print(b)
        b.update({'k1': ['v1'], 'k2': ['v2']}); print(b)

        c.update({'k1': {'k11': 'v11'}}); print(c)
        c.update({'k1': {'k11': 'v'}}); print(c)
        c.update({'k1': {'k11': 'v11', 'k12': 'v12'}}); print(c)
        c.update({'k2': {'k21': ['v']}}); print(c)
        c.update({'k1': {'k13':{'k131':'v131'}}, 'k2': {'k21': ['v']}}); print(c)
        c.update({'k1': {'k13':{'k131':'v'}}, 'k2': {'k21': ['v']}}); print(c)

        d.update({'k1': JDict({'k11': 'v'})}); print(d)
        d.update({'k1': JDict({'k11': 'v11', 'k12': 'v12'})}); print(d)
    """
    def __init__(self, theDict={}):
        OrderedDict.__init__(self, theDict)
        # self = OrderedDict(theDict)  # -->not necessary or wrong??

    def update(self, theDict, prev=None):
        # prev for recursive call
        if prev == None: prev = self

        for (newKey,newVal) in list(theDict.items()):
            # if no newKey, this is easy, simply initialize
            if newKey not in prev:
                prev[newKey] = newVal
            else:
            # key exisiting
                key = newKey
                if type(newVal) in [list]:
                    # assuming the existing val is also a list
                    prev[key].extend(newVal)
                elif type(newVal) in [OrderedDict, dict, JDict]:
                    # assuming the existing val is a dictionary
                    # recursive call, starting from the key
                    self.update(newVal,prev[key])
                else:
                    # simply update key:newVal
                    prev[key] = newVal

    def sort(self, reverse=False):
        """
        sort by keys, default ascending order
        returns a new JDict object
        """
        # self = JDict(sorted(self.items(),reverse=reverse)) will not work, see
        # http://stackoverflow.com/questions/1216356/
        return JDict(sorted(list(self.items()),reverse=reverse))

class TimeStamp(object):
    """
    Essentially pd timestamp with default timezone that has a "candy coat"
    ts=TimeStamp(pd_timestamp_obj) # default, current time 
    ts.datestr
    ts.datets           # pd timestamp object that refers to the very beginning of the date (00:00:00)
    ts.datetimestamps   # in seconds (e.g., for Zacks event date). Refer to very beginning of the date (00:00:00)
    ts.timestampms      # in milli seconds (e.g., for TD stock price). A date without time will be calculationed using 00:00:00
    ts.ts               # access the pd timestamp object
    ts.fromtimestamp(1610079792506).datestr  # returns TimeStamp obj then .datestr
    ts.todatetime()     # parse to a datetime/timestamp as US/Eastern, returns TimeStamp obj

    Still sometimes, naive datetime is preferred, eg, pd.to_datetime(ts.datestr) - pd.offsets.Day(120); 
    pd.to_datetime(ts.datestr) does not accept a tz parameter
    But note that, pd.date_range('2020-09-16','2020-11-16',freq='B',tz='US/Central') accepts a tz parameter
    """
    def __init__(self, timestamp=None, tz="US/Eastern"):
        import pandas as pd
        if timestamp is None: 
            self.ts = pd.Timestamp.now(tz=tz)
        else:
            self.ts = timestamp

    @property
    def datestr(self):
        return self.ts.strftime("%Y-%m-%d")

    @property
    def datets(self,tz='US/Eastern'):
        import pandas as pd
        date = self.ts.strftime("%Y-%m-%d")
        ts = pd.to_datetime(date).tz_localize(tz=tz)
        return ts

    @property
    def datetimestamps(self,tz='US/Eastern'):
        import pandas as pd
        date = self.ts.strftime("%Y-%m-%d")
        ts = pd.to_datetime(date).tz_localize(tz=tz)
        return int(ts.timestamp())

    @property
    def timestampms(self):
        # int() would introduce very tiny error when convert back
        return int(self.ts.timestamp()*1000)

    @classmethod
    def fromtimestamp(cls,timestamp,unit='ms',local='UTC',convertto='US/Eastern',*args,**kwargs):
        import pandas as pd
        ts = pd.to_datetime(timestamp,unit=unit,*args,**kwargs).tz_localize(local).tz_convert(convertto)
        return TimeStamp(ts)
    
    @classmethod
    def todatetime(cls,arg,tz='US/Eastern',*args,**kwargs):
        import pandas as pd
        ts = pd.to_datetime(arg,*args,**kwargs).tz_localize(tz)
        return TimeStamp(ts)

class Moment(object):
    """A datetime like class, but with convenient attributes and methods
    Moment(moment=datetime.datetime(2020,5,30,15,30,00,100)).moment                  -> naive 15:30 datetime
    Moment('US/Mountain',moment=datetime.datetime(2020,5,30,15,30,00,100)).moment    -> 15:30 in mountain time
    Moment('US/Mountain').moment                                                     -> current time converted to mountain
    Moment().Convert('US/Mountain').moment                                           -> current time converted to mountain
    Moment().moment                                                                  -> current naive time

    Has common datetime attributes as strings.
    Wraps timedelta for easy calling.
    Formats the datetime to a string.
    Transforms a datetime-like string to a Moment object.
    Converts to specified timezone

    Attributes:
        moment, returns the datetime instance

        (the following all return string, while ignoring tzinfo)
        datetime, returns the string of date and time
        date, returns the string of the date
        time
        timezone
        year
        month
        day
        weekday
        hour
        ampm
        minute
        second
        microsecond
    """
    def __init__(self, timezone=None, moment=None):
        """Generates the current datetime in specified timezone, or local naive datetime if omitted, (when moment omitted).
        
        Args:
            timezone is the specified timezone string
                UTC: Coordinated Universal Time (not a time zone, but a time standard)
                GMT: Greenwich Mean Time (a time zone)
                but Moment('UTC').datetime = Moment('GMT').datetime

                US timezone:
                US/Alaska, US/Aleutian, US/Arizona, US/Central, US/East-Indiana, US/Eastern, 
                US/Hawaii, US/Indiana-Starke, US/Michigan, US/Mountain, US/Pacific, US/Pacific-New, US/Samoa

                Use the codes to look up:
                from pytz import all_timezones
                print len(all_timezones)
                for zone in all_timezones:
                    if 'US' in zone:
                        print zone

            moment could be naive (without tzinfo) or aware (with tzinfo) datetime

        """
        if not timezone:
            if not moment:
                moment = datetime.datetime.now()
            else:
                moment = moment
        else:
            timezone = pytz.timezone(timezone)
            if not moment:
                # if moment is not provided, use now
                moment = datetime.datetime.now()
                # first: assign local timezone
                tz = tzlocal.get_localzone()
                # moment = tz.localize(moment)
                import zoneinfo
                moment = moment.replace(tzinfo=zoneinfo.ZoneInfo(str(tz)))
                # second: convert to desired timezone
                moment = timezone.normalize(moment.astimezone(timezone))
            else:
                if not moment.tzinfo:
                    # if provided moment is naive
                    # simply attach timezone to it
                    moment = timezone.localize(moment)
                else:
                    # if provided moment is aware
                    # convert to desired timezone
                    # add normalize to correct potential DST after Moment().Shift() Methods
                    moment = timezone.normalize(moment.astimezone(timezone))

        self.moment = moment
        self.timezone = str(moment.tzinfo) if moment.tzinfo else None

        self.datetime = moment.strftime("%Y-%m-%d_%H-%M-%S")
        # in the format of 1/2/2014 for single digits, and 10/14/2014 for two digits
        self.date = str(moment.month)+"/"+str(moment.day)+"/"+str(moment.year)
        self.time = moment.strftime("%H:%M:%S")
        self.year = moment.strftime("%Y")
        self.month = moment.strftime("%B")
        self.day = moment.strftime("%d")
        self.weekday = moment.strftime("%A")
        self.hour = moment.strftime("%I")
        self.ampm = moment.strftime("%p")
        self.minute = moment.strftime("%M")
        self.second = moment.strftime("%S")
        self.microsecond = moment.strftime("%f")

    def Shift(self, *args, **kwargs):
        """Shifts time to the past or future.

        Args:
            the same as timedelta: [days[, seconds[, microseconds[, milliseconds[, minutes[, hours[, weeks]]]]]]]
            Use negative values to shift to the past

        Returns:
            an instance of Moment

        Raises:
            None
        """
        future = self.moment + datetime.timedelta(*args, **kwargs)
        return Moment(timezone=self.timezone, moment=future)

    def Format(self, datetimeFormat):
        """Formats the datetime to a string.

        Args:
            the format of the desired string.

        Returns:
            a string

        Raises:
            None
        """
        datetimeString = self.moment.strftime(datetimeFormat)
        return datetimeString

    def Convert(self, timezone=None):
        """Convert to a specified timezone.
        timezone is a specified timezone string, see the class help
        returns an instance of Moment
        """
        # if naive
        if self.timezone is None:
            tz = tzlocal.get_localzone()
            # attach local tz
            # dt = tz.localize(self.moment)
            import zoneinfo
            dt = self.moment.replace(tzinfo=zoneinfo.ZoneInfo(str(tz)))
        else:
            dt = self.moment
        timezone = pytz.timezone(timezone)
        newdt = timezone.normalize(dt.astimezone(timezone))
        return Moment(timezone=None, moment=newdt)

    @classmethod
    def Transform(cls, datetimeString, datetimeFormat, timezone=None):
        """Transforms a datetime-like string to a Moment object.

        Args:
             a datetime-like string, e.g. "21/11/06 16:30"
             the format information of the string, e.g. "%d/%m/%y %H:%M"

                    %a      Sun, Mon, ..., Sat (en_US);
                    %A      Sunday, Monday, ..., Saturday (en_US);
                    %w      0 is Sunday and 6 is Saturday.   0, 1, ..., 6     
                    %d      01, 02, ..., 31      
                    %b      Jan, Feb, ..., Dec (en_US);
                    %B      January, February, ..., December (en_US);
                    %m      01, 02, ..., 12      
                    %y      Year    00, 01, ..., 99      
                    %Y      Year    1970, 1988, 2001, 2013   
                    %H      Hour (24-hour clock)    00, 01, ..., 23      
                    %I      Hour (12-hour clock)    01, 02, ..., 12      
                    %p      AM, PM (en_US);
                    %M      Minute 00, 01, ..., 59      
                    %S      Second 00, 01, ..., 59
                    %f      Microsecond     000000, 000001, ..., 999999
                    %z      UTC offset in the form +HHMM or -HHMM   (empty), +0000, -0400, +1030
                    %Z      Time zone name  (empty), UTC, EST, CST
                    %j      Day of the year     001, 002, ..., 366   
                    %U      Week number of the year (Sunday as the first day of the week)
                    %W      Week number of the year (Monday as the first day of the week)
                    %c      Locale's appropriate date and time representation.  Tue Aug 16 21:30:00 1988 (en_US);
                    %X      Locale's appropriate date representation.   08/16/88 (None); 08/16/1988 (en_US);
                    %X      Locale's appropriate time representation.   21:30:00 (en_US);
                    %%      A literal '%' character.    %
            
            timezone is the specified timezone string, see the class help

        Returns:
            a Moment object

        Raises:
           None
        """
        # strptime returns naive datetime
        # hack: python 2.7 does not recognize %z
        # input: 
        #       datetimeString = 'Thu, 18 Jun 2015 16:52:23 -0400'
        #       datetimeFormat = '%a, %d %b %Y %H:%M:%S %z'
        # output: 
        #       moment (datetime with tzinfo) which is like:
        #       datetime.datetime(2015, 6, 18, 16, 52, 23, tzinfo=-0400)
        if '%z' in datetimeFormat:
            from datetime import timedelta, tzinfo
            class FixedOffset(tzinfo):
                """Fixed offset in minutes: `time = utc_time + utc_offset`."""
                def __init__(self, offset):
                    self.__offset = timedelta(minutes=offset)
                    hours, minutes = divmod(offset, 60)
                    #NOTE: the last part is to remind about deprecated POSIX GMT+h timezones
                    #  that have the opposite sign in the name;
                    #  the corresponding numeric value is not used e.g., no minutes
                    self.__name = '<%+03d%02d>%+d' % (hours, minutes, -hours)
                def utcoffset(self, dt=None):
                    return self.__offset
                def tzname(self, dt=None):
                    return self.__name
                def dst(self, dt=None):
                    return timedelta(0)
                def __repr__(self):
                    # offset in minutes
                    offset = self.utcoffset().total_seconds()/60
                    sign = '-' if offset < 0 else '+'
                    # /60.  add .  to get decimal points
                    return "%s%02d%02d" % (sign, abs(offset) / 60., abs(offset) % 60)
            import re
            naive_datetimeString = re.sub('[-+]\d{4}','',datetimeString)
            naive_dt = datetime.datetime.strptime(naive_datetimeString, datetimeFormat.replace('%z',''))
            offset_str = re.findall('[-+]\d{4}',datetimeString)[0]
            offset = int(offset_str[-4:-2])*60 + int(offset_str[-2:])
            if offset_str[0] == "-":
                offset = -offset
            moment = naive_dt.replace(tzinfo=FixedOffset(offset))
            # print moment
            # datetime.datetime(2015, 6, 18, 16, 52, 23, tzinfo=-0400)
        # hack end
        else:
            moment = datetime.datetime.strptime(datetimeString, datetimeFormat)
        return Moment(timezone=timezone, moment=moment)

def daylight(dt=None, timezone="US/Eastern"):
    """
    dt: default current utc datetime
    https://stackoverflow.com/a/55735278/2292993
    daylight(datetime(2022, 3, 13, 1), timezone="US/Eastern") # False
    daylight(datetime(2022, 3, 13, 2), timezone="US/Eastern") # False
    daylight(datetime(2022, 3, 13, 3), timezone="US/Eastern") # True
    daylight(datetime(2022, 3, 14), timezone="US/Eastern")    # True
    """
    import pytz
    from datetime import datetime
    if dt is None:
        dt = datetime.utcnow()
    timezone = pytz.timezone(timezone)
    timezone_aware_date = timezone.localize(dt, is_dst=False)
    return timezone_aware_date.tzinfo._dst.seconds != 0

def lines(path='.', pattern='\.py$|.ini$|\.c$|\.h$|\.m$', recursive=True):
    """Counts lines of codes, counting empty lines as well.
    lines(path='.', pattern='\.py$|.ini$|\.c$|\.h$|\.m$', recursive=True)
    """

    # modified from https://liangsun.org/posts/python-code-for-counting-loclines-of-code/
    # Created by Liang Sun <i@liangsun.org> in 2012
    # This code is for Python 2.x
    COUNT_EMPTY_LINE = True

    def read_line_count(fname):
        count = 0
        for line in open(fname).readlines():
            if COUNT_EMPTY_LINE or len(line.strip()) > 0:
                count += 1
        return count

    line_count = 0
    file_count = 0
    files = fls(path,pattern) if recursive else ls(path,pattern)
    for file in files:
        try:
            file_count += 1
            c = read_line_count(file)
            print("%s : %d" % (os.path.basename(file), c))
            line_count += c
        except:
            pass

    print('-----------------------------')
    print('File counted: %d' % file_count)
    print('Line counted: %d' % line_count)
    print('Done!')

def keygen(length=8, complexity=3):
    """generate a random key
    keygen(length=8, complexity=3)

    length: how long is the generated key

    complexity:
    1 = digits only
    2 = digits + lowercase  i.e., 'abcdefghijklmnopqrstuvwxyz'
    3 = digits + uppercase (default)
    4 = digits + lowercase + uppercase
    5 = digits + lowercase + uppercase + punctuation
    6 = digits + lowercase + uppercase + punctuation + whitespace
    or alternatively, complexity could be a list of chars like, '0123456789abcdefABCDEF'
    in which case, the key will only be composed of chars from the list
    """
    import string
    import random
    if complexity == 1:
        chars = string.digits
    elif complexity == 2:
        chars = string.digits + string.ascii_lowercase
    elif complexity == 3:
        chars = string.digits + string.ascii_uppercase
    elif complexity == 4:
        chars = string.digits + string.ascii_lowercase + string.ascii_uppercase
    elif complexity == 5:
        chars = string.digits + string.ascii_lowercase + string.ascii_uppercase + string.punctuation
    elif complexity == 6:
        chars = string.digits + string.ascii_lowercase + string.ascii_uppercase + string.punctuation + string.whitespace
    elif complexity not in [1, 2, 3, 4, 5, 6]:
        chars = str(complexity)
    return ''.join(random.choice(chars) for x in range(length))

def hashes(str_or_filepath, reference=None):
    """
    str_or_filepath:
        if str, hashes a str with a random salt using pbkdf2_hmac
        if filepath, calculate/print a file's md5 32, sha1 32 (and optionally compare with a reference); can handle big files in a memory efficient way
    reference: 
        if str, when reference provided returns T/F; otherwise return hashed 
        if filepath, when reference provided, print matched/not matched; otherwise print hashed only
    """
    import hashlib
    if exists(str_or_filepath):
        filename = str_or_filepath
        md5 = hashlib.md5()
        sha1 = hashlib.sha1()
        with open(filename, 'rb') as f:
            for chunk in iter(lambda: f.read(128 * md5.block_size), b''):
                md5.update(chunk)
                sha1.update(chunk)
        if reference:
            if reference.lower() == md5.hexdigest().lower():
                print('md5 32: ' + md5.hexdigest() + ' (matched)!')
            else:
                print('md5 32: ' + md5.hexdigest() + ' (NOT MATCHED)!')
            if reference.lower() == sha1.hexdigest().lower():
                print('sha1 32: ' + sha1.hexdigest() + ' (matched)!')
            else:
                print('sha1 32: ' + sha1.hexdigest() + ' (NOT MATCHED)!')
        else:
            print('md5 32: ' + md5.hexdigest())
            print('sha1 32: ' + sha1.hexdigest())
    else:
        # https://nitratine.net/blog/post/how-to-hash-passwords-in-python/
        import binascii
        if reference:
            reference = reference.encode('utf-8')
            reference = binascii.a2b_hex(reference)
            salt = reference[:32] 
        else:
            salt = os.urandom(32)
        password = str_or_filepath
        hashed = hashlib.pbkdf2_hmac(
            'sha256', # The hash digest algorithm for HMAC
            password.encode('utf-8'), # Convert the password to bytes
            salt, # Provide the salt
            100000, # It is recommended to use at least 100,000 iterations of SHA-256 
            # dklen=128 # Get a 128 byte key 
        )
        if reference:
            return hashed==reference[32:]        
        else:
            storage = salt + hashed
            # hexadecimal representation of the binary data
            storage = binascii.b2a_hex(storage)
            # byte to string
            storage = storage.decode('utf-8')
            return storage

def readt(path, *args, **kwargs):
    """
    read text file
    returns a list of lines
    """
    with open(path, 'r') as f: 
        lines = f.readlines()
    return lines

def readx(path, sheet=0, r=[1,], c=None, *args, **kwargs):
    """
    (path, sheet=0, r=[1,], c=None, *args, **kwargs)
    Read xlsx, xls file into a list (see returns for details), using openpyxl, xlrd library
    Args:
        path, a xlsx, xls file
        sheet, either sheet number (the first is 0) or sheet name (e.g., Sheet1)
        r, None=start-end; 3=4th row (zero based); [3,4,5]=listed rows; [1] or [1,]=from 2nd row to end (ie, skipping first row, or header), [3] or [3,]=from 4th row to end
        c, the same format as r
    Returns:
        if multiple rows/columns, returns a list of list [[row1],[row2]]. To further slice returned result, result[r][c]
        if a single row/column, returns a list
        if a single cell, returns that cell as string, number depending on that cell type
        Note: the cells/elements in a list/column/row could have different data types, exactly the same as the actual data type in the excel file
              eg, a number is formated as string in excel file, then the returned type is string, u'2017'
              empty cell returns u'', but cell with spaces returns u'  '
    Raises:
       None
    """
    if path.endswith('.xls'):
        import xlrd
        wbobj = xlrd.open_workbook(path)
        # convert sheet index to sheet name
        if type(sheet) in [int]:
            sheets = wbobj.sheet_names()
            sheet = sheets[sheet]
        else:
            sheet = sheet
        sheetobj = wbobj.sheet_by_name(sheet)

        if r == None:
            r = list(range(0,sheetobj.nrows))     # all rows
        elif type(r) in [int]:
            r = [r]     # a single row
        elif type(r) in [list]:
            if len(r) == 1:
                r = list(range(r[0],sheetobj.nrows))  # from r to end
            else:
                r = r   # multiple rows
        else:
            raise Exception('Invalid row number(s)')

        if c == None:
            c = list(range(0,sheetobj.ncols))     # all cols
        elif type(c) in [int]:
            c = [c]    # a single col
        elif type(c) in [list]:
            if len(c) == 1:
                c = list(range(c[0],sheetobj.ncols))  # from c to end
            else:
                c = c   # multiple cols
        else:
            raise Exception('Invalid col number(s)')

        if len(r)==1 and len(c)==1:
            result = sheetobj.cell_value(r[0],c[0])
        elif len(r)==1 and len(c)>1:
            result = sheetobj.row_values(r[0])
            result = [result[i] for i in c]
        elif len(r)>1 and len(c)==1:
            result = sheetobj.col_values(c[0])
            result = [result[i] for i in r]
        elif len(r)>1 and len(c)>1:
            result = []
            for rr in r:
                row = sheetobj.row_values(rr)
                row = [row[i] for i in c]
                result.append(row)
    elif path.endswith('.xlsx'):
        from openpyxl import load_workbook
        wb = load_workbook(path)
        # convert sheet index to sheet name
        if type(sheet) in [int]:
            sheets = wb.sheetnames
            sheet = sheets[sheet]
        else:
            sheet = sheet
        ws = wb[sheet]

        if r == None:
            r = list(range(0,ws.max_row))     # all rows
        elif type(r) in [int]:
            r = [r]     # a single row
        elif type(r) in [list]:
            if len(r) == 1:
                r = list(range(r[0],ws.max_row))  # from r to end
            else:
                r = r   # multiple rows
        else:
            raise Exception('Invalid row number(s)')

        if c == None:
            c = list(range(0,ws.max_column))     # all cols
        elif type(c) in [int]:
            c = [c]    # a single col
        elif type(c) in [list]:
            if len(c) == 1:
                c = list(range(c[0],ws.max_column))  # from c to end
            else:
                c = c   # multiple cols
        else:
            raise Exception('Invalid col number(s)')

        # openpyxl is 1 based
        if len(r)==1 and len(c)==1:
            result = ws.cell(row=r[0]+1,column=c[0]+1).value
        elif len(r)==1 and len(c)>1:
            for row in ws.iter_rows(min_row=r[0]+1,max_row=r[0]+1,values_only=True):
                result = list(row)
            result = [result[i] for i in c]
        elif len(r)>1 and len(c)==1:
            for col in ws.iter_cols(min_col=c[0]+1,max_col=c[0]+1,values_only=True):
                result = list(col)
            result = [result[i] for i in r]
        elif len(r)>1 and len(c)>1:
            result = []
            for row in ws.iter_rows(min_row=min(r)+1,max_row=max(r)+1,min_col=min(c)+1,max_col=max(c)+1,values_only=True):                
                rc = [row[i] for i in c]
                result.append(rc)
    return result

def savex(path, data, header=None, delimiter=",", sheet_name='Sheet1', *args, **kwargs):
    """
    (path, data, header=None, sheet_name='Sheet1', *args, **kwargs)
    Write a list of list to a xlsx (xlsxwriter), xls(xlwt), csv file
    Args:
        path
            the path to the excel file (xlsx/xls), or csv(.csv, comma separated); explicitly specify .xlsx/xls or .csv 
            xls (but not xlsx) may be more compatible with old-version software (spss v20)
            however, xls limited to 256 columns by 65536 rows
            in which case, consider csv as an alternative
            or xlsx which has limits of 16,384 columns by 1,048,576 rows
            csv writing is much faster than xlsx
        data
            a list of list (each sublist is a row)
        header
            a list with each column's name, eg ['company','brand','price']
        delimiter
            only for writing csv, ignored for excel
        Example
            stats = [['subj','average','contrast','hemisphere']] as header
            loop:
                avg = ez.esp0(cmd)
                stats.append([subj,avg[0],con,hem])
    Returns:
        File is overwritten!
        Each element of list could be of different type, resulting in different format for each cell in excel file
        None and empty strings '' are written blank
        Strings and string of number are written as string
        Numbers as Numbers
        Returns Nothing
        For csv file, everything a string
    Raises:
       None
    """
    if path.endswith('.xls'):
        import xlwt
        book = xlwt.Workbook()
        sheet = book.add_sheet(sheet_name)

        if header:
            for c, val in enumerate(header):
                sheet.write(0, c, val, *args, **kwargs)
            header = 1
        else:
            header = 0

        for r, row in enumerate(data,header):
            for c, val in enumerate(row):
                sheet.write(r, c, val, *args, **kwargs)
        book.save(path)
    elif path.endswith('.xlsx'):
        import xlsxwriter
        xbook = xlsxwriter.Workbook(path)
        xsheet = xbook.add_worksheet(sheet_name)

        if header:
            for c, val in enumerate(header):
                xsheet.write(0, c, val, *args, **kwargs)
            header = 1
        else:
            header = 0

        for r, row in enumerate(data,header):
            # https://xlsxwriter.readthedocs.io/worksheet.html#worksheet-write-row
            xsheet.write_row(r, 0, row, *args, **kwargs)
        xbook.close()
    else:
        import csv
        with open(path,"w") as f:
            wr = csv.writer(f, delimiter=delimiter, *args, **kwargs)
            if header: wr.writerow(header)
            wr.writerows(data)
writex = savex

def savexlist(xlist,file='Data.xlsx',engine='openpyxl',date_format=None,datetime_format=None,mode='w',*args,**kwargs):
    """
    xlist: a list of pandas df
    *args,**kwargs pass to pandas.DataFrame.to_excel (sheet_name ignored, instead auto: Sheet1, Sheet2)
    engine "has" to be openpyxl
    """
    import pandas as pd
    with pd.ExcelWriter(file,engine=engine,date_format=date_format,datetime_format=datetime_format,mode=mode) as writer:
        for i in range(0,len(xlist)):
            xlist[i].to_excel(writer, sheet_name=f'Sheet{i+1}',*args,**kwargs)
writexlist=savexlist

def readp(*args,**kwargs):
    import pandas as pd
    return pd.read_parquet(*args,**kwargs)

def readplist(plist,ignore_index=False,*args,**kwargs):
    """
    plist: [list of parquet files]
    ignore_index: internally read each parquet file and then pd.concat
    """
    import pandas as pd
    res = []
    for p in plist:
        res.append(pd.read_parquet(p,*args,**kwargs))
    res = pd.concat(res,ignore_index=ignore_index)
    return res

def mergep(path,filter,filename,infer_date_from_filename=False,ignore_index=False):
    """
    path filter: passed to ez.ls(path,filter)
    filename: passed to ez.jp(path,filename) as final parquet file name
    infer_date_from_filename: T/F, parse movers_2021-12-21.par (split by _)
    ignore_index: passed to readplist internally when not infer date from filename
    return fs (for easy removal)
    """
    import pandas as pd
    fs = ls(path,filter,full=True)
    if infer_date_from_filename:
        ps = pd.DataFrame()
        for f in fs:
            [_,fn,ext] = sp(f)
            date = fn.split('_')[1]
            date = pd.to_datetime(date)
            p = readp(f).assign(date=date)
            ps = pd.concat([ps,p])
        ps = ps.set_index('date')
    else:
        ps = readplist(fs,ignore_index=ignore_index)
    ps.to_parquet(jp(path,filename)) 
    return fs

def savep(df,*args,**kwargs):
    return df.to_parquet(*args,**kwargs)
writep = savep

def savek(obj,file='data.pickle',*args,**kwargs):
    import pickle
    with open(file, 'wb') as handle:
        pickle.dump(obj, handle, protocol=pickle.HIGHEST_PROTOCOL, *args, **kwargs)
writek = savek

def readk(file='data.pickle',*args, **kwargs):
    import pickle
    with open(file, 'rb') as handle:
        obj = pickle.load(handle,*args,**kwargs)
    return obj

def readj(file='data.json',*args,**kwargs):
    import json
    with open(file, 'r') as f:
        res=json.loads(f.read(),*args,**kwargs)
    if isinstance(res, dict): res=JDict(res)
    return res

def savej(x,file='data.json',indent=4,*args,**kwargs):
    """
    save a dictionary, a list, a value
    """
    import json
    with open(file, 'w') as f:
        json.dump(x,f,indent=indent,*args,**kwargs)
writej = savej

def savet(lines,file='data.txt',*args,**kwargs):
    """
    save a list of lines to text
    """
    with open(file, 'w') as f: 
        f.writelines(lines)
writet = savet

def savew(lines,file='data.docx',heading=None,heading_level=1,*args,**kwargs):
    """
    save a list of line to word
    heading: optional heading immediately before the lines
    heading_level: 0-9 0=title, 1=h1, 2=h2 etc
    """
    # https://www.geeksforgeeks.org/working-with-text-in-python-docx-module/
    import docx
    doc = docx.Document()
    if heading: doc.add_heading(heading,heading_level)
    # https://stackoverflow.com/a/29312199/2292993
    lines = re.sub(r'[^\x00-\x7F]+|\x0c',' ', lines) # remove all non-XML-compatible characters
    for l in lines:
        doc.add_paragraph(l.rstrip('\n'))
    doc.save(file)
writew = savew

def hanzifreq(filename, size=10, outfile=None, encoding='utf8'):  
    """
    (filename, size, outfile, encoding)
    size: top # items; if None, print(save) all items
    outfile: path; if None, no save
    does not count punctuation (ie, ignore punctuation)
    print(save) freq table (top # items) and return a list [(char,times)]  (always all items)
    
    see also http://www.tagxedo.com/ (support Chinese characters)
    """
    # modified from http://blog.csdn.net/xm1331305/article/details/8090639

    import codecs  
    from time import time  
    from operator import itemgetter  
    import sys

    begin = time()

    count = {}  
    for line in codecs.open(filename, 'r', encoding):  
        for chr in line:  
            if '\u4E00' <= chr <= '\u9FA5' or  '\uF900' <= chr <= '\uFA2D':  
                count[chr] = 1 + count.get(chr, 0)

    hanzifreq = sorted(iter(count.items()), key=itemgetter(1), reverse=True)
    result = hanzifreq  # for return
    if size: hanzifreq = hanzifreq[:size]

    print('\n'.join(['%s\t%s' % (chr, times) for chr, times in hanzifreq])) 
    if outfile:
        with codecs.open(outfile, mode='w', encoding='utf-8') as outFile:
            outFile.write('\n'.join(['%s,%s' % (chr, times) for chr, times in hanzifreq]))
    
    print('Done! Elapsed %s seconds.' % (time()-begin))
    return result

def chunk(xs, n):
    '''
    Split the list, xs, into n chunks. 
    gives exactly n chunks, but they may not be evenly sized, since the last chunk gets padded with any surplus.
    eg, 
    chunk('abcdefghij', 3) ->     ['abc', 'def', 'ghij']
    chunk([1,2,3,4],3)     ->     [[1], [2], [3, 4]]
    chunk([1,2,3,4,5,6],6) ->     [[1], [2], [3], [4], [5], [6]]
    chunk([1,2,3,4,5,6],1) ->     [[1, 2, 3, 4, 5, 6]]
    '''
    # http://wordaligned.org/articles/slicing-a-list-evenly-with-python
    L = len(xs)
    assert 0 < n <= L
    s, r = divmod(L, n)
    chunks = [xs[p:p+s] for p in range(0, L, s)]
    chunks[n-1:] = [xs[-r-s:]]
    return chunks

def chunk2(initialList, chunkSize):
    """
    This function chunks a list into sub lists 
    that have a length equals to chunkSize.

    Example:
    lst = [3, 4, 9, 7, 1, 1, 2, 3]
    print(chunk2(lst, 3)) 
    returns
    [[3, 4, 9], [7, 1, 1], [2, 3]]
    """
    # https://stackoverflow.com/a/28786255/2292993
    finalList = []
    for i in range(0, len(initialList), chunkSize):
        finalList.append(initialList[i:i+chunkSize])
    return finalList

def pinyinauthor(names):
    """
    convert Yuan Zhang Zhu, Jian Zhu --> Zhu, Y. Z., Zhu, J. for APA reference list
    can auto get rid of numbers in the author names
    """
    names = re.sub('(?<=[(%s)])(%s)*|^(%s)+|(%s)+$' % ('\s','\s','\s','\s'), '', names, count=0)
    names = re.sub('\d', '', names, count=0)
    names = re.sub(',{2,}', ',', names, count=0)
    names = names.split(', ')
    newnames = ''; i = 0
    for name in names: 
        i += 1
        name = name.split(' ')
        # name.reverse()
        # move family to first
        name.insert(0,name.pop())
        if i==len(names):
            newnames = newnames + name[0].title() + ', ' + '. '.join([x[0].upper() for x in name[1:]]) + '.'
        else:
            newnames = newnames + name[0].title() + ', ' + '. '.join([x[0].upper() for x in name[1:]]) + '., '
    # apa: Journal article with DOl, more than seven authors
    newnames = newnames.split('., ')
    if len(newnames) > 7:
        newnames = newnames[0:6] + ['. . . '+newnames[-1]]
    newnames = '., '.join(newnames)
    return newnames

def xlcolconv(col):
    """
    Excel-style column name converter: 
    Number from or to letters, e.g., 1 = A, 26 = Z, 27 = AA, 703 = AAA.
    input a number (27), output letters ('AA')
    input letters ('AA'), output number (27)
    """
    if type(col) in [int]:
        n = col
        letters = ''
        while n > 0:
            n, r = divmod (n - 1, 26)
            letters = chr(r + ord('A')) + letters
        return letters
    else:
        letters = col
        letters = letters.upper()
        n = 0
        for c in letters:
            n = n * 26 + 1 + ord(c) - ord('A')
        return n

def encoding_detect(character_or_file_path, n_lines=20):
    """
    try to guess a character's or file's encoding
    character_or_file_path: a character string, or a file path
    n_lines: only reads the first few lines, useful for big file, n_lines > total lines in file works fine. Ignored if pass a character
    returns 'utf-8', 'ascii' etc, or None if could not detect
    """
    import chardet
    if exists(character_or_file_path): 
        # Open the file as binary data
        with open(file_path, 'rb') as f:
            # Join binary lines for specified number of lines
            rawdata = b''.join([f.readline() for _ in range(n_lines)])
    else:
        rawdata = str.encode(character_or_file_path)
    result = chardet.detect(rawdata)['encoding']
    return result

def encoding_convert(file_path, source_encoding=None, backup=False):
    """
    convert file to utf-8 (new file is the same file_path)
    source_encoding: if not provided, will guess. if guess fails, print out message
    backup: T/F. if true, create file_path.bak
    """
    if backup:
        import shutil
        shutil.copyfile(file_path, file_path + '.bak')

    # do not use codecs.open or io.open, simply supports encoding parameter in one function
    # manual decoding/encoding
    with open(file_path, 'rb') as f: 
        content = f.read()

    if source_encoding is None:
        import chardet
        source_encoding = chardet.detect(content)['encoding']

    if source_encoding is None:
        print("I cannot guess source encoding (try GB2312 for MS-Windows?)", file_path)
        return None
    else:
        print(source_encoding, file_path)

    target_encoding = 'utf-8'
    if source_encoding != target_encoding:
        content = content.decode(source_encoding)
        with open(file_path, 'wb') as f: 
            f.write(content.encode(target_encoding))

def opens(filepath):
    """
    filepath: a str of file or list of files
    opens with default program in Mac, Win, Linux
    """
    if type(filepath) in [list]:
        for f in filepath:
            opens(f)
    else:
        import subprocess, os, sys
        if sys.platform.startswith('darwin'):
            subprocess.call(('open', filepath))
        elif os.name == 'nt': # For Windows
            os.startfile(filepath)
        elif os.name == 'posix': # For Linux
            subprocess.call(('xdg-open', filepath))

def vx(df):
    import tempfile
    _, tmp = tempfile.mkstemp(suffix='.xlsx')
    df.to_excel(tmp)
    opens(tmp)

def getkmvar(var):
    # https://wiki.keyboardmaestro.com/action/Execute_a_Shell_Script
    # Keyboard Maestro only allows Get (cannot set) env variable from shell script
    # eg., getkmvar('PathFinderSelection')
    # returns stripped string, or none if not exist or empty in user input in KM
    import os
    var = 'KMVAR_'+var
    try:
        value = os.environ[var].strip()
    except KeyError:
        value = None
    return value

def setkmvar(var,value):
    """
    var: KM variable in gui, e.g., MouseLoc
    value: a string or number--number will be converted to a string in KMVAR(?)
    Keyboard Maestro does not allow to set env variable from shell script, but applescript would be OK
    """
    # https://wiki.keyboardmaestro.com/manual/Scripting
    applescript = f"""
    tell application "Keyboard Maestro Engine"
      -- set calcResult to getvariable "Calculation Result"
      -- If the Keyboard Maestro Variable does not exist, the AppleScript Variable will be set to empty string
     
      setvariable "{var}" to "{value}"
      -- If the Keyboard Maestro Variable does not exist, it will be created
    end tell
    """
    espA(applescript)
    return None

def runkm(uuid,timeout=3600):
    """
    run km macro
    uuid: Edit -> Copy macro as uuid
    timeout: seconds (Applescripts default timeout of 2 mins)
    """
    # https://stackoverflow.com/a/22446535/2292993
    applescript = f"""
    tell application "Keyboard Maestro Engine"
        with timeout of {timeout} seconds
            do script "{uuid}"
        end timeout
    end tell
    """
    espA(applescript)
    return None

def msgbox(msg):
    """
    a trick to call Keyboard Maestro to display large text
    """
    setkmvar('PyMsg',msg)
    runkm("7D38B1AD-720C-491F-8C07-2C33D9B2A958")
    return None

def menubar(msg):
    """
    a trick to call Keyboard Maestro to display menubar text
    """
    setkmvar('KM_MENUBAR_TEXT',msg)
    return None

def setpassword(acount,password):
    """
    under Keychains-->login-->km.py
    """
    import keyring
    keyring.set_password("km.py", acount, password)

def getpassword(acount):
    """
    under Keychains-->login-->km.py
    """
    import keyring
    return keyring.get_password("km.py", acount)

def office_pptx_replace_font(pptx_path: str, replacedict: dict, save_path=None):
    """
    partially working so far, it will replace any font, regardless of the specified condition
    replacedict: {'Arial':'Garamond', 'Any':'Garamond'} auto translated to regex
    save_path: if None, overwrite the original file, style also kept
    requires pip install python-pptx    
    Currently not implemented for charts or graphics.
    office_pptx_replace('input.pptx',{'Any':'Garamond'},'output.pptx') 
    """
    # https://stackoverflow.com/questions/37924808
    import re
    from pptx import Presentation
    from pptx.util import Pt
    prs = Presentation(pptx_path)

    # todo: change font size also not working...

    # todo: prs.slide_masters only one master!
    # slide.slide_layout.slide_master.placeholders also not working...
    # for master in prs.slide_masters:
    #     for shape in master.placeholders:
    #         for key, replace in replacedict.items():
    #             regex = re.compile(key);
    #             if shape.has_text_frame:
    #                 text_frame = shape.text_frame
    #                 for paragraph in text_frame.paragraphs:
    #                     for run in paragraph.runs:
    #                         # todo: I cannot make the condition if... to work, or python-pptx not work?
    #                         run.font.name = replace

    for slide in prs.slides:
        for shape in slide.shapes:
            for key, replace in replacedict.items():
                regex = re.compile(key)
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            # todo: I cannot make the condition if... to work, or python-pptx not work?
                            run.font.name = replace
                            # if run.font.name is not None:
                            #     # font = run.font
                            #     # font.name = 'Calibri'
                            #     # font.size = Pt(18)
                            #     # font.bold = True
                            #     # font.italic = None  # cause value to be inherited from theme
                            #     # font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
                            #     # # If you prefer, you can set the font color to an absolute RGB value. Note that this will not change color when the theme is changed:
                            #     # font.color.rgb = RGBColor(0xFF, 0x7F, 0x50)
                            #     # # A run can also be made into a hyperlink by providing a target URL:
                            #     # run.hyperlink.address = 'https://github.com/scanny/python-pptx'
                            #     if regex.search(run.font.name) or key=='Any':
                            #         run.font.name = replace
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            text_frame = cell.text_frame
                            for paragraph in text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.name = replace
                                    # if run.font.name is not None:
                                    #     if regex.search(run.font.name) or key=='any':
                                    #         run.font.name = replace

    if save_path is None: save_path=pptx_path
    prs.save(save_path)

def office_pptx_replace(pptx_path: str, replacedict: dict, save_path=None):
    """
    replacedict: {'src':'replace','src2':'replace2'} auto translated to regex
    save_path: if None, overwrite the original file, style also kept
    requires pip install python-pptx    
    Currently not implemented for charts or graphics.
    office_pptx_replace('input.pptx',{'string to replace': 'replacement text'},'output.pptx') 
    """
    # https://stackoverflow.com/questions/37924808
    import re
    from pptx import Presentation
    prs = Presentation(pptx_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            for key, value in replacedict.items():
                regex = re.compile(key); replace = value
                if shape.has_text_frame:
                    if regex.search(shape.text):
                        text_frame = shape.text_frame
                        for paragraph in text_frame.paragraphs:
                            for run in paragraph.runs:
                                cur_text = run.text
                                new_text = regex.sub(replace,cur_text)
                                run.text = new_text
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if regex.search(cell.text):
                                new_text = regex.sub(replace,cell.text)
                                cell.text = new_text

    if save_path is None: save_path=pptx_path
    prs.save(save_path)

def office_docx_replace(docx_path, replacedict, save_path=None):
    """
    replacedict: {'src':'replace','src2':'replace2'} auto translated to regex
    save_path: if None, overwrite the original file, style also kept
    requires pip install python-docx
    """
    import re
    from docx import Document
    # https://stackoverflow.com/a/42829667/2292993
    # doc_obj change in place?
    def _docx_replace_regex(doc_obj, regex, replace):
        for p in doc_obj.paragraphs:
            if regex.search(p.text):
                inline = p.runs
                # Loop added to work with runs (strings with same style)
                for i in range(len(inline)):
                    if regex.search(inline[i].text):
                        text = regex.sub(replace, inline[i].text)
                        inline[i].text = text

        for table in doc_obj.tables:
            for row in table.rows:
                for cell in row.cells:
                    _docx_replace_regex(cell, regex, replace)

    doc = Document(docx_path)
    for key,value in replacedict.items():
        regex = re.compile(key)
        replace = value
        _docx_replace_regex(doc, regex, replace)
    if save_path is None: save_path=docx_path
    doc.save(save_path)

def office_pdf_from(inputfile,outputdir=None):
    """
    call libreoffice headless mode to convert to pdf
    inputfile: single file, or *.docx
    outputdir: if none, same directory as inputfile
    outputfile will be the same name with .pdf
    some format will be different due to libreoffice's rendering
    """
    import sys, os, subprocess

    # TODO: Provide support for more platforms
    if sys.platform == 'darwin':
        libreoffice=os.path.expanduser('~/Dropbox/Apps/PDF/LibreOffice.app/Contents/MacOS/soffice')
    else:
        libreoffice='libreoffice'

    if outputdir is None: outputdir=splitpath(inputfile)[0]
    args = [libreoffice, '--headless', '--convert-to', 'pdf', '--outdir', outputdir, inputfile]
    # subprocess.run(args, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=timeout)
    esp(' '.join(args))

def office_doc2docx(inputfile,outputdir=None):
    """
    call libreoffice headless mode to convert to docx
    inputfile: single file, or *.doc
    outputdir: if none, same directory as inputfile
    outputfile will be the same name with .pdf
    some format will be different due to libreoffice's rendering
    """
    import sys, os, subprocess

    # TODO: Provide support for more platforms
    if sys.platform == 'darwin':
        libreoffice=os.path.expanduser('~/Dropbox/Apps/PDF/LibreOffice.app/Contents/MacOS/soffice')
    else:
        libreoffice='libreoffice'

    if outputdir is None: outputdir=splitpath(inputfile)[0]
    args = [libreoffice, '--headless', '--convert-to', 'docx', '--outdir', outputdir, inputfile]
    # subprocess.run(args, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=timeout)
    esp(' '.join(args))

def office_pdf_unlock(inputpdfs):
    """
    inputpdfs: ['pdf1','pdf2'], or 'pdf1 pdf2', or 'pdf1'
    new file in the same folder suffixed _compressed
    less effective than acrobat pro's Save Reduced Size
    """
    if type(inputpdfs) not in [list]:
        # hack for filename with space
        # assume "quoted form of " applescript returns single quote
        inputpdfs = inputpdfs.split("' '")  
        # strip "" from both ends that might be present from keyboardmaestro
        inputpdfs = [e.strip("'").strip('"') for e in inputpdfs]
    for pdf in inputpdfs:
        [path,file,ext]=splitpath(pdf)
        outputpdf=joinpath(path,file+'_unlocked'+ext)

        # # windows version
        # gswin32c -dSAFER -dBATCH -dNOPAUSE -sDEVICE=pdfwrite -sFONTPATH=%windir%/fonts;xfonts;. -sPDFPassword= -dPDFSETTINGS=/prepress -dPassThroughJPEGImages=true -sOutputFile=OUTPUT.pdf INPUT.pdf
        cmd = f"/usr/local/bin/gs -q -dNOPAUSE -dBATCH -sDEVICE=pdfwrite -sOutputFile='{outputpdf}' -c .setpdfwrite -f '{pdf}'"
        execute(cmd, False)

        # # to download ghostscript:
        # mac: https://pages.uoregon.edu/koch/
        # windows: https://www.ghostscript.com/download/gsdnld.html
        # +++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++
        # internal help
        # +++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++
        # http://www.verypdf.com/app/pdf-password-remover-mac/index.html
        # An owner password is always used for control the printing, copying, modifying permissions of PDF. 
        # With removing the owner password, you will directly access these permissions of the PDF.
        # A user password is to control the permission of opening a PDF. 
        # If there is a user password, it is needed to unlock a pdf to begin with.

        # use ghostscript to remove owner password
        # http://www.commandlinefu.com/commands/view/4345/remove-security-limitations-from-pdf-documents-using-ghostscript
        # http://www.localizingjapan.com/blog/2013/02/23/unlocking-secured-password-protected-pdf-files/

        # "$*" All of the positional parameters, seen as a single word
        # "$@" Same as $*, but each parameter is a quoted string, that is, 
        # the parameters are passed on intact, without interpretation or expansion. 
        # This means, among other things, that each parameter in the argument list is seen as a separate word.

def office_pdf_crop(inputpdfs,parameters=''):
    """
    inputpdfs: ['pdf1','pdf2'], or 'pdf1 pdf2', or 'pdf1'
    new file in the same folder suffixed _cropped (or provide -o outfile option, but currently only works for single pdf input)
    less effective than acrobat pro's Save Reduced Size
    # https://github.com/abarker/pdfCropMargins
    parameters: if empty '', default of retaining 10% of the existing margins (not necessarily implying -u)
                '-u -s -a4 left bottom right top' separated by space. 
                left bottom right top are the margin to crop in unit of bp (72 bp = 1 inch)
                -u Crop all the pages uniformly
                -s force each page to the same size
                -a4 Decrease the margin sizes individually with four absolute offset values
                -gui requires pip install PySimpleGUI
    """
    if type(inputpdfs) not in [list]:
        # hack for filename with space
        # assume "quoted form of " applescript returns single quote
        inputpdfs = inputpdfs.split("' '")  
        # strip "" from both ends that might be present from keyboardmaestro
        inputpdfs = [e.strip("'").strip('"') for e in inputpdfs]
    from pdfCropMargins import crop
    for pdf in inputpdfs:
        # https://github.com/abarker/pdfCropMargins
        # crop defaults to pwd to output the generated file?
        # strip quotes from applescript input (quoted form of filepath)
        path=os.path.split( os.path.abspath(pdf.strip("'").strip('"')) )[0]
        os.chdir(path)
        parameters = parameters.split(' ') if parameters!='' else []
        # not working for kindle dx..., but works for goodreader
        crop(parameters+[pdf])
        # ['-b','m','-b','c','-b','t','-b','a','-b','b']
        # https://www.prepressure.com/pdf/basics/page-boxes
        # When you use the Crop tool in Acrobat to resize a page it displays correctly in Acrobat, however when you place it into another application (eg InDesign) it loads the pdf displaying the original size.
        # When you use the Crop tool and then save the PDF, the cropping is applied like a mask… the dimension can be restored at any time.
        # To make the cropped size permanent… you need to save the page as am EPS file and redistill it.
        # But Briss works for kindle...

def office_pdf_compress(inputpdfs):
    """
    inputpdfs: ['pdf1','pdf2'], or 'pdf1 pdf2', or 'pdf1'
    new file in the same folder suffixed _compressed
    less effective than acrobat pro's Save Reduced Size
    """
    if type(inputpdfs) not in [list]:
        # hack for filename with space
        # assume "quoted form of " applescript returns single quote
        inputpdfs = inputpdfs.split("' '")  
        # strip "" from both ends that might be present from keyboardmaestro
        inputpdfs = [e.strip("'").strip('"') for e in inputpdfs]
    import fitz
    for pdf in inputpdfs:
        doc=fitz.open(pdf)
        [path,file,ext]=splitpath(pdf)
        outputpdf=joinpath(path,file+'_compressed'+ext)
        doc.save(outputpdf,garbage=4,clean=True,deflate=True)
        doc.close()

def office_pdf_merge(inputpdfs,mergeorder=None,outputpdf=None,compress=True):
    """
    inputpdfs: ['pdf1','pdf2'], or 'pdf1 pdf2', or 'pdf1'
    mergeorder: reorder inputpdfs
        # if called from Keyboard Maestro
        # Specify pdf files merge order
        # e.g., '4 3 2 1' (1 based, total numbers should be equal to files selected)
        # Default Empty, no particular order assigned
    outputpdf: path for merged file. if None, auto name
    compress: compress merged file
    requires pip install PyMuPDF
    """
    if type(inputpdfs) not in [list]:
        # hack for filename with space
        # assume quoted form
        inputpdfs = inputpdfs.split("' '")  
        # strip "" from both ends that might be present from keyboardmaestro
        inputpdfs = [e.strip("'").strip('"') for e in inputpdfs]
    import fitz

    if mergeorder is not None: 
        mergeorder = mergeorder.strip("'").strip('"').split(' ')
        if len(mergeorder)>1:
            mergeorder = [int(e)-1 for e in mergeorder]
            inputpdfs = [inputpdfs[e] for e in mergeorder]

    pdf1 = inputpdfs[0]
    doc1 = fitz.open(pdf1)         # must be a PDF
    for pdf2 in inputpdfs[1:]:
        doc2 = fitz.open(pdf2)                 # must be a PDF
        pages1 = len(doc1)                     # save doc1's page count
        toc1 = doc1.getToC(simple=False)        # save TOC 1
        toc2 = doc2.getToC(simple=False)        # save TOC 2
        # https://pymupdf.readthedocs.io/en/latest/document/#Document.getToC
        # [lvl, title, page, dest] page is 1 based
        # dest – (dict) included only if simple=False. Contains details of the link destination.
        # dest contains additional fine position of bookmark
        if len(toc1)==0: toc1=[[1,splitpath(pdf1)[1],1]]
        if len(toc2)==0: toc2=[[1,splitpath(pdf2)[1],1]]
        doc1.insertPDF(doc2)                   # doc2 at end of doc1
        for t in toc2:                         # increase toc2 page numbers
            t[2] += pages1                     # by old len(doc1)
        doc1.setToC(toc1 + toc2)               # now result has total TOC

    if outputpdf is None: outputpdf=joinpath(splitpath(pdf1)[0],'AllCombined.pdf')
    # create time zone value in PDF format
    cdate = fitz.getPDFnow()
    pdf_dict = {"creator": "PDF Joiner",
               "producer": "PyMuPDF",
               "creationDate": cdate,
               "modDate": cdate,
               "title": splitpath(outputpdf)[1],
               "author": "Jerry",
               "subject": "",
               "keywords": ""}
    doc1.setMetadata(pdf_dict)      # put in meta data

    # https://pymupdf.readthedocs.io/en/latest/document/#Document.save
    # doc.save cannot overwrite
    if exists(outputpdf): rm(outputpdf)
    if compress:
        doc1.save(outputpdf,garbage=4,clean=True,deflate=True)
    else:
        doc1.save(outputpdf)
    doc1.close()

def office_pdf_autoname(inputpdfs):
    """
    inputpdfs: ['pdf1','pdf2'], or 'pdf1 pdf2', or 'pdf1'
    new file in the same folder: year_author_journal_title
    if cannot rename for whatever reasons, silently skip
    """
    if type(inputpdfs) not in [list]:
        # hack for filename with space
        # assume quoted form
        inputpdfs = inputpdfs.split("' '")  
        # strip "" from both ends that might be present from keyboardmaestro
        inputpdfs = [e.strip("'").strip('"') for e in inputpdfs]
    import fitz
    for pdf in inputpdfs:
        try:
            doc=fitz.open(pdf)

            meta = doc.metadata
            #) year
            year = meta['creationDate']
            if year is None:
                year = ''
            else:
                # https://stackoverflow.com/a/40916306/2292993
                # translate/remap/replace from '' to '', and delete string.punctuation
                year = (year[2:6] if ('D:' in year) else year[-4:]).translate(str.maketrans('','', string.punctuation))
            #) author
            author = meta['author']
            if author is None: 
                author = ''
            else:
                # last name
                author = author.strip(' ').split(' ')[-1].translate(str.maketrans('','', string.punctuation))
                author = author.encode('ascii',errors='ignore').decode('utf-8')  # some weird author names, get rid of nonascii

            #) journal
            journal = meta['subject']
            if journal is None: 
                journal = ''
            else:
                journal = journal.strip(' ').split(',')[0].translate(str.maketrans('','', '!"#$%&\'()*+,-./:;<=>?@[\\]^`{|}~')).replace(' ','_')
                journal = journal.encode('ascii',errors='ignore').decode('utf-8')
                journal = journal[:20] if len(journal) > 20 else journal

            #) title
            output = JDict()
            # first 4 pages
            i = 0; j = min(3,len(doc))
            for page in doc:
                # text = page.getText("text")
                html_text = page.getText("html")
                from bs4 import BeautifulSoup
                soup = BeautifulSoup(html_text,"lxml")
                # https://stackoverflow.com/a/39016902/2292993
                patt = re.compile("font-size:(\d+)")
                for tag in soup.select("[style*=font-size]"):
                    output.update({float(patt.search(tag["style"]).group(1)): [tag.text.strip()]})
                i += 1
                if i>j: break
            # largest font size (title may have multiple lines, therefore join)
            largestfontsize = max(output.keys())
            # ignore if too large (may be a watermark)
            if largestfontsize>40: output.pop(largestfontsize)
            title = ' '.join(output[max(output.keys())])
            title = title.translate(str.maketrans('','', '!"#$%&\'()*+,-./:;<=>?@[\\]^`{|}~')).replace(' ','_')
            title = title.encode('ascii',errors='ignore').decode('utf-8')
            title = title[:20] if len(title) > 20 else title

            fullName = pdf
            newName = '_'.join([year,author,journal,title]) + '.pdf'
            if not newName[0].isdigit(): newName = Moment().datetime + '.pdf'
            newFullName = joinpath(splitpath(fullName)[0],newName)
            if (newName=='___.pdf' or ('%' in newName) or exists(newFullName)): continue
            print( sprintf('Rename: %s --> %s', splitpath(fullName)[1], splitpath(newFullName)[1]) )
            os.rename(fullName, newFullName)
        except:
            continue

def applescript_pages_replace(searchWord, replacementString):
    """
    >>> depricated: use docx_replace <<<

    replace one word at a time (but all occurrences) in pages active document
    For the best result, searchWord to be recommended using all capitalized "DRFULLNAME"
    case sensitive, can NOT have _, word bounded by spaces, just one word, not two words or more
    this way, no recursive replacement would happen, other issues could be avoided
    ("DRFULLNAME", "My name is")
    """
    applescript = '''
    -- does not work for symbols, eg, {  } or word pairs separated by separator
    -- https://iworkautomation.com/pages/body-text-replace.html
    on pagesReplace(searchWord, replacementString)
        my replaceWordWithStringInBodyText(searchWord, replacementString)
    end 

    on replaceWordWithStringInBodyText(searchWord, replacementString)
        tell application "Pages"
            reopen
            activate
            tell the front document
                tell body text
                    -- start at the end and go to the beginning
                    repeat with i from the (count of paragraphs) to 1 by -1
                        tell paragraph i
                            repeat
                                try
                                    if exists searchWord then
                                        set (last word where it is searchWord) to replacementString
                                    else
                                        exit repeat
                                    end if
                                on error errorMessage
                                    exit repeat
                                end try
                            end repeat
                        end tell
                    end repeat
                end tell
            end tell
            -- return true
        end tell
    end replaceWordWithStringInBodyText

    my pagesReplace("%(searchWord)s","%(replacementString)s")
    '''
    def myesp(cmdString):
        import os, inspect, tempfile, subprocess
        caller = inspect.currentframe().f_back
        cmd =  cmdString % caller.f_locals
        
        fd, path = tempfile.mkstemp(suffix='.applescript')
        try:
            with os.fdopen(fd, 'w') as tmp:
                tmp.write(cmd.replace('"','\"').replace("'","\'")+'\n\n')
            subprocess.call('osascript ' + path, shell=True, executable="/bin/bash")
        finally:
            os.remove(path)
        return None
    myesp(applescript)

def applescript_pages_pdfactive():
    """
    >>> depricated: use libreoffice headless <<<
    convert current/active Pages document to pdf to the same folder as the active document, quit Pages after conversion (will save or remind to save)
    pdf name is same as page document name, if pdf exists, overwrites
    if active document does not have a valid folder (eg, never saved), save to download folder
    """
    applescript = '''
use framework "Foundation"
use scripting additions   
------------------------------------------------------------------------------
# Auth: Christopher Stone
# dCre: 2017/05/22 15:10
# dMod: 2017/05/22 15:18 
# Appl: System Events
# Task: Attempt to acquire the POSIX Path of the file associated with the front window of the front app.
# Libs: None
# Osax: None
# Tags: @Applescript, @Script, @ASObjC, @System_Events, @Acquire, @POSIX_Path, @Path, @File, @Associated, @Front, @Window, @Front_App
------------------------------------------------------------------------------
tell application "Pages"
    activate
end 
-- have to wait till Pages frontmost, otherwise the code would fail
delay 3
try
    tell application "System Events"
        set frontProcess to first process whose frontmost is true
        set procName to name of frontProcess
        tell frontProcess
            tell front window
                tell attribute "AXDocument"
                    set fileUrlOfFrontWindow to its value
                end tell
            end tell
        end tell
    end tell
    
on error
    set fileUrlOfFrontWindow to missing value
end try

if fileUrlOfFrontWindow is not missing value and fileUrlOfFrontWindow is not "file:///Irrelevent" then
    set posixPath to (current application's class "NSURL"'s URLWithString:fileUrlOfFrontWindow)'s |path|() as text
    -- parent folder instead of full path name
    set the defaultDestinationFolder to POSIX file (posixPath & "/..")
else
    set the defaultDestinationFolder to (path to downloads folder)
end if
------------------------------------------------------------------------------
    -- https://iworkautomation.com/pages/document-export.html
    -- also https://github.com/gitmalong/pages2docx
    -- https://gist.github.com/loop/7207134ed7ff7a288ee1
    on pagesCurrent2pdf(defaultDestinationFolder)
    --property exportFileExtension : "pdf"
    --property useEncryptionDefaultValue : false    

    set usePDFEncryption to false
    tell application "Pages"
        activate
        try
            if not (exists document 1) then error number -128
            
            if usePDFEncryption is true then
                -- PROMPT FOR PASSWORD (OPTIONAL)
                repeat
                    display dialog "Enter a password for the PDF file:" default answer "" buttons {"Cancel", "No Password", "OK"} default button 3 with hidden answer
                    copy the result to {button returned:buttonPressed, text returned:firstPassword}
                    if buttonPressed is "No Password" then
                        set usePDFEncryption to false
                        exit repeat
                    else
                        display dialog "Enter the password again:" default answer "" buttons {"Cancel", "No Password", "OK"} default button 3 with hidden answer
                        copy the result to {button returned:buttonPressed, text returned:secondPassword}
                        if buttonPressed is "No Password" then
                            set usePDFEncryption to false
                            exit repeat
                        else
                            if firstPassword is not secondPassword then
                                display dialog "Passwords do no match." buttons {"Cancel", "Try Again"} default button 2
                            else
                                set providedPassword to the firstPassword
                                set usePDFEncryption to true
                                exit repeat
                            end if
                        end if
                    end if
                end repeat
            end if
            
            -- DERIVE NAME AND FILE PATH FOR NEW EXPORT FILE
            set documentName to the name of the front document
            if documentName ends with ".pages" then set documentName to text 1 thru -7 of documentName
            
            tell application "Finder"
                set exportItemFileName to documentName & "." & "pdf"
                -- set incrementIndex to 1
                -- repeat until not (exists document file exportItemFileName of defaultDestinationFolder)
                --     set exportItemFileName to documentName & "-" & (incrementIndex as string) & "." & "pdf"
                --     set incrementIndex to incrementIndex + 1
                -- end repeat
            end tell
            set the targetFileHFSPath to (defaultDestinationFolder as string) & exportItemFileName

            tell application "System Events"
                if exists file targetFileHFSPath then
                    set script1 to "rm '" & (POSIX path of targetFileHFSPath) & "'"
                    do shell script script1
                end if
            end tell
            
            -- EXPORT THE DOCUMENT
            with timeout of 1200 seconds
                if usePDFEncryption is true then
                    export front document to file targetFileHFSPath as PDF with properties {password:providedPassword} with replacing
                else
                    export front document to file targetFileHFSPath as PDF with replacing
                end if
            end timeout
            
        on error errorMessage number errorNumber
            if errorNumber is not -128 then
                display alert "EXPORT PROBLEM" message errorMessage
            end if
            error number -128
        end try
    end tell

    tell application "Pages"
        quit
        -- or alternatively,
        -- quit without saving
    end tell

    -- -- SHOW THE NEW PDF FILE
    -- tell application "Finder"
    --     activate
    --     reveal document file targetFileHFSPath
    -- end tell
    end
    my pagesCurrent2pdf(defaultDestinationFolder)
    '''
    def myesp(cmdString):
        import os, inspect, tempfile, subprocess
        caller = inspect.currentframe().f_back
        cmd =  cmdString % caller.f_locals
        
        fd, path = tempfile.mkstemp(suffix='.applescript')
        try:
            with os.fdopen(fd, 'w') as tmp:
                tmp.write(cmd.replace('"','\"').replace("'","\'")+'\n\n')
            subprocess.call('osascript ' + path, shell=True, executable="/bin/bash")
        finally:
            os.remove(path)
        return None
    myesp(applescript)

def applescript_preview_moveactive(theFolder,quitpreview=0):
    """
    (theFolder)
    save the current Preview pdf, move the pdf to theFolder, quit Preview or reopen the pdf
    """
    # https://stackoverflow.com/a/16071855/2292993
    if quitpreview:
        applescript = '''
        on previewCurrentMove(theFolder)
            tell application "Preview"
                activate
                save front document
                set theFile to ((path of front document) as text)
                -- close front document
                quit
            end tell

            do shell script "mv " & quoted form of theFile & " " & theFolder
        end previewCurrentMove
        my previewCurrentMove("%(theFolder)s")
        '''
    else:
        applescript = '''
        on previewCurrentMove(theFolder)
            tell application "Preview"
                activate
                save front document
                set theFile to ((path of front document) as text)
            end tell

            --known issue: this method would not notify Preview
            do shell script "mv " & quoted form of theFile & " " & theFolder
        end previewCurrentMove
        my previewCurrentMove("%(theFolder)s")
        '''
    def myesp(cmdString):
        import os, inspect, tempfile, subprocess
        caller = inspect.currentframe().f_back
        cmd =  cmdString % caller.f_locals
        
        fd, path = tempfile.mkstemp(suffix='.applescript')
        try:
            with os.fdopen(fd, 'w') as tmp:
                tmp.write(cmd.replace('"','\"').replace("'","\'")+'\n\n')
            subprocess.call('osascript ' + path, shell=True, executable="/bin/bash")
        finally:
            os.remove(path)
        return None
    # https://stackoverflow.com/a/7788702/2292993
    import psutil
    # if preview is running, case sensitive
    if ("Preview" in (p.name() for p in psutil.process_iter())):
        myesp(applescript)

def applescript_finder_alias(theFrom, theTo):
    """
    (theFrom, theTo)
    create a short/alias
    theFrom, theTo: relative or abs path, both folder or both file
    """
    # https://apple.stackexchange.com/questions/51709
    applescript = '''
    tell application "Finder"
       make new alias to %(theType)s (posix file "%(theFrom)s") at (posix file "%(todir)s")
       set name of result to "%(toname)s"
    end tell
    '''
    def myesp(cmdString):
        import os, inspect, tempfile, subprocess
        caller = inspect.currentframe().f_back
        cmd =  cmdString % caller.f_locals
        
        fd, path = tempfile.mkstemp(suffix='.applescript')
        try:
            with os.fdopen(fd, 'w') as tmp:
                tmp.write(cmd.replace('"','\"').replace("'","\'")+'\n\n')
            subprocess.call('osascript ' + path, shell=True, executable="/bin/bash")
        finally:
            os.remove(path)
        return None
    import os
    theFrom = os.path.abspath(theFrom)
    theTo = os.path.abspath(theTo)
    if os.path.isfile(theFrom): 
        theType = 'file'
    else:
        theType = 'folder'
    todir = os.path.dirname(theTo)
    toname = os.path.basename(theTo)
    myesp(applescript)

def applescript_mail(emails,subjectline,titles,body,attaches=[],sendout=0):
    """
    >>> recommend: applescript_outlook <<<

    (emails,subjectline,titles,body,attaches,sendout)
    emails: if Multiple emails, 'a@a.com, b@b.com'  # with or without spaces after , or ;
    subjectline: 
    titles: "Dear Dr. Zhu"  # do not add comma at the end
    body: "\nblabla"        # need to add a newline before the body, internally: titles & "," & body & "\n\n"
    attaches: ['file/path/to/a.pdf','file/path/to/b.pdf'] # posix filepath in a python list # relative or full path
    sendout: 1/0

    note: (still issue with exchange/apple mail even the following configuration)
    In Apple Mail, go to the menu Edit > Attachments
        Make sure the following settings are checked:
        Always Send Windows-Friendly Attachments
        Always Insert Attachments at End of Message
    Otherwise: Exchange creating ATT00001 attachments
    http://kb.mit.edu/confluence/pages/viewpage.action?pageId=4981187
    """
    # workaround for warnings for multiple emails
    import re
    emails = re.split('[,;]',emails)
    emails = ','.join(['"{}"'.format(e.strip()) for e in emails])
    # add double quote
    import os
    attaches = ','.join(['"{}"'.format(os.path.abspath(a)) for a in attaches])
    applescript = '''
    -- emails, attaches is a list {}
    on applemail(emails,subjectline,titles,body,attaches,sendout)
        repeat with i from 1 to count attaches
            set item i of attaches to (POSIX file (item i of attaches) as alias)
        end repeat

        tell application "Mail"
            set theSubject to subjectline
            set theContent to titles & "," & body & "\n\n"
            set theAddress to emails -- the receiver 

            set msg to make new outgoing message with properties {subject: theSubject, content: theContent, visible:true}
            -- tell msg to make new to recipient at end of every to recipient with properties {address:theAddress}
            repeat with i from 1 to count theAddress
                tell msg to make new to recipient at end of every to recipient with properties {address:item i of theAddress}
            end repeat
            repeat with theAttachmentFile in attaches
                tell msg to make new attachment with properties {file name:theAttachmentFile as alias} at after the last paragraph
            end repeat

            delay 3
            if sendout as number is equal to 1 then
                send msg
            else
                activate
            end if
        end tell
    end
    my applemail({%(emails)s}, "%(subjectline)s", "%(titles)s", "%(body)s", {%(attaches)s}, %(sendout)d)
    '''
    def myesp(cmdString):
        import os, inspect, tempfile, subprocess
        caller = inspect.currentframe().f_back
        cmd =  cmdString % caller.f_locals
        
        fd, path = tempfile.mkstemp(suffix='.applescript')
        try:
            with os.fdopen(fd, 'w') as tmp:
                tmp.write(cmd.replace('"','\"').replace("'","\'")+'\n\n')
            subprocess.call('osascript ' + path, shell=True, executable="/bin/bash")
        finally:
            os.remove(path)
        return None
    myesp(applescript)

def applescript_outlook(emails,subjectline,titles,body,attaches=[],sendout=0):
    """
    (emails,subjectline,titles,body,attaches,sendout)
    emails: if Multiple emails, 'a@a.com, b@b.com'  # with or without spaces after , or ;
    subjectline: 
    titles: "Dear Dr. Zhu"  # do not add comma at the end
    body: "\nblabla"        # need to add a newline before the body, internally: titles & "," & body & "\n\n"
                            well, actually outlook defaults html, I replace \n with <br> internally 
                            see https://discussions.apple.com/thread/5929457
    attaches: ['file/path/to/a.pdf','file/path/to/b.pdf'] # posix filepath in a python list # relative or full path
    sendout: 1/0, integer

    no Exchange issue creating ATT00001 attachments with outlook!
    >>> todo: if applescript not working, consider python+outlook <<<
    """
    body = body.replace('\n','<br>')
    # workaround for warnings for multiple emails
    import re
    emails = re.split('[,;]',emails)
    emails = ','.join(['"{}"'.format(e.strip()) for e in emails])
    import os
    # outlook add attachement in reverse order
    attaches.reverse()
    # add double quote
    attaches = ','.join(['"{}"'.format(os.path.abspath(a)) for a in attaches])
    # https://stackoverflow.com/a/35469878/2292993
    # https://stackoverflow.com/a/30900060/2292993
    applescript = '''
    -- emails, attaches is a list {}
    on outlookmail(emails,subjectline,titles,body,attaches,sendout)
        repeat with i from 1 to count attaches
            set item i of attaches to (POSIX file (item i of attaches) as alias)
        end repeat

        tell application "Microsoft Outlook"
            set theSubject to subjectline
            set theContent to titles & "," & body & "<br><br>"
            set theAddress to emails -- the receiver 

            --set msg to make new outgoing message with properties {subject:theSubject, plain text content:theContent}            
            set msg to make new outgoing message with properties {subject:theSubject, content:theContent}            
            repeat with i from 1 to count theAddress
                tell msg to make new recipient with properties {email address:{address:item i of theAddress}} at end of to recipients of msg
            end repeat
            repeat with theAttachmentFile in attaches
                tell msg to make new attachment with properties {file:theAttachmentFile as alias}
            end repeat

            delay 3
            if sendout as number is equal to 1 then
                send msg
            else
                open msg
                activate
            end if
        end tell
    end
    my outlookmail({%(emails)s}, "%(subjectline)s", "%(titles)s", "%(body)s", {%(attaches)s}, %(sendout)d)
    '''
    def myesp(cmdString):
        import os, inspect, tempfile, subprocess
        caller = inspect.currentframe().f_back
        cmd =  cmdString % caller.f_locals
        
        fd, path = tempfile.mkstemp(suffix='.applescript')
        try:
            with os.fdopen(fd, 'w') as tmp:
                tmp.write(cmd.replace('"','\"').replace("'","\'")+'\n\n')
            subprocess.call('osascript ' + path, shell=True, executable="/bin/bash")
        finally:
            os.remove(path)
        return None
    myesp(applescript)
# applescript_outlook('jerryzhujian9@gmail.com; jerryzhujian9@gmail.com','hello','Dear Zhu','\n\nBest,\nJerry',['/Applications/Calculator.app/Contents/version.plist','/Applications/Calculator.app/Contents/Info.plist'],0)

def playfirefoxmacro(macro_name,vars=[]):
    """
    macro_name: e.g, 'USBankCashPlus'
    vars is a list of strings; in firefox, refer to ${!cmd_var1}, ${!cmd_var2} etc
    """
    # https://github.com/A9T9/RPA/tree/master/command-line/python
    def _PlayMacro(macro, vars=[], path_autorun_html = None):
        assert os.path.exists(path_autorun_html)
        i=1; varstr=''
        for v in vars:
            varstr+=f'&cmd_var{i}={v}'
            i+=1
        args = r'file://' + path_autorun_html + '?macro=' + macro + varstr + '&closeRPA=1&direct=1&continueInLastUsedTab=0&nodisplay=0'
        cmd = f'osascript -e \'tell application "Firefox" to open location "{args}"\''
        print(cmd)
        proc = execute0(cmd)
    _PlayMacro(macro_name, vars=vars, path_autorun_html = fullpath(r'~/Dropbox/Apps/KeyboardMaestro/ui.vision.html'))

def SetClip(content):
    """
    content: The text to be copied to the clipboard.
    """
    import pyperclip
    pyperclip.copy(content)
setclip=SetClip    

def GetClip():
    import pyperclip
    return pyperclip.paste()
getclip=GetClip

def Mail(to, subject, body=None, text=None, attachments=None, bcc=None, cc=None, email=None, password=None):
    """Mail(to, subject, body, attachments=None, bcc=None, cc=None, email=None, password=None)
    to/bcc/cc: ['a@a.com','b@b.com'] or 'a@a.com, b@b.com'
    body: html code or text
    text: specified for weird text that is incompatible with html; when set, body is ignored
    attachments: 'file_in_working_dir.txt' or ['a.txt','b.py','c.pdf']
    email, password: ignored if pysecrets exists
    """
    from yagmail import SMTP
    import html  # html is a python built-in lib (at least included in anaconda)?
    # https://yagmail.readthedocs.io/en/latest/usage.html
    try:
        # import os, sys
        # HERE = os.path.dirname(os.path.abspath(__file__))
        # sys.path.insert(0, HERE)
        
        # EMAIL = "someone@gmail.com", PASSWORD = "abcdefghijkl"
        gclient = SMTP(EMAIL,PASSWORD)
    except:
        gclient = SMTP(email,password)
    contents=None
    if text is not None: body=None
    if text is not None: contents=html.escape(text)
    if body is not None: contents=body
    gclient.send(subject=subject,to=to,cc=cc,bcc=bcc,contents=contents,attachments=attachments)
    gclient.close()
    return None
mail = Mail
gmail = Mail

def o365auth(id=None, secret=None):
    try:
        credentials = (O365ID, O365SECRET)
    except:
        credentials = (id, secret)
    # https://stackoverflow.com/questions/8469122 
    # readline lifts Maximum characters that can be stuffed into raw_input() in Python
    import readline 
    from O365 import Account, FileSystemTokenBackend
    # save token to installation directory
    token_backend = FileSystemTokenBackend(token_path=os.path.dirname(os.path.abspath(__file__)), token_filename='o365_token.txt')
    account = Account(credentials, token_backend=token_backend)
    scopes=['basic', 'message_all','onedrive_all','sharepoint_dl','tasks_all','calendar_all','calendar_shared_all']
    if not account.is_authenticated:  # will check if there is a token and has not expired
        # ask for a login
        # console based authentication See Authentication for other flows
        account.authenticate(scopes=scopes)
    return account

def outlook(to, subject, body=None, attachments=None, bcc=None, cc=None, reply_to=None, id=None, secret=None):
    """outlook(to, subject, body, attachments=None, bcc=None, cc=None, reply_to=None, id=None, secret=None)
    to/bcc/cc: ['a@a.com','b@b.com'] or 'a@a.com, b@b.com'
    reply_to: 'a@a.com'
    body: html code or text. Best practice: three quotes f string (\\n will be auto replaced with <br>) or write/format in Outlook/Word, then Paste as text, Wrap text with <pre></pre>
    attachments: 'file_in_working_dir.txt' or ['a.txt','b.py','c.pdf']
    id, secret: ignored if pysecrets exists
    """
    account = o365auth(id,secret)
    m = account.new_message()
    m.to.add(to)
    m.subject = subject
    m.body = body.replace('\n','<br>')
    if cc is not None: m.cc.add(cc)
    if bcc is not None: m.bcc.add(bcc)
    if reply_to is not None: m.reply_to.add(reply_to)
    if attachments is not None: m.attachments.add(attachments)
    m.send()

def onedrive_connect(file_or_folder_path,id=None,secret=None):
    """
    retrieve a file/folder object for further operation from a path
    retrieve as is, if already a file/folder object
    
    file_or_folder_path: str, '/University/Teaching/Class', '/University/Teaching/Class/Quiz/QuizCompletionAutomation.xlsx'
                         or already a file/folder object
        
    Returns a file/folder item, which can be further processed
        item.copy(target=None, name=None)
        item.move(target)
        item.delete() # to recycle bin
        item.name
    """
    if type(file_or_folder_path) in [str]:
        account = o365auth(id,secret)
        storage = account.storage()  # here we get the storage instance that handles all the storage options.
        # get the default drive
        drive = storage.get_default_drive()
        item = drive.get_item_by_path(file_or_folder_path)
    else:
        item = file_or_folder_path
    return item

def onedrive_move(source,targetFolder,id=None,secret=None):
    """
    remotely move an item (supports file & folder) to another Folder
    source: a folder or file path/object
    targetFolder: a remote folder (fullpath)
    """
    source = onedrive_connect(source,id=id,secret=secret)
    target = onedrive_connect(targetFolder,id=id,secret=secret)
    return source.move(target)

def onedrive_copy(source,targetFolder,id=None,secret=None):
    """
    remotely copy an item (supports file & folder) to another Folder
    source: a folder or file path/object
    targetFolder: a remote folder (fullpath)
    """
    source = onedrive_connect(source,id=id,secret=secret)
    target = onedrive_connect(targetFolder,id=id,secret=secret)
    return source.copy(target)

def onedrive_delete(remote,id=None,secret=None):
    """
    remotely delete an item (supports file & folder) to Recycle Bin
    source: a folder or file path/object
    """
    remote = onedrive_connect(remote,id=id,secret=secret)
    return remote.delete()

def onedrive_rename(remote,newname,id=None,secret=None):
    """
    remotely rename an item (supports file & folder) to the same folder
    remote: a folder or file path/object
    newname: a str (not a full path)
    internally copy to a new name and then delete the original item.
    """
    remote = onedrive_connect(remote,id=id,secret=secret)
    remote.copy(name=newname)
    return remote.delete()

def onedrive_download(remote,local,id=None,secret=None):
    """
    download from remote to local (supports file & folder)
    remote: a folder or file path/object
    local: when remote is folder, a local folder path
           when remote is file, a local file path or folder
    """
    item = onedrive_connect(remote,id=id,secret=secret)
    if item.is_folder: item.download_contents(local)
    if item.is_file: 
        [to_path, name, ext] = splitpath(local)
        if ext=="":
            to_path = local
            name = item.name
        else:
            to_path = to_path
            name = name+ext
        item.download(to_path=to_path, name=name, chunk_size='auto', convert_to_pdf=False, output=None)
    return local

def onedrive_upload(local,remote,id=None,secret=None):
    """
    upload from local to remote (supports file only, overwrite if conflict)
    local: a local file path
    remote: a remote file path
    """
    [remote_folder, name, ext]=splitpath(remote)
    folder = onedrive_connect(remote_folder,id=id,secret=secret)
    folder.upload_file(item=local, item_name=name+ext, chunk_size=5242880, upload_in_chunks=False, stream=None, stream_size=None, conflict_handling=None)
    return remote

def onedrive_share(path,share_type='view',share_scope='anonymous',share_password=None,share_expiration_date=None,id=None,secret=None):
    """
    path: '/University/Teaching/Class'
    share_type: 'view','edit','embed'
    share_scope: 'anonymous','organization'
    share_password: str or None
    share_expiration_date: '2022-02-14', 
        if not specified, auto max determined by adminstrator
        if not specified, call the function a second time will auto renew expiration date (as long as share does not change)
    """
    item = onedrive_connect(path,id=id,secret=secret)
    permission = item.share_with_link(share_type=share_type,share_scope=share_scope,share_password=share_password,share_expiration_date=share_expiration_date)
    return permission.share_link

def onedrive_lsfile(folder_path,limit=None,query=None,order_by=None,batch=None,id=None,secret=None):
    """
    folder_path: '/University/Teaching/Class'
    returns a list of file objects; not recursive
    """
    folder = onedrive_connect(folder_path,id=id,secret=secret)
    res = []
    for item in folder.get_items(limit=limit,query=query,order_by=order_by,batch=batch):
        if item.is_folder:continue
        res.append(item)
    return res

def onedrive_lsfolder(folder_path,limit=None,query=None,order_by=None,batch=None,id=None,secret=None):
    """
    folder_path: '/University/Teaching/Class'
    Returns a list of folder objects; not recursive
    """
    folder = onedrive_connect(folder_path,id=id,secret=secret)
    res = []
    for item in folder.get_items(limit=limit,query=query,order_by=order_by,batch=batch):
        if item.is_file:continue
        res.append(item)
    return res

def onedrive_excel(xlsx,sheet=1,id=None,secret=None):
    """
    xlsx: file path ('/University/Teaching/Class/Quiz/QuizCompletionAutomation.xlsx') or file object
    sheet: sheet number (1 based) or name ('Sheet1')
    
    Returns a worksheet object, which can be further processed
            cell = ws.get_cell(row,col) # row,col are 0 based!
            cell = ws.get_range('A1')   # alternatively
            cell.values = 1
            cell.update()
            
            cell.clear()
            
            table = ws.get_table('Quiz')
            for row in table.get_rows():
                print(row.values)
            # see more: https://github.com/O365/python-o365#excel
    """
    xlsx = onedrive_connect(xlsx,id=None,secret=secret)
    
    from O365.excel import WorkBook
    excel_file = WorkBook(xlsx)
    if type(sheet) in [int]:
        sheet=excel_file.get_worksheets()[sheet-1].name
    ws = excel_file.get_worksheet(sheet)
    return ws
    
def onedrive_readx(xlsx,sheet=1,id=None,secret=None):
    """
    xlsx: file path ('/University/Teaching/Class/Quiz/QuizCompletionAutomation.xlsx') or file object
    sheet: sheet number (1 based) or name ('Sheet1')
    
    Returns a list of list representing rows (that have been used)
    """
    ws = onedrive_excel(xlsx=xlsx,sheet=sheet,id=id,secret=secret)
    return ws.get_used_range().values
    
def getpasswordbw(item,what='usrpwd',sync=False,verbose=0,debug=False):
    """
    item: search string (not case sensitive, better be unique), or item id
    https://bitwarden.com/help/article/cli/#get
    get one at a time: item|username|password|uri|totp|exposed|attachment|folder|collection|organization|org-collection|template|fingerprint
    special customized what='usrpwd'

    in case not working, try native bitwarden commands:
    bw status
    bw login
    bw unlock
    bw sync
    bw get 
    """
    # todo: implement bw on linux
    oldwhat=what
    if what=='usrpwd': what='item'

    # local var could be the same name as global var
    # but if change var, better use local var
    # if change global var, define: global var
    try:
        _EMAIL = EMAIL
        _PASSWORD = PASSWORD + '+'
    except:
        _EMAIL = ''; _PASSWORD=''

    machine = getos()
    debug_msg = ''
    if machine=='Darwin':
        bw = '/usr/local/bin/bw'
        sync = '' if sync else '# ' # comment out
        out = execute0(f'{bw} status',verbose=verbose)
        status = re.search('"status":"(\w+)"',out[0]).group(1)
        if status == 'unauthenticated':
            cmd = f"""
            export BW_USER={_EMAIL}
            export BW_PASSWORD={_PASSWORD}
            {bw} login $BW_USER $BW_PASSWORD
            """
            debug_msg+=cmd
            execute(cmd,verbose=verbose)
            status = 'locked' # login first
        if status == 'locked' or status == 'unlocked':  # always unlock to get session id
            cmd = f"""
            export BW_PASSWORD={_PASSWORD}
            export BW_SESSION=$({bw} unlock --passwordenv BW_PASSWORD --raw)
            {sync}{bw} sync --quiet
            {bw} get {what} {item}
            {bw} lock --quiet
            """
            debug_msg+=cmd
            out = execute0(cmd,verbose=verbose)
            if what=='item':
                import json
                out = json.loads(out[-1])
            else:
                out = out[-1]
    elif machine=='Windows':
        bw = '%WINAPPS%/bw.exe'
        sync = '' if sync else 'rem ' # comment out
        out = execute0(f'{bw} status',verbose=verbose)
        status = re.search('"status":"(\w+)"',out[0]).group(1)
        if status == 'unauthenticated':
            cmd = f"""
            set BW_USER={_EMAIL}
            set BW_PASSWORD={_PASSWORD}
            {bw} login %BW_USER% %BW_PASSWORD%
            """
            debug_msg+=cmd
            execute(cmd,verbose=verbose)
            status = 'locked' # login first
        if status == 'locked' or status == 'unlocked':  # always unlock to get session id
            cmd = f"""
            set BW_PASSWORD={_PASSWORD}
            FOR /F "tokens=*" %%g IN ('{bw} unlock --passwordenv BW_PASSWORD --raw') do (set BW_SESSION=%%g)
            {sync}{bw} sync --quiet
            {bw} get {what} {item}
            {bw} lock --quiet
            """
            debug_msg+=cmd
            out = execute0(cmd,verbose=verbose)
            if what=='item':
                import json
                out = json.loads(out[-1])
            else:
                out = out[-1]
    
    if debug: print(debug_msg); print(out)
    if oldwhat=='usrpwd':
        (usr,pwd)=(out['login']['username'],out['login']['password'])
        return (usr,pwd)
    else:
        return out

def send(*keys,delay=[0.025,0.25],times=1):
    """
    KEYS.COMMAND,KEYS.SHIFT,KEYS.ARROW_LEFT
    KEYS.COMMAND,'a'
    'abc', KEYS.ARROW_DOWN
    'abc\ndef\tg'

    delay in seconds after each key press, 0, [0.025,0.25] <-random between
    times
    delay and times have to be named paramater, cannot be omitted as position parameter because of *keys
    """
    from pynput.keyboard import Key, Controller
    keyboard = Controller()

    # https://www.selenium.dev/selenium/docs/api/py/_modules/selenium/webdriver/common/keys.html#Keys

    # cmd,shift,ctrl,tab,space,backspace,delete,esc,enter,
    # up/down/left/right,home,end,page_up,page_down,
    # insert,menu,pause,print_screen,scroll_lock,
    # f1,f2,...,f20
    # media_next,media_play_pause,media_previous,media_volume_down,media_volume_mute,media_volume_up
    # https://pynput.readthedocs.io/en/latest/keyboard.html#pynput.keyboard.Key

    # remap certain keys between selenium and pynput
    remaps={
        u'\ue003':Key.backspace,   # BACKSPACE
        u'\ue004':Key.tab,   # TAB
        u'\ue006':Key.enter,   # RETURN
        u'\ue007':Key.enter,   # ENTER
        u'\ue008':Key.shift,   # SHIFT, LEFT_SHIFT
        u'\ue009':Key.ctrl,   # CONTROL, LEFT_CONTROL
        u'\ue00a':Key.alt,   # ALT, LEFT_ALT
        u'\ue00c':Key.esc,   # ESCAPE
        u'\ue00d':Key.space,   # SPACE
        u'\ue00e':Key.page_up,   # PAGE_UP
        u'\ue00f':Key.page_down,   # PAGE_DOWN
        u'\ue010':Key.end,   # END
        u'\ue011':Key.home,   # HOME
        u'\ue012':Key.left,   # LEFT, ARROW_LEFT
        u'\ue013':Key.up,   # UP, ARROW_UP
        u'\ue014':Key.right,   # RIGHT, ARROW_RIGHT
        u'\ue015':Key.down,   # DOWN, ARROW_DOWN
        u'\ue017':Key.delete,   # DELETE

        u'\ue031':Key.f1,   # F1
        u'\ue032':Key.f2,   # F2
        u'\ue033':Key.f3,   # F3
        u'\ue034':Key.f4,   # F4
        u'\ue035':Key.f5,   # F5
        u'\ue036':Key.f6,   # F6
        u'\ue037':Key.f7,   # F7
        u'\ue038':Key.f8,   # F8
        u'\ue039':Key.f9,   # F9
        u'\ue03a':Key.f10,   # F10
        u'\ue03b':Key.f11,   # F11
        u'\ue03c':Key.f12,   # F12

        u'\ue03d':Key.cmd,   # META
        u'\ue03d':Key.cmd,   # COMMAND
    }

    if type(delay) not in [tuple,list]: delay=[delay,delay]

    for t in range(0,times):
        if (len(keys)==1) and (keys[0] not in remaps):
            keyboard.type(*keys)
        else:
            for k in keys:
                if k in remaps:
                    keyboard.press(remaps[k])
                else:
                    keyboard.press(k)
            for k in list(reversed(keys)):
                if k in remaps:
                    keyboard.release(remaps[k])
                else:
                    keyboard.release(k)
                sleep(random.uniform(delay[0],delay[1]))

# def ocr(img=None,lang=['en'],gpu=False,*args,**kwargs):
#     # cannot work with certain opencv version, use instead: https://github.com/schappim/macOCR
#     """
#     img: filepath, OpenCV image object (numpy array), image file as bytes, URL to raw image (default None=clipboard img)
#     lang: ['en','ch_sim','ch_tra']
#     https://www.jaided.ai/easyocr/documentation/
#     """
#     from PIL import ImageGrab
#     import easyocr
#     if img is None: img = ImageGrab.grabclipboard()
#     reader = easyocr.Reader(lang,gpu=gpu)
#     result = reader.readtext(img,detail=True,paragraph=False,*args,**kwargs)
#     # result=[([[189, 75], [469, 75], [469, 165], [189, 165]], '愚园路', 0.3754989504814148),
#     #  ([[86, 80], [134, 80], [134, 128], [86, 128]], '西', 0.40452659130096436),
#     #  ([[517, 81], [565, 81], [565, 123], [517, 123]], '东', 0.9989598989486694),
#     #  ([[78, 126], [136, 126], [136, 156], [78, 156]], '315', 0.8125889301300049),
#     #  ([[514, 126], [574, 126], [574, 156], [514, 156]], '309', 0.4971577227115631),
#     #  ([[226, 170], [414, 170], [414, 220], [226, 220]], 'Yuyuan Rd.', 0.8261902332305908),
#     #  ([[79, 173], [125, 173], [125, 213], [79, 213]], 'W', 0.9848111271858215),
#     #  ([[529, 173], [569, 173], [569, 213], [529, 213]], 'E', 0.8405593633651733)]
#     top=[]; left=[]; out={}
#     for r in result:
#         top=r[0][0][1]; left=r[0][0][0]; text=r[1]
#         if top not in out.keys():
#             # /10 to scale
#             out[top]=' '*int(left/10)+text
#         else:
#             text_old = out[top].lstrip()
#             text_new = text
#             left_old = len(out[top]) - len(out[top].lstrip())
#             left_new = int(left/10)
#             if left_new > left_old:
#                 part1 = ' '*int(left_old/10)+text_old; part2 = ' '*int(left_new - len(part1))+text_new
#             else:
#                 # unlike to trigger, since result ordered by top already
#                 part1 = ' '*int(left_new/10)+text_new; part2 = ' '*int(left_old - len(part1))+text_old
#             out[top] = part1+part2
#     result=[text for (top, text) in sorted(out.items())]
#     return result
#     # for r in result:
#     #     print(r)

# def findimg(image, area=None, p=0.8, show=False):
#     '''
#     Searchs for an image on the screen
    
#     image: path to the image file (see opencv imread for supported types)
#     area: [x,y,w,h]
#     p: precision the higher, the lesser tolerant and fewer false positives are found default is 0.8
#     show: show all matches (highlighted)

#     returns :
#     [[x,y,w,h]] or [] if not found
#     '''
#     # modified from https://github.com/drov0/python-imagesearch/blob/master/python_imagesearch/imagesearch.py
#     import cv2
#     import numpy as np
#     import mss
#     from PIL import Image

#     # https://python-mss.readthedocs.io/examples.html
#     sct = mss.mss()
#     if area is None: 
#         area=sct.monitors[0]
#     else:
#         area={'left': area[0], 'top': area[1], 'width': area[2], 'height': area[3]}
#     im = sct.grab(area)
#     # raw is bgra
#     # Image.frombytes("RGB", im.size, im.bgra, "raw", "BGRX").show()
#     img = np.array(im)
#     img = cv2.cvtColor(img, cv2.COLOR_BGRA2BGR)
#     # cv2.imshow("OpenCV", img)
#     # Image.fromarray(cv2.cvtColor(img, cv2.COLOR_BGR2RGB)).show()
#     template = cv2.imread(image, cv2.IMREAD_COLOR) 
#     # cv2.imshow("OpenCV", template)
#     # Image.fromarray(cv2.cvtColor(template, cv2.COLOR_BGR2RGB)).show()
#     if template is None:
#         raise FileNotFoundError('Image file not found: {}'.format(image))
#     res = cv2.matchTemplate(img, template, cv2.TM_CCOEFF_NORMED)
    
#     matches = []
#     # https://docs.opencv.org/3.4/d4/dc6/tutorial_py_template_matching.html
#     h, w, channels = template.shape
#     loc = np.where( res >= p )
#     for pt in zip(*loc[::-1]):
#         cv2.rectangle(img, pt, (pt[0] + w, pt[1] + h), (0,0,255), 2)
#         matches.append([ pt[0],pt[1],w,h,res[pt[1],pt[0]] ])
#     if show:
#         # cv2.imshow("OpenCV", img)
#         # cv2.imwrite('res.png',img)
#         # https://stackoverflow.com/a/43234001/2292993
#         Image.fromarray(cv2.cvtColor(img, cv2.COLOR_BGR2RGB)).show()
#     # min_val, max_val, min_loc, max_loc = cv2.minMaxLoc(res)
#     # if max_val < p:
#     #     return None
#     # return max_loc
#     return matches

def _mouse(init_pos, fin_pos, deviation=25, delay=2):
    from random import randint, choice
    from math import ceil, sqrt
    #time parameter
    speed=delay
    ts = [t/(speed * 100.0) for t in range(speed * 101)]
    #bezier centre control points between (int(deviation / 2)) and (deviaion) of travel distance, plus or minus at random
    control_1 = (init_pos[0] + choice((-1, 1)) * abs(ceil(fin_pos[0]) - ceil(init_pos[0])) * 0.01 * randint(int(deviation / 2), deviation),
                init_pos[1] + choice((-1, 1)) * abs(ceil(fin_pos[1]) - ceil(init_pos[1])) * 0.01 * randint(int(deviation / 2), deviation)
                    )
    control_2 = (init_pos[0] + choice((-1, 1)) * abs(ceil(fin_pos[0]) - ceil(init_pos[0])) * 0.01 * randint(int(deviation / 2), deviation),
                init_pos[1] + choice((-1, 1)) * abs(ceil(fin_pos[1]) - ceil(init_pos[1])) * 0.01 * randint(int(deviation / 2), deviation)
                    )
    xys = [init_pos, control_1, control_2, fin_pos]
    def make_bezier(xys):
        # xys should be a sequence of 2-tuples (Bezier control points)
        def pascal_row(n):
            # This returns the nth row of Pascal's Triangle
            result = [1]
            x, numerator = 1, n
            for denominator in range(1, n//2+1):
                # print(numerator,denominator,x)
                x *= numerator
                x /= denominator
                result.append(x)
                numerator -= 1
            if n&1 == 0:
                # n is even
                result.extend(reversed(result[:-1]))
            else:
                result.extend(reversed(result)) 
            return result
        n = len(xys)
        combinations = pascal_row(n - 1)
        def bezier(ts):
            # This uses the generalized formula for bezier curves
            # http://en.wikipedia.org/wiki/B%C3%A9zier_curve#Generalization
            result = []
            for t in ts:
                tpowers = (t**i for i in range(n))
                upowers = reversed([(1-t)**i for i in range(n)])
                coefs = [c*a*b for c, a, b in zip(combinations, tpowers, upowers)]
                result.append(
                    list(sum([coef*p for coef, p in zip(coefs, ps)]) for ps in zip(*xys)))
            return result
        return bezier
    bezier = make_bezier(xys)
    mouse_points = bezier(ts)
    #round floats to ints
    mouse_points = [[round(v) for v in x] if type(x) is not str else x for x in mouse_points]
    return mouse_points

def _move(b, a=None, deviation=25, delay=2, rate=1000, draw=False):
    '''
    GENERATE BEZIER CURVE POINTS
    Takes a (init_pos, default current position) --> b (fin_pos) as a 2-tuple representing xy coordinates
    deviation (int)
        deviation controls how straight the lines drawn my the cursor
        are. Zero deviation gives straight lines
        Accuracy is a percentage of the displacement of the mouse from point A to
        B, which is given as maximum control point deviation.
        Naturally, deviation of 10 (10%) gives maximum control point deviation
        of 10% of magnitude of displacement of mouse from point A to B, 
        and a minimum of 5% (int(deviation / 2))    
    delay: an int multiplier for speed. The lower, the faster. 1 is fastest.
    rate: another param related to speed, total travel time in seconds = distance/rate. The higher, the faster.
    draw: boolean deciding whether or not to show the curve the mouse makes
    '''
    # https://github.com/vincentbavitz/bezmouse
    # see also https://pyautogui.readthedocs.io/en/latest/mouse.html#tween-easing-functions
    from time import sleep
    from PIL import Image
    from random import randint, choice
    from math import ceil, sqrt
    from pynput.mouse import Button, Controller
    mouse = Controller()
    if a is None: a = mouse.position
    init_pos=a; fin_pos=b
    mouse_points=_mouse(init_pos, fin_pos, deviation, delay)
    '''
    Moves mouse in accordance with a list of points (continuous curve)
    mouse_points
        list of 2-tuples or lists of ints or floats representing xy coords
    draw
        a boolean deciding whether or not to show the curve the mouse makes
    '''
    if draw == True:
        REL_ORIGIN=[0,0]
        drawpoints = [(v[0] - REL_ORIGIN[0], v[1] - REL_ORIGIN[1]) for v in mouse_points if type(v) is not str]
        from PIL import ImageGrab
        img = ImageGrab.grab()
        width_px, height_px = img.size  # primary screen size
        def draw_points(points, width=width_px, height=height_px):
            '''
            Draws yellow crosses to a image for all coordinates in "points"
            show the image
            '''
            img = Image.new("RGB", (width, height))
            pix = img.load()
            try:
                for coords in points:
                    pix[coords[0], coords[1]] = (255, 255, 0)
                    pix[coords[0] + 1, coords[1] + 1] = (255, 255, 0)
                    pix[coords[0] + 1, coords[1] - 1] = (255, 255, 0)
                    pix[coords[0] - 1, coords[1] + 1] = (255, 255, 0)
                    pix[coords[0] - 1, coords[1] - 1] = (255, 255, 0)
                img.show()
            except:
                pass
        draw_points(drawpoints)    
    distance=sqrt((b[0]-a[0])**2+(b[1]-a[1])**2)
    duration=distance/rate
    for coord in mouse_points:
        sleep(duration/len(mouse_points))
        mouse.position=coord

def moveclick(location, radius=5, n=2, wait=0.25, deviation=25, delay=2, rate=1000, draw=False, *args, **kwargs):
    '''
    location: (x,y) or [x,y], [x,y,w,h,p] or [x,y,w,h]
    radius: works only if location is an area. if an area is 100*100 with an radius of 5 it may be at 46,50 the first time and then 55,53 etc
    wait: seconds before click after move to the location
    deviation (int)
        deviation controls how straight the lines drawn my the cursor
        are. Zero deviation gives straight lines
        Accuracy is a percentage of the displacement of the mouse from point A to
        B, which is given as maximum control point deviation.
        Naturally, deviation of 10 (10%) gives maximum control point deviation
        of 10% of magnitude of displacement of mouse from point A to B, 
        and a minimum of 5% (int(deviation / 2))    
    delay: an int multiplier for speed. The lower, the faster. 1 is fastest.
    rate: another param related to speed, total travel time in seconds = distance/rate. The higher, the faster.
    draw: boolean deciding whether or not to show the curve the mouse makes
    '''
    from random import randint
    if len(location)==2:
        xx,yy=location[0],location[1]
    else:
        x,y,w,h=location[0],location[1],location[2],location[3]
        xx = x + randint(int(w/2-radius), int(w/2+radius))
        yy = y + randint(int(h/2-radius), int(h/2+radius))
    _move((xx, yy), a=None, deviation=deviation, delay=delay, rate=rate, draw=draw, *args, **kwargs)
    import time
    time.sleep(wait)
    from pynput.mouse import Button, Controller
    mouse = Controller()
    mouse.click(Button.left, n)

def click(n=2):
    from pynput.mouse import Button, Controller
    mouse = Controller()
    mouse.click(Button.left, n)

def rclick(n=1):
    from pynput.mouse import Button, Controller
    mouse = Controller()
    mouse.click(Button.right, n)

def scroll(dx,dy):
    from pynput.mouse import Button, Controller
    mouse = Controller()
    mouse.scroll(dx, dy)

def position():
    from pynput.mouse import Button, Controller
    mouse = Controller()
    return mouse.position

def size():
    # https://stackoverflow.com/questions/3129322/
    mm_per_inch = 25.4
    px_per_inch =  72.0 #most common
    import mss
    sct=mss.mss()
    return sct.monitors

####************************************************************************************************
                                     ####*GSheet*####
####************************************************************************************************
# the first place to check authorization is go to the following gspread link:
# https://gspread.readthedocs.io/en/latest/oauth2.html#enable-api-access
# must grant permission to service account (pysheet@python-243720.iam.gserviceaccount.com) 
# as editor of each individual google sheet mannually
class GSheet():
    """
    .ws (worksheet), .wb (workbook) return the orginal google sheet/sheets objects
    """
    def __init__(self,url,sheet_name=None):
        try:
            import gspread
            # not sure where I got these json file (maybe downloaded from google?)
            client = gspread.service_account_from_dict(GSHEET_KEY)
        except:
            pass

        sheets = client.open_by_url(url)
        if sheet_name is None:
            sheet = sheets.sheet1
        else:
            sheet = sheets.worksheet(sheet_name)
        self.wb = sheets
        self.ws = sheet

    @property
    def title(self):
        return self.ws.title
    
    @property
    def url(self):
        return self.ws.url

    @property
    def Ncol(self):
        """
        physical n col, may be blank
        """
        return self.ws.col_count

    @property
    def Nrow(self):
        return self.ws.row_count

    @property
    def columns(self):
        try:
            return self.ws.row_values(1)
        except:
            return []

    @property
    def values(self):
        """as list of list"""
        return self.ws.get_all_values()
    
    @property
    def records(self):
        """as dict"""
        return self.ws.get_all_records()    

    @property
    def nrow(self):
        values = self.values
        return len(values)

    @property
    def ncol(self):
        values = self.values
        tmp = [len(r) for r in values]
        if len(tmp)==0:
            return 0
        else:
            import numpy as np
            return np.max(tmp)
    
    def clear(self,range):
        # https://gspread.readthedocs.io/en/latest/api.html#gspread.models.Spreadsheet.values_clear
        # sheets method
        range = f"'{self.title}'!{range}"
        return self.wb.values_clear(range)

    def getrc(self,irow,icol,option='FORMATTED_VALUE'):
        """
        irow,icol: int row/col number
        option: 'FORMATTED_VALUE', 'UNFORMATTED_VALUE', 'FORMULA'
        returns a single cell value (trimmed if string)
        """
        val = self.ws.cell(irow,icol,value_render_option=option).value
        if type(val) in [str]: val = val.strip()
        return val

    def getval(self,*args,option='FORMATTED_VALUE',**kwargs):
        """
        range could be A:A, 2:2
        option: 'FORMATTED_VALUE', 'UNFORMATTED_VALUE', 'FORMULA'
        returns a list of list
        """
        # https://gspread.readthedocs.io/en/latest/api.html#gspread.models.Worksheet.get
        return list(self.ws.get(*args, value_render_option=option, **kwargs))
    
    def getc(self,col,*args,option='FORMATTED_VALUE',**kwargs):
        """
        col int, 1 based
        option: 'FORMATTED_VALUE', 'UNFORMATTED_VALUE', 'FORMULA'
        returns a list
        """
        # https://gspread.readthedocs.io/en/latest/api.html#gspread.models.Worksheet.get
        return self.ws.col_values(col, *args, value_render_option=option, **kwargs)

    def getr(self,row,*args,option='FORMATTED_VALUE',**kwargs):
        """
        row int, 1 based
        option: 'FORMATTED_VALUE', 'UNFORMATTED_VALUE', 'FORMULA'
        returns a list
        """
        # https://gspread.readthedocs.io/en/latest/api.html#gspread.models.Worksheet.get
        return self.ws.row_values(row, *args, value_render_option=option, **kwargs)

    def insert(self,values,row=1,option="USER_ENTERED"):
        """
        values (list) – List of list. eg, [[11,12],[21,22]]
        index (int) – Start row to update (one-based). Defaults to 1 (one).
        option: "RAW", "USER_ENTERED"
        """
        # https://gspread.readthedocs.io/en/latest/api.html#gspread.models.Worksheet.insert_rows
        return self.ws.insert_rows(values=values, row=row, value_input_option=option)

    def update(self,range,values,option="USER_ENTERED"):
        """
        option: "RAW", "USER_ENTERED"
        Note that update range can be bigger than values array: update('A2:B4', [[42], [43]]) -> Updates A2 and A3 with values 42 and 43
        """
        # https://gspread.readthedocs.io/en/latest/api.html#gspread.models.Worksheet.update
        return self.ws.update(range,values,value_input_option=option)

    def append(self,values,option="USER_ENTERED"):
        """
        option: "RAW", "USER_ENTERED"
        append to the end of data
        """
        # https://gspread.readthedocs.io/en/latest/api.html#gspread.models.Worksheet.append_rows
        return self.ws.update(f'A{self.nrow+1}',values,value_input_option=option)

    def switch(self,sheet_name):
        """switch to a particular worksheet if there are multiple ones"""
        self.ws = self.wb.worksheet(sheet_name)
        return None

    def export(self,xlsx_filepath):
        """export all worksheets to a xlsx (formula, color, etc will be discarded; essentially read into list of list and then write xlsx)"""
        import xlsxwriter
        xbook = xlsxwriter.Workbook(xlsx_filepath)
        for ws in self.wb.worksheets():
            data = ws.get_all_values()
            xsheet = xbook.add_worksheet(ws.title)
            for r, row in enumerate(data,0):
                # https://xlsxwriter.readthedocs.io/worksheet.html#worksheet-write-row
                xsheet.write_row(r, 0, row)
        xbook.close()
        return None

# sheet = client.open("Master-Letter for Jerry's Teaching Application").sheet1
# sheets = client.open_by_url('https://docs.google.com/spreadsheets/d/1VexrXFPEdh5oAFP9CuhpsRCVqiHJhbAd_p-FL3kiu0Q/edit')
# sheets = client.open_by_url(gsheeturl)
# sheet = sheets.worksheet("Transactions")

# API:
# https://gspread.readthedocs.io/en/latest/api.html#models
# a row does not have to include all columns. Just like when you type in GSheet. 
# The rest of columns that are not included will not be changed by append/insert.
# USER_ENTERED
# https://stackoverflow.com/questions/27125967/google-spreadsheets-gspread-append-row-issue
# https://gspread.readthedocs.io/en/latest/api.html?highlight=append#gspread.models.Worksheet.insert_row
# https://github.com/burnash/gspread/issues/524
# sheet.append_row(line,'USER_ENTERED')
# sheet.update_acell('G2', ''); sheet.update_acell('H2', ''); sheet.update_acell('I2', '')
# sheet.insert_row(line,index=2,value_input_option='USER_ENTERED')

####************************************************************************************************
                                     ####*Dropbox*####
####************************************************************************************************
import dropbox
# https://dropbox-sdk-python.readthedocs.io/en/latest/api/dropbox.html#

def dropbox_spaceusage():
    # returns in GB [left, used, total (individual)]
    dbx = dropbox.Dropbox(DROPBOX_ACCESS_TOKEN)
    usage = dbx.users_get_space_usage()
    usage = [usage.used/1024/1024/1024, usage.allocation.get_individual().allocated/1024/1024/1024]
    return [round(usage[1]-usage[0],2), round(usage[0],2), round(usage[1],2)]

def upload(localfile,cloudfile=None):
    # shortcut if only one cloudpath passed in, assuming current local working directory
    # e.g., '/Investment/fz/fz/fz/data_td.json'
    if cloudfile is None:
        path,file= os.path.split(localfile)
        cloudfile=localfile
        localfile=file
    dbx = dropbox.Dropbox(DROPBOX_ACCESS_TOKEN)
    with open(localfile, "rb") as f:
        dbx.files_upload(f.read(),cloudfile,dropbox.files.WriteMode.overwrite,mute=True)
        # print(cloudfile+' upload done.')

def download(cloudfile,localfile=None):
    # shortcut if only one cloudpath passed in, assuming current local working directory
    # e.g., '/Investment/fz/fz/fz/data_td.json'
    if localfile is None:
        path,file= os.path.split(cloudfile)
        localfile=file
    dbx = dropbox.Dropbox(DROPBOX_ACCESS_TOKEN)
    dbx.files_download_to_file(localfile,cloudfile)
    return localfile

def _files_list_folder_all(dbx,folder,recursive=False,limit=2000):
    cur = dbx.files_list_folder(folder,recursive=recursive,limit=limit)
    items = cur.entries
    while cur.has_more==True:
        cur = dbx.files_list_folder_continue(cur.cursor)
        items.extend(cur.entries)
    return items

def lsfile(folder,recursive=False):
    dbx = dropbox.Dropbox(DROPBOX_ACCESS_TOKEN)
    items = _files_list_folder_all(dbx,folder,recursive=recursive,limit=2000)
    files = []
    for item in items:
        if type(item)==dropbox.files.FolderMetadata: continue
        # e.g., '/Investment/fz/fz/fz/data_td.json', 'data_td.json'
        files.append([item.path_display,item.name])
    return files

def lsfolder(folder,recursive=False):
    dbx = dropbox.Dropbox(DROPBOX_ACCESS_TOKEN)
    items = _files_list_folder_all(dbx,folder,recursive=recursive,limit=2000)
    files = []
    for item in items:
        if type(item)==dropbox.files.FileMetadata: continue
        # e.g., '/Investment/fz/fz/fz/firefox', 'firefox'
        files.append([item.path_display,item.name])
    return files

def delfilefolder(file):
    # Delete the file or folder at a given path. If the path is a folder, all its contents will be deleted too. A successful response indicates that the file or folder was deleted
    dbx = dropbox.Dropbox(DROPBOX_ACCESS_TOKEN)
    return dbx.files_delete(file)

####************************************************************************************************
                                     ####*GCal*####
####************************************************************************************************
# Reference: 
# https://github.com/kuzmoyev/google-calendar-simple-api
# https://google-calendar-simple-api.readthedocs.io/en/latest/getting_started.html
class GCal():
    """
    
    """
    def __init__(self,default_calendar: str = 'primary'):
        """
        default_calendar: 
                Users email address or name/id of the calendar. Default: primary calendar of the user
                If user's email or "primary" is specified, then primary calendar of the user is used.
                You don't need to specify this parameter in this case as it is a default behaviour.
        """
        # work in library path
        oldpwd = os.getcwd()
        os.chdir(os.path.dirname(os.path.abspath(__file__)))

        from gcsa.google_calendar import GoogleCalendar
        savej(GCAL_KEY,'gcal_secrets.json')
        
        calendar = GoogleCalendar(default_calendar,credentials_path='gcal_secrets.json',token_path='gcal_token.pickle',save_token=True)

        os.chdir(oldpwd)
        self.gcal = calendar

    def addevent(self,*args,**kwargs):
        """
        summary: str,
        start: Union[datetime.date, datetime.datetime, beautiful_date.beautiful_date.BeautifulDate],
        end: Union[datetime.date, datetime.datetime, beautiful_date.beautiful_date.BeautifulDate] = None,
        *,
        timezone: str = 'America/Chicago',
        event_id: str = None,
        description: str = None,
        location: str = None,
        recurrence: Union[str, List[str]] = None,
        color_id: str = None,
        visibility: str = 'default',
        attendees: Union[gcsa.attendee.Attendee, str, List[gcsa.attendee.Attendee], List[str]] = None,
        attachments: Union[gcsa.attachment.Attachment, List[gcsa.attachment.Attachment]] = None,
        conference_solution: Union[gcsa.conference.ConferenceSolution, gcsa.conference.ConferenceSolutionCreateRequest] = None,
        reminders: Union[gcsa.reminders.Reminder, List[gcsa.reminders.Reminder]] = None,
        default_reminders: bool = False,
        minutes_before_popup_reminder: int = None,
        minutes_before_email_reminder: int = None,
        guests_can_invite_others: bool = True,
        guests_can_modify: bool = False,
        guests_can_see_other_guests: bool = True,
        transparency: str = None
        """
        from gcsa.event import Event

        event = Event(*args,**kwargs)
        self.gcal.add_event(event)

####************************************************************************************************
                                     ####*GDrive*####
####************************************************************************************************
# Reference: 
# https://docs.iterative.ai/PyDrive2/quickstart/#authentication
# https://qiita-com.translate.goog/ftnext/items/60ced8bc432bec6101f0?_x_tr_sl=auto&_x_tr_tl=en&_x_tr_hl=en&_x_tr_pto=wapp
# https://stackoverflow.com/a/24542604/2292993 (outdated)

# Steps:
# 1) Create and save as gdrive_secrets.json (I saved the contents to pysecrets.py and dump it to physical json file)
# 2) Create settings.yaml
# 3) Authorize via web with gdrive_secrets.json for the first time
# 4) Get saved gdrive_credentials.json which can be used in the future without need to authroize via web
def gauth():
    # work in library path
    oldpwd = os.getcwd()
    os.chdir(os.path.dirname(os.path.abspath(__file__)))

    from pydrive2.auth import GoogleAuth
    from pydrive2.drive import GoogleDrive

    # jerry: dump gdrive_secrets.json
    savej(GDRIVE_KEY,'gdrive_secrets.json')
    settings="""
            client_config_file: gdrive_secrets.json
            save_credentials: True
            save_credentials_backend: file
            save_credentials_file: gdrive_credentials.json
            get_refresh_token: True
            """
    with open('settings.yaml', 'w') as f:
        f.write(settings)

    googleauth = GoogleAuth()
    googleauth.LocalWebserverAuth()
    gdrive = GoogleDrive(googleauth)

    os.chdir(oldpwd)
    return gdrive

def gdownload(id,filename):
    """
    export/download a file
    id: get from link
    the downloaded file may not have the exact same content formatting
    """
    gdrive = gauth()
    # https://developers.google.com/drive/api/guides/ref-export-formats
    gfile = gdrive.CreateFile({'id': id})
    FORMATS = { 
                '.txt': 'text/plain',
                '.rtf': 'application/rtf',
                '.pdf': 'application/pdf',
                '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                '.png': 'image/png',
                '.svg': 'image/svg+xml',
                '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
               }
    gfile.GetContentFile(filename, FORMATS[splitpath(filename)[-1]])
    return filename

def gcopy(id,filetitle,folderid=None):
    """
    make a copy of file in gdrive
    id: get from link
    filetitle: new file title/name
    folderid: copy to a folder other than root [optional], get from link
    """
    # https://developers.google.com/drive/api/v2/reference/files/copy
    # https://stackoverflow.com/questions/43865016/python-copy-a-file-in-google-drive-into-a-specific-folder
    # https://github.com/iterative/PyDrive2/issues/53
    gdrive = gauth()
    # https://github.com/googlearchive/PyDrive/blob/986ad8fb4f68dde6ffaa7c61c951d9573314dea0/pydrive/auth.py#L61
    gdrive.auth.Authorize() # needed, otherwise gdrive.auth.service is None
    if folderid is None:
        gdrive.auth.service.files().copy(fileId=id,
                                         body={'title': filetitle}).execute()
    else:
        gdrive.auth.service.files().copy(fileId=id,
                                         body={"parents": [{"id": folderid}],'title': filetitle}).execute()
    return filetitle

def gls(folderid=None):
    """
    ls files and folders in a folder (not including trashed)
    folderid: optional, if None (default), search root folder
    returns a list of tuples (title,id)
    """
    gdrive = gauth()
    res = []

    # https://stackoverflow.com/questions/40224559/list-of-file-in-a-folder-drive-api-pydrive
    # Build string dynamically (need to use escape characters to support single quote syntax)
    if folderid is None:
        str = "'root' in parents and trashed=false"
    else:
        str = "\'" + folderid + "\'" + " in parents and trashed=false"    

    # Starting iterating over files
    filelist = gdrive.ListFile({'q': str}).GetList()
    for file in filelist:
        # print('title: %s, id: %s' % (file['title'], file['id']))
        res.append((file['title'], file['id']))
    
    return res

####************************************************************************************************
                                     ####*OrderedSet*####
####************************************************************************************************
"""
https://pypi.org/project/ordered-set/
An OrderedSet is a custom MutableSet that remembers its order, so that every
entry has an index that can be looked up.

Based on a recipe originally posted to ActiveState Recipes by Raymond Hettiger,
and released under the MIT license.
"""
import itertools as it
from collections import deque

try:
    # Python 3
    from collections.abc import MutableSet, Sequence
except ImportError:
    # Python 2.7
    from collections import MutableSet, Sequence

SLICE_ALL = slice(None)
__version__ = "3.1"


def is_iterable(obj):
    """
    Are we being asked to look up a list of things, instead of a single thing?
    We check for the `__iter__` attribute so that this can cover types that
    don't have to be known by this module, such as NumPy arrays.

    Strings, however, should be considered as atomic values to look up, not
    iterables. The same goes for tuples, since they are immutable and therefore
    valid entries.

    We don't need to check for the Python 2 `unicode` type, because it doesn't
    have an `__iter__` attribute anyway.
    """
    return (
        hasattr(obj, "__iter__")
        and not isinstance(obj, str)
        and not isinstance(obj, tuple)
    )


class OrderedSet(MutableSet, Sequence):
    """
    An OrderedSet is a custom MutableSet that remembers its order, so that
    every entry has an index that can be looked up.

    Example:
        >>> OrderedSet([1, 1, 2, 3, 2])
        OrderedSet([1, 2, 3])
    """

    def __init__(self, iterable=None):
        self.items = []
        self.map = {}
        if iterable is not None:
            self |= iterable

    def __len__(self):
        """
        Returns the number of unique elements in the ordered set

        Example:
            >>> len(OrderedSet([]))
            0
            >>> len(OrderedSet([1, 2]))
            2
        """
        return len(self.items)

    def __getitem__(self, index):
        """
        Get the item at a given index.

        If `index` is a slice, you will get back that slice of items, as a
        new OrderedSet.

        If `index` is a list or a similar iterable, you'll get a list of
        items corresponding to those indices. This is similar to NumPy's
        "fancy indexing". The result is not an OrderedSet because you may ask
        for duplicate indices, and the number of elements returned should be
        the number of elements asked for.

        Example:
            >>> oset = OrderedSet([1, 2, 3])
            >>> oset[1]
            2
        """
        if isinstance(index, slice) and index == SLICE_ALL:
            return self.copy()
        elif is_iterable(index):
            return [self.items[i] for i in index]
        elif hasattr(index, "__index__") or isinstance(index, slice):
            result = self.items[index]
            if isinstance(result, list):
                return self.__class__(result)
            else:
                return result
        else:
            raise TypeError("Don't know how to index an OrderedSet by %r" % index)

    def copy(self):
        """
        Return a shallow copy of this object.

        Example:
            >>> this = OrderedSet([1, 2, 3])
            >>> other = this.copy()
            >>> this == other
            True
            >>> this is other
            False
        """
        return self.__class__(self)

    def __getstate__(self):
        if len(self) == 0:
            # The state can't be an empty list.
            # We need to return a truthy value, or else __setstate__ won't be run.
            #
            # This could have been done more gracefully by always putting the state
            # in a tuple, but this way is backwards- and forwards- compatible with
            # previous versions of OrderedSet.
            return (None,)
        else:
            return list(self)

    def __setstate__(self, state):
        if state == (None,):
            self.__init__([])
        else:
            self.__init__(state)

    def __contains__(self, key):
        """
        Test if the item is in this ordered set

        Example:
            >>> 1 in OrderedSet([1, 3, 2])
            True
            >>> 5 in OrderedSet([1, 3, 2])
            False
        """
        return key in self.map

    def add(self, key):
        """
        Add `key` as an item to this OrderedSet, then return its index.

        If `key` is already in the OrderedSet, return the index it already
        had.

        Example:
            >>> oset = OrderedSet()
            >>> oset.append(3)
            0
            >>> print(oset)
            OrderedSet([3])
        """
        if key not in self.map:
            self.map[key] = len(self.items)
            self.items.append(key)
        return self.map[key]

    append = add

    def update(self, sequence):
        """
        Update the set with the given iterable sequence, then return the index
        of the last element inserted.

        Example:
            >>> oset = OrderedSet([1, 2, 3])
            >>> oset.update([3, 1, 5, 1, 4])
            4
            >>> print(oset)
            OrderedSet([1, 2, 3, 5, 4])
        """
        item_index = None
        try:
            for item in sequence:
                item_index = self.add(item)
        except TypeError:
            raise ValueError(
                "Argument needs to be an iterable, got %s" % type(sequence)
            )
        return item_index

    def index(self, key):
        """
        Get the index of a given entry, raising an IndexError if it's not
        present.

        `key` can be an iterable of entries that is not a string, in which case
        this returns a list of indices.

        Example:
            >>> oset = OrderedSet([1, 2, 3])
            >>> oset.index(2)
            1
        """
        if is_iterable(key):
            return [self.index(subkey) for subkey in key]
        return self.map[key]

    # Provide some compatibility with pd.Index
    get_loc = index
    get_indexer = index

    def pop(self):
        """
        Remove and return the last element from the set.

        Raises KeyError if the set is empty.

        Example:
            >>> oset = OrderedSet([1, 2, 3])
            >>> oset.pop()
            3
        """
        if not self.items:
            raise KeyError("Set is empty")

        elem = self.items[-1]
        del self.items[-1]
        del self.map[elem]
        return elem

    def discard(self, key):
        """
        Remove an element.  Do not raise an exception if absent.

        The MutableSet mixin uses this to implement the .remove() method, which
        *does* raise an error when asked to remove a non-existent item.

        Example:
            >>> oset = OrderedSet([1, 2, 3])
            >>> oset.discard(2)
            >>> print(oset)
            OrderedSet([1, 3])
            >>> oset.discard(2)
            >>> print(oset)
            OrderedSet([1, 3])
        """
        if key in self:
            i = self.map[key]
            del self.items[i]
            del self.map[key]
            for k, v in self.map.items():
                if v >= i:
                    self.map[k] = v - 1

    def clear(self):
        """
        Remove all items from this OrderedSet.
        """
        del self.items[:]
        self.map.clear()

    def __iter__(self):
        """
        Example:
            >>> list(iter(OrderedSet([1, 2, 3])))
            [1, 2, 3]
        """
        return iter(self.items)

    def __reversed__(self):
        """
        Example:
            >>> list(reversed(OrderedSet([1, 2, 3])))
            [3, 2, 1]
        """
        return reversed(self.items)

    def __repr__(self):
        if not self:
            return "%s()" % (self.__class__.__name__,)
        return "%s(%r)" % (self.__class__.__name__, list(self))

    def __eq__(self, other):
        """
        Returns true if the containers have the same items. If `other` is a
        Sequence, then order is checked, otherwise it is ignored.

        Example:
            >>> oset = OrderedSet([1, 3, 2])
            >>> oset == [1, 3, 2]
            True
            >>> oset == [1, 2, 3]
            False
            >>> oset == [2, 3]
            False
            >>> oset == OrderedSet([3, 2, 1])
            False
        """
        # In Python 2 deque is not a Sequence, so treat it as one for
        # consistent behavior with Python 3.
        if isinstance(other, (Sequence, deque)):
            # Check that this OrderedSet contains the same elements, in the
            # same order, as the other object.
            return list(self) == list(other)
        try:
            other_as_set = set(other)
        except TypeError:
            # If `other` can't be converted into a set, it's not equal.
            return False
        else:
            return set(self) == other_as_set

    def union(self, *sets):
        """
        Combines all unique items.
        Each items order is defined by its first appearance.

        Example:
            >>> oset = OrderedSet.union(OrderedSet([3, 1, 4, 1, 5]), [1, 3], [2, 0])
            >>> print(oset)
            OrderedSet([3, 1, 4, 5, 2, 0])
            >>> oset.union([8, 9])
            OrderedSet([3, 1, 4, 5, 2, 0, 8, 9])
            >>> oset | {10}
            OrderedSet([3, 1, 4, 5, 2, 0, 10])
        """
        cls = self.__class__ if isinstance(self, OrderedSet) else OrderedSet
        containers = map(list, it.chain([self], sets))
        items = it.chain.from_iterable(containers)
        return cls(items)

    def __and__(self, other):
        # the parent implementation of this is backwards
        return self.intersection(other)

    def intersection(self, *sets):
        """
        Returns elements in common between all sets. Order is defined only
        by the first set.

        Example:
            >>> oset = OrderedSet.intersection(OrderedSet([0, 1, 2, 3]), [1, 2, 3])
            >>> print(oset)
            OrderedSet([1, 2, 3])
            >>> oset.intersection([2, 4, 5], [1, 2, 3, 4])
            OrderedSet([2])
            >>> oset.intersection()
            OrderedSet([1, 2, 3])
        """
        cls = self.__class__ if isinstance(self, OrderedSet) else OrderedSet
        if sets:
            common = set.intersection(*map(set, sets))
            items = (item for item in self if item in common)
        else:
            items = self
        return cls(items)

    def difference(self, *sets):
        """
        Returns all elements that are in this set but not the others.

        Example:
            >>> OrderedSet([1, 2, 3]).difference(OrderedSet([2]))
            OrderedSet([1, 3])
            >>> OrderedSet([1, 2, 3]).difference(OrderedSet([2]), OrderedSet([3]))
            OrderedSet([1])
            >>> OrderedSet([1, 2, 3]) - OrderedSet([2])
            OrderedSet([1, 3])
            >>> OrderedSet([1, 2, 3]).difference()
            OrderedSet([1, 2, 3])
        """
        cls = self.__class__
        if sets:
            other = set.union(*map(set, sets))
            items = (item for item in self if item not in other)
        else:
            items = self
        return cls(items)

    def issubset(self, other):
        """
        Report whether another set contains this set.

        Example:
            >>> OrderedSet([1, 2, 3]).issubset({1, 2})
            False
            >>> OrderedSet([1, 2, 3]).issubset({1, 2, 3, 4})
            True
            >>> OrderedSet([1, 2, 3]).issubset({1, 4, 3, 5})
            False
        """
        if len(self) > len(other):  # Fast check for obvious cases
            return False
        return all(item in other for item in self)

    def issuperset(self, other):
        """
        Report whether this set contains another set.

        Example:
            >>> OrderedSet([1, 2]).issuperset([1, 2, 3])
            False
            >>> OrderedSet([1, 2, 3, 4]).issuperset({1, 2, 3})
            True
            >>> OrderedSet([1, 4, 3, 5]).issuperset({1, 2, 3})
            False
        """
        if len(self) < len(other):  # Fast check for obvious cases
            return False
        return all(item in self for item in other)

    def symmetric_difference(self, other):
        """
        Return the symmetric difference of two OrderedSets as a new set.
        That is, the new set will contain all elements that are in exactly
        one of the sets.

        Their order will be preserved, with elements from `self` preceding
        elements from `other`.

        Example:
            >>> this = OrderedSet([1, 4, 3, 5, 7])
            >>> other = OrderedSet([9, 7, 1, 3, 2])
            >>> this.symmetric_difference(other)
            OrderedSet([4, 5, 9, 2])
        """
        cls = self.__class__ if isinstance(self, OrderedSet) else OrderedSet
        diff1 = cls(self).difference(other)
        diff2 = cls(other).difference(self)
        return diff1.union(diff2)

    def _update_items(self, items):
        """
        Replace the 'items' list of this OrderedSet with a new one, updating
        self.map accordingly.
        """
        self.items = items
        self.map = {item: idx for (idx, item) in enumerate(items)}

    def difference_update(self, *sets):
        """
        Update this OrderedSet to remove items from one or more other sets.

        Example:
            >>> this = OrderedSet([1, 2, 3])
            >>> this.difference_update(OrderedSet([2, 4]))
            >>> print(this)
            OrderedSet([1, 3])

            >>> this = OrderedSet([1, 2, 3, 4, 5])
            >>> this.difference_update(OrderedSet([2, 4]), OrderedSet([1, 4, 6]))
            >>> print(this)
            OrderedSet([3, 5])
        """
        items_to_remove = set()
        for other in sets:
            items_to_remove |= set(other)
        self._update_items([item for item in self.items if item not in items_to_remove])

    def intersection_update(self, other):
        """
        Update this OrderedSet to keep only items in another set, preserving
        their order in this set.

        Example:
            >>> this = OrderedSet([1, 4, 3, 5, 7])
            >>> other = OrderedSet([9, 7, 1, 3, 2])
            >>> this.intersection_update(other)
            >>> print(this)
            OrderedSet([1, 3, 7])
        """
        other = set(other)
        self._update_items([item for item in self.items if item in other])

    def symmetric_difference_update(self, other):
        """
        Update this OrderedSet to remove items from another set, then
        add items from the other set that were not present in this set.

        Example:
            >>> this = OrderedSet([1, 4, 3, 5, 7])
            >>> other = OrderedSet([9, 7, 1, 3, 2])
            >>> this.symmetric_difference_update(other)
            >>> print(this)
            OrderedSet([4, 5, 9, 2])
        """
        items_to_add = [item for item in other if item not in self]
        items_to_remove = set(other)
        self._update_items(
            [item for item in self.items if item not in items_to_remove] + items_to_add
        )

def gif(imgs,out='animation.gif',crop=None,duration=300,loop=0):
    """
    imgs: a list of image path
    out: path of gif to save
    crop: [left, up, right, bottom] of the image to crop
    duration: a single value or a list of values for each frame (in ms)
    loop: >=0, 0 forever
    """
    from PIL import Image
    # Create the frames
    frames = []
    for i in imgs:
        im = Image.open(i)
        print(f"Input image {i}: w {im.size[0]} x h {im.size[1]}")
        if crop is not None: im = im.crop(crop)
        frames.append(im)

    # Save into a GIF file
    frames[0].save(out, format='GIF',
                   append_images=frames[1:],
                   save_all=True,
                   duration=duration, loop=loop)
    print(f"Output gif: {out}")

def pdf2gif(pdf,out='animation.gif',dpi=300,crop=None,duration=300,loop=0):
    """
    pdf: a multiple-page pdf file with each page being a frame (can auto trim blank margins)
         draw.io workflow (as of 20.3.0, slightly misaligned thus causing flickering in gif; recommend to use powerpoint instead): 
         -set up paper size: custom, 4 in x 4 in
         -draw, will auto expand to another page
         -print to pdf (easier than export as pdf)
    out: path of gif to save (I have not implemented transparent background; Pillow seems to be buggy in this feature)
    dpi: None or int
    crop: [left, up, right, bottom] of the image to crop
    duration: a single value or a list of values for each frame (in ms)
    loop: >=0, 0 forever
    """
    # auto trim to 10% of blank margins
    office_pdf_crop(pdf,'-u')
    root, ext = splitext(pdf)
    pdf = root+'_cropped'+ext

    # https://pymupdf.readthedocs.io/en/latest/recipes-images.html#recipesimages
    import fitz  # import the bindings
    doc = fitz.open(fullpath(pdf))  # open document

    from PIL import Image
    frames = []

    for page in doc:  # iterate through the pages
        pix = page.get_pixmap(dpi=dpi)  # render page to an image
        # pix.save("page-%i.png" % page.number)  # store image as a PNG
        # https://github.com/pymupdf/PyMuPDF/issues/322
        im=Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        print(f"Input image {page.number+1}: w {im.size[0]} x h {im.size[1]}")
        if crop is not None: im = im.crop(crop)
        frames.append(im)

    # Save into a GIF file
    # Pillow seems to be buggy with transparent background; not implementing this feature 
    frames[0].save(out, format='GIF',
                   append_images=frames[1:],
                   save_all=True,
                   duration=duration, loop=loop)
    print(f"Output gif: {out}")

# +++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++
# debugging
# +++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++
if __name__ == "__main__":
    pass